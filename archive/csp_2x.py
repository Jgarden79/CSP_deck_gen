#import standard libraries
import eikon as ek
import datetime
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
pd.options.mode.chained_assignment = None  # default='warn'
import config
import configparser as cp
cfg=cp.ConfigParser()
cfg.read('eikonmjr.cfg.txt')
ek.set_app_key(cfg['eikon']['app_id'])
import warnings
warnings.simplefilter(action='ignore', category=FutureWarning)
from pptx import Presentation
from pptx.util import Inches, Pt
import chart_studio
chart_studio.tools.set_credentials_file(username='JGarden79', api_key='eWGoAmjzRp3GIVRTFfSR')
import plotly.io as pio
import plotly.graph_objects as go
from plotly.subplots import make_subplots
from IPython.display import display_html
from PyPDF2 import PdfFileMerger
import options_dat as od
import os

path = 'Images'
isExist = os.path.exists(path) # check if images exists
if not isExist:
     # Create a new directory because it does not exist
    os.makedirs(path)
path_2 = 'html'
isExist = os.path.exists(path_2) # check if images exists
if not isExist:
     # Create a new directory because it does not exist
    os.makedirs(path_2)
#
for f in os.listdir(path):
    os.remove(os.path.join(path,f))

for f in os.listdir(path_2):
    if f != ".git":
        os.remove(os.path.join(path_2,f))

path_3 = 'assets'
for f in os.listdir(path_3):
    os.remove(os.path.join(path_3,f))

def add_disc(file_name_in=list, file_name_out =str ):
    pdfs = []
    for i in file_name_in:
        x = 'images/{}.pdf'.format(i)
        pdfs.append(x)
    pdfs.append('OptionsTemplateSpreadsheetDisclosures.pdf')
    merger = PdfFileMerger(strict=False)
    for pdf in pdfs:
        merger.append(open(pdf, "rb"))
    with open('images/{}.pdf'.format(file_name_out), "wb") as fout:
        merger.write(fout)

def dipsplay_side_by_side(*args):
    html_str = ''
    for df in args:
        html_str += df.to_html()
    display_html(html_str.replace('table', 'table style="display:inline"'), raw=True)


def plot_trade(
    chart_data: pd.DataFrame,
    trade_title: str,
    trade_subtitle: str,
    table_description: str,
    chart_description: str,
    trade_summary_data: list,
    fee:float,
    firm: str = 'L'
):
    """
    Creates, formats, and returns a standard plotly figure for a structure option trade given the
    trade data and details.

    Args:
        chart_data (pd.DataFrame): data used to populate the trade return table and chart
        trade_title (str): trade title displayed on the upper-left corner
        trade_subtitle (str): trade subtitle
                Currently used to display the ending date of the trade.
        table_description (str): description placed at the bottom of the trade return data table
                Currently used to describe the trade details and bounds.
        chart_description (str): description placed at the bottom of the trade graph
                Currently used to display the underlying price and timestamp at the time of the
                figure's creation.
        trade_summary_data (list): data displayed in the summary table below the trade title
                Currently displayes the trades structure, term, underlyer, cap, downside before
                protection, and protection level.
        firm (str, optional): firm the trade is being generated for and affects the formatting of
                the resulting figure. Limited to values of 'L' (Lido) or 'O' (Oakhurst). This
                parameter can be updated to an enum in future refactorings.
                Defaults to 'L' (Lido).

    Raises:
        ValueError: if the firm parameter is not one of the handled cases. Allowable values include:
                'L': Lido
                'O': Oakhurst

    Returns:
        plotly.graph_objects._figure.Figure: plotly figure with required formatting that can be
                saved, exported to html, included on a tear-sheet, etc.
    """

    # Set the brand details
    if firm.upper() == 'L':
        # brand settings for lido
        color_scheme = ['#2073C9', '#1A508C', '#041525', '#E8C245', 'white']
        fonts = ['Noe Display', 'Untitled Sans']
    elif firm.upper() == 'O':
        # brand settings for Oakhurst
        color_scheme = ['#6B8D73', '#4F5151', '#333333']
        fonts = ['Open Sans Light', 'Open Sans Light']
    elif firm.upper() == 'C':
        color_scheme = ['#2073C9', '#1A508C', '#041525', '#E8C245', 'white']
        fonts = ['Noe Display', 'Untitled Sans']
    else:
        _err_str = f"Unhandled firm exception. '{firm}' is not a recognized value for firm."
        raise ValueError(_err_str)
    # Create the figure subplots
    fig = make_subplots(cols=3, rows=2, vertical_spacing=0.01, horizontal_spacing=0.05,
                        specs=[
                                [
                                    {'type': 'table', 't': 0.2, 'l': 0.015, },
                                    {'type': 'table', 'colspan': 2, 't': 0.08, 'l': 0, 'r': 0.02},
                                    None,
                                ],
                                [
                                    None,
                                    {'type': 'scatter', 'colspan': 2, 'l': 0.08, 'r': 0.02},
                                    None,
                                ]
                            ])

    # Add the Trade Return
    # fig.add_trace(go.Scatter(
    #     x=chart_data[chart_data.columns[0]],
    #     y=chart_data[chart_data.columns[1]],
    #     name='Trade Return',
    #     mode='lines',
    #     line=dict(color=color_scheme[0])), row=2, col=2)
    # add the market
    fig.add_trace(go.Scatter(
        x=chart_data[chart_data.columns[0]],
        y=chart_data[chart_data.columns[2]],
        name='Underlying Asset Price Return',
        mode='lines',
        line=dict(color="#1A508C")), row=2, col=2)
    fig.add_trace(go.Scatter(
        x=chart_data[chart_data.columns[0]],
        y=chart_data[chart_data.columns[3]],
        name='Net Return',
        mode='lines',
        line=dict(color=color_scheme[3])), row=2, col=2)
    # add the breakeven line
    fig.add_trace(go.Scatter(
        x=chart_data[chart_data.columns[0]],
        y=(chart_data[chart_data.columns[2]] - chart_data[chart_data.columns[2]]),
        name='Break Even',
        mode='lines',
        line=dict(color='red')), row=2, col=2)

    # Format the axes of the graph
    fig.update_xaxes(showgrid=False, gridwidth=1, tickfont={'color': 'black'}, ticklabelposition='inside',
                     title_font={'color': 'black', 'family': fonts[1]})
    fig.update_yaxes(title_text='Return', showgrid=True, gridwidth=1, gridcolor='black', tickformat='.0%', dtick=0.1,
                     range=[chart_data[chart_data.columns[2]].min(), chart_data[chart_data.columns[2]].max()],
                     tickfont={'color': 'black'}, title_font={'color': 'black', 'family': fonts[1]})

    # Set the layout of the figure, titles, descriptions, etc.
    fig.update_layout(
        autosize=False,
        paper_bgcolor='rgba(0,0,0,0)',
        plot_bgcolor='rgba(0,0,0,0)',
        margin=dict(l=10, r=10, t=10, b=30, autoexpand=False),
        width=816,
        height=1050,
        legend=dict(font={'color': 'black', 'family': fonts[1], 'size': 15, }, x=0.357, y=0.0, orientation='h'),
        title=dict(text=trade_title, font=dict(family=fonts[0], size=31, color=color_scheme[-1]), x=0.03, y=0.92),
        annotations=[
            go.layout.Annotation(
                showarrow=False,
                text=trade_subtitle,
                font=dict(family=fonts[1], size=20, color=color_scheme[-1]),
                xref='paper',
                yref='paper',
                x=0.04,
                y=0.87),
            go.layout.Annotation(
                showarrow=False,
                text=chart_description,
                font=dict(family=fonts[1], size=12, color='black'),
                xref='paper',
                yref='paper',
                x=0.92,
                y=0.02),
            go.layout.Annotation(
                showarrow=False,
                text=table_description,
                align='left',
                xanchor='left',
                xref='paper',
                x=0.35,
                yanchor='top',
                yref='paper',
                y=0.54,
                font=dict(
                    family=fonts[1],
                    size=13,
                    color="Black")),
        ]
    )

    # Add the data table to the figure
    fig.add_trace(go.Table(
        columnorder=[1, 2, 3, 4],
        columnwidth=[0.25, 0.25, 0.25, 0.25],
        header=dict(
            values=[f"<b>{chart_data.columns[i]}</b>" for i in [0,2,1,3]],
            line_color='Black',
            height=20,
            fill_color='white',
            font=dict(size=12, color='black', family=fonts[1])
        ),
        cells=dict(
            values=[
                [f"{i:.2f}" for i in chart_data[chart_data.columns[0]]][0::len(chart_data) // 15],
                chart_data[chart_data.columns[2]][0::len(chart_data) // 15],
                chart_data[chart_data.columns[1]][0::len(chart_data) // 15],
                chart_data[chart_data.columns[3]][0::len(chart_data) // 15],
            ],
            fill_color=['white', 'white', 'white'],
            line_color='Black',
            font=dict(size=11, color=['Black', 'Black'], family=fonts[1]),
            height=20,
            align=['center', 'center', 'center'],
            format=[[None], ['.2%'], ['.2%'], ['.2%']]
        )
    ), row=1, col=2)

    # Add the trade summary table to the figure
    fig.add_trace(go.Table(
        columnorder=[1, 2],
        columnwidth=[0.33, 0.67],
        header=dict(
            values=['', '<b>Trade</b>'],
            fill_color='rgba(239, 236, 229, 0)',
            line_color=color_scheme[-1],
            height=14,
            align=['left', 'center'],
            font=dict(size=11, color= color_scheme[-1], family=fonts[1])
        ),
        cells=dict(
            values=[
                [f"<b>{i[0]}:</b>" for i in trade_summary_data],    # catagories
                [f"<b>{i[1]}</b>" for i in trade_summary_data],     # inputs
            ],
            fill_color=['rgba(239, 236, 229, 0)', 'rgba(239, 236, 229, 0)'],
            line_color=color_scheme[-1],
            font=dict(size=11, color=[color_scheme[-1], color_scheme[-1],], family=fonts[1]),
            height=14,
            align=['left', 'left']
        )
    ), row=1, col=1)

    fig.update_layout(hovermode="x unified")

    return fig

class CSP_2x():

    '''
    Class object to calculate and visualize CSP 2x Trade.

    ...

    Parameters
    ..........

    :param RIC: <str> The Reuters identification code for the security in question
    :param get_data: <str> 'y' or 'n' only - y will pull fresh data from tr, n will use stored data
    :param expiry: <str> Expiration date of the desired options chain
    :param yld_tgt: <int> target yield of the trade, defined as: net options cost / underlying last
    :return: will return output based on the Output param of trade in question

    Attributes
    ..........

        RIC : str
            The Reuters identification code for the security in question
        expiry : str
            Expiration date of the desired options chain
        rng : flt
            Range in percentage of the payoff diagrams
        sym : str
            Ticker symbol of the security in question
        und : df
            DataFrame containing information on the security in question ('Instrument','CF_LAST','YIELD')
        last : flt
            Most recent price of the security in question
        firm : str
            Firm name displayed on PDFs and charts (Lido Advisors or Oakhurst Advisors)
        output : str
            Output style for calc_trade function ('chart', line', or 'rades')

        sorted_trade_chain_calls : df
            Dataframe of imported CALL options data

        sorted_trade_chain_puts : df
            DataFrame of imported PUT options data

        ATM : series
            Series of ATM Call data

        LPUT : series
            Series of Long Put data

        SPUT : series
            Series of Short Put data

        SCALL : series
            Series of Short Call data

        notional : flt
            Notional value of the trade

        inflows : flt
            Income from Short Put premium and expected dividend cash flow (if used)

        tgt_call_val : flt
            Calculated premium required for Short Call

        net_opt_cost : flt
            The net premium cost of the selected options (LCALL prem. - SPUT prem. - SCALL prem.)

        drag : flt
            If dividends are used, he difference between the expected dividend cash flow and the net options premium
            Otherwise, the net options premiumn

        drag_pct : flt
            Drag expressed as a percentage of the underlying stock / ETF price (drag / last)

        cap : flt
            Maximum gain in percent of the strategy

        plug : flt
            Total dividends expected during the duration of the trade. 0 if exclude_div=='y'.

        QTR_div_yld : flt
            Quarterly dividend yield of the underlying security

        months_left : flt
            Time until options expiration in months

        ds_before : flt
            Downside before protection (percent of underlying last)

        upside : flt
            Maximum gain of the trade (percent)

        annual_up : flt
            Annualized maximum gain (percent)

        protection : flt
            Protection of the trade (difference between LPUT and SPUT in percent)

        annual_p : flt
            Annualized protection (percent)

        cost : flt
            Total cost of the trade including purchasing the underlying security

    Methods
    ......

    calc_trade(output,firm)

        Calculates trade parameters and outputs Pdf, trades or trade line.

        :param output: <str> "line" - will provide terms, "trades" will provide trades, "chart" will provide fact sheet
        :param firm: <str> "L" for Lido, "O" for Oakhurst

    payoff_plot()

        Creates Pdf of strategy payoff diagram at expiration for presentation decks

    plotly_plot()

        Generates interactive payoff diagram in html with plotly

    data_table()

        Creates styled DataFrame and exports to pre-formatted Excel file for use in presentation decks

    historical_chart(period='3y')

         Creates Pdf of historical daily price chart for underlying stock or ETF using non-adjusted close prices from YFinance
         Intended for use in presentation decks

         :param period: <str> The lookback period of daily prices ('1y', '3y', '5y', '10y')

    '''

    def __init__(self,RIC,get_data,expiry,tgt_yld=0,fee=0.0125):
        self.RIC=RIC
        self.get_data=get_data
        self.expiry=expiry
        self.tgt_yld=tgt_yld
        self.fee=fee
        self.rng=0.5 # range for payoff chart (+/- %)
        if self.get_data == 'y':  # get data if necessary
            od.get_options([self.RIC])
        else:  # otherwise pass
            pass
        self.sym = pd.read_pickle('assets/{}_sym.pkl'.format(self.RIC)).iloc[0]['ticker']  # import the symbol
        self.und = pd.read_pickle('assets/{}_cached_und.pkl'.format(self.sym))  ##Gets Underlying Price###
        self.last = self.und.iloc[0]['CF_LAST']  # get most recent price
    def __repr__(self):
        return "Cap_and_Cushion(RIC={},get_data={},expiry={},tgt_yld={},fee={}".format(
                    self.RIC,
                    self.get_data,
                    self.expiry,
                    self.tgt_yld,
                    self.fee)

    def calc_trade(self,output='line',firm='L'):
        '''Calculates strategy trades'''
        self.output=output
        self.firm=firm
        # options chain set up
        chain = pd.read_pickle('assets/{}_cached_chain.pkl'.format(self.sym))  # read in the options chaing to pandas
        chain['MID'] = (chain['CF_BID'] + chain["CF_ASK"]) / 2
        trade_chain = chain[chain['EXPIR_DATE'] == self.expiry]  # isolate the chain we want
        sorted_trade_chain = trade_chain.sort_values(['PUTCALLIND', 'STRIKE_PRC'],
                                                 axis=0, )  # sort the chain by type and strike
        self.sorted_trade_chain_calls = sorted_trade_chain[sorted_trade_chain['PUTCALLIND'] == 'CALL']  # create a call df
        self.sorted_trade_chain_puts = sorted_trade_chain[sorted_trade_chain['PUTCALLIND'] == 'PUT ']  # create a put df
        # ATM CALL
        self.LCALL = self.sorted_trade_chain_calls.loc[
                self.sorted_trade_chain_calls['STRIKE_PRC'].sub(self.last).abs().idxmin()]  # get atm call
        self.notional = self.last  # trade's notional is the last traded price of the underlying
        self.scall_tgt = (self.LCALL['MID'] / 2) + (self.tgt_yld * self.last)
        self.SCALL = self.sorted_trade_chain_calls.loc[
                self.sorted_trade_chain_calls['MID'].sub(self.scall_tgt).abs().idxmin()] # get short put
        date = datetime.datetime.strptime(self.expiry, '%Y-%m-%d').date()  # convert date to datetime
        today = datetime.datetime.now().date()  # get today's date
        time_left = date - today  # days left
        self.adj_time_left = time_left / datetime.timedelta(days=1)  # convert to flt
        adj_time_left_div = np.floor(time_left / datetime.timedelta(days=1) / 365 * 4)
        self.net_opt_cost=self.LCALL['MID'] - (self.SCALL['MID'] * 2) #net cost of options (total)
        self.net_opt_cost_dlrs=self.net_opt_cost*100
        self.trade_delta = 1 + self.LCALL['DELTA'] - (self.SCALL['DELTA'] * 2)
        #caclualte management fee
        self.mgt_fee_pct = self.fee * (self.adj_time_left / 365)
        self.LCALL['Trans'] = 'BCO'  # add trade type
        self.SCALL['Trans'] = 'SCO'  # add trade type
        Trade = pd.concat([self.LCALL.to_frame().T,
                           self.SCALL.to_frame().T])  # create Trade df
        exp_date = pd.to_datetime(Trade['EXPIR_DATE'].iloc[0]).strftime('%y%m%d')  # get options formated date
        option_type = [s[0] for s in Trade["PUTCALLIND"].to_list()]  # isolate first letter of option type
        strikes = Trade['STRIKE_PRC'].to_list()  # # isolate strikes in list
        option_sym = ['{}{}{}{}'.format(self.sym, self.expiry,
                                        option_type[i],
                                        int(strikes[i])) for i in range(0, len(Trade))]  # create symbols
        Trade['Symbol'] = option_sym  # add to df
        trade_details = Trade.filter(['Trans', 'Symbol', 'STRIKE_PRC', "MID"])  # create trade_det df
        trade_details['Qty'] = [1,2]
        self.upside = ((self.SCALL['STRIKE_PRC'] / self.last) - 1) * 2  # upside calculation
        self.annual_up = self.upside * (365 / self.adj_time_left)
        self.protection = 0
        self.ds_before = "n/a"
        self.annual_p = 0
        self.months_left = np.round((self.adj_time_left / 365) * 12, 1)
        trade_line = pd.DataFrame({"Underlying Asset": self.sym,
                                   "Asset Price": self.last,
                                   'Minimum ($)': round(self.last * 100, -2),
                                   "Expiration Date": Trade['EXPIR_DATE'].iloc[0],
                                   "Months Left": self.months_left,
                                   'Potential Upside (%)': '{:.2%}'.format(self.upside),
                                   'Annual Potential Upside (%)': '{:.2%}'.format(self.annual_up),
                                   'Downside Before Protection (%)': self.ds_before,
                                   'Protection (%)': '{:.2%}'.format(self.protection),
                                   'Annual Protection (%)': '{:.2%}'.format(self.annual_p),
                                   'Spread (%)': '{:.2%}'.format(self.upside + self.protection),
                                   'Trade Cost (%)': '{:.2%}'.format(self.net_opt_cost/self.last),
                                   'Collateral': '{}'.format(self.sym)},
                                  index=[0])  # creates line
        # begin work on plot data:
        chrt_rng = np.linspace(self.LCALL['STRIKE_PRC'] * (1-self.rng), self.LCALL['STRIKE_PRC'] * (1+self.rng), 50, dtype=float)  # set chart range
        chrt_rng = np.round(chrt_rng, 2)
        lcall_ev = [np.maximum(p - self.LCALL['STRIKE_PRC'], 0) for p in chrt_rng]  # calc scall end val        
        scall_ev = [np.maximum(p - self.SCALL['STRIKE_PRC'], 0) * -2 for p in chrt_rng]  # calc scall end val
        perf_df = pd.DataFrame({"{} Price".format(self.sym): chrt_rng,
                                "LCALL": lcall_ev,
                                "SCALL": scall_ev,
                                "UND": chrt_rng,
                                })  # Create the df
        perf_df = perf_df.set_index("{} Price".format(self.sym))  # Set the mkt px as the index
        perf_df['Trade'] = perf_df.sum(axis=1)  # calculate total value
        self.cost = self.LCALL['MID'] - (self.SCALL['MID'] * 2) + self.last  # total trade cost including underlying
        perf_df['Trade Return'] = (perf_df['Trade'] / (self.cost)) - 1  # trade return
        perf_df['Trade Return - Net'] = perf_df['Trade Return'] - self.mgt_fee_pct
        perf_df = perf_df.sort_index(ascending=False)  # reorder in descending
        perf_df['{} Price Return'.format(self.sym)] = [(p / self.last) - 1 for p in perf_df.index]  # add underlying performance
        rets_tab = perf_df.filter(perf_df.columns[-3:]).reset_index()  # reset index
        rets_tab['Trade Return Net'] = [i- float(self.fee * self.adj_time_left) for i in perf_df['Trade Return'].to_list()]
        fc_date = pd.to_datetime(Trade['EXPIR_DATE'].iloc[0]).strftime("%m-%d-%Y")

        if output == 'line':  # line output
            return trade_line
        elif output == 'trades':  # trades output
            print(trade_details)
        elif output == 'execute':
            op_exec = trade_details
            op_exec['units'] = [1 for i in range(0, len(op_exec))]
            op_exec.to_csv('assets/temp_op.csv')
        else:
            # Figure objects
            trade_title = f"{self.sym}<br>2x Trade"
            trade_subtitle = f"Ending {fc_date}"
            table_description = f"The investor gets up to 200% of the price appreciation on {self.sym} beginning at a price<br>"\
                f"of {int(self.last)}, capped at {int(self.SCALL['STRIKE_PRC'])}."
            chart_description = f"{self.sym} Price (Created on: {today} when {self.sym} was at {self.last})"
            trade_summary_data = [
                ('Structure', 'Accelerated Upside with Cap'),
                ('Term', f"{self.months_left} Months"),
                ('Underlying Asset', self.sym),
                ('Cap', f"{self.upside:.2%}"),
                ('Downside Before Protection', self.ds_before),
                ('Protection', f"{self.protection:.2%}"),
                ('Fee', f"{self.fee:.2%} Annually")
            ]
            # Create the figure
            fig = plot_trade(chart_data=rets_tab, trade_title=trade_title, trade_subtitle=trade_subtitle,
                             table_description=table_description, chart_description=chart_description,
                             trade_summary_data=trade_summary_data, firm=self.firm,fee=self.fee)
            pio.write_html(fig, file='html/{}_{}_{}_2X.html'.format(self.sym, self.expiry, self.firm),
                           auto_open=False)
            pio.write_image(fig, file='Images/{}_{}_{}_2X.png'.format(self.sym, self.expiry, self.firm),format='png', scale=6)
            prs = Presentation('rebrand_test.pptx')
            prs.save('Images/{}_{}_{}_2X.pptx'.format(self.sym, self.expiry, self.firm))
            shapes = prs.slides[0].shapes
            pic = shapes.add_picture('Images/{}_{}_{}_2X.png'.format(self.sym, self.expiry, self.firm), Inches(0.001), Inches(0.001),
                                     height=Inches(10.99), width=Inches(8.5))
            prs.save('Images/{}_{}_{}_2X.pptx'.format(self.sym, self.expiry, self.firm))
            return

    def payoff_plot(self):

        '''Prints strategy payodd diagram at expiration for presentation decks'''

        # Functions to calculate options payoffs at expiration
        def call_payoff(stock_range,strike,premium):
            return np.where(stock_range>strike,stock_range-strike,0)-premium
        # Define stock price range at expiration
        up_dn = 0.5
        stock_range=np.arange((1-up_dn)*self.last,(1+up_dn)*self.last,1)
        # Calculate payoffs for individual legs
        long_call=self.LCALL['STRIKE_PRC']
        long_call_prem=self.LCALL['MID']
        short_call=self.SCALL['STRIKE_PRC']
        short_call_prem=self.SCALL['MID']
        payoff_long_call=call_payoff(stock_range,long_call,long_call_prem)
        payoff_short_call=call_payoff(stock_range,short_call,short_call_prem)*-2
        # Calculate Strategy Payoff
        strategy=(((payoff_long_call+payoff_short_call+stock_range)/self.last)-1)
        strategy_net=strategy - self.mgt_fee_pct
        buy_hold_ret=((stock_range/(self.last))-1)

        self.cap = ((self.SCALL['STRIKE_PRC'] / self.last) - 1) * 2
        protection =0 # no protection
        #Create Visualization
        plt.style.use('fivethirtyeight')
        fig,ax=plt.subplots(figsize=(14,6))
        from matplotlib.ticker import PercentFormatter
        plt.yticks(np.arange(-up_dn,up_dn+0.01,0.2),fontsize=16)
        plt.gca().yaxis.set_major_formatter(PercentFormatter(1))
        plt.xticks(fontsize=16)
        plt.gca().xaxis.set_major_formatter('${x:1.0f}')
        plt.plot(stock_range,buy_hold_ret,c='grey',lw=3,ls='dashed',label='{} Price Return'.format(self.sym))
        plt.plot(stock_range,strategy,c='green',lw=3,label='Strategy Return (gross)')
        plt.plot(stock_range,strategy_net,lw=2,label='Strategy Return (net)')
        plt.vlines(x=self.last,ymin=buy_hold_ret.min(),ymax=buy_hold_ret.max(),linestyle='dashed',color='grey',lw=2)
        plt.annotate('Current Price: ${:.2f}'.format(self.last),
                xy=(self.last,(buy_hold_ret.max()-0)*0.5),fontsize=12,
                rotation=90,
                horizontalalignment='right',verticalalignment='center')
        plt.hlines(y=0,xmin=stock_range.min(),xmax=stock_range.max(),color='gray')
        plt.legend(loc='upper left',fontsize=14)
        plt.ylabel('Return at Expiration',fontsize=16,fontweight='bold')
        plt.xlabel('{} Price at Expiration'.format(self.sym),fontsize=16,fontweight='bold')
        plt.suptitle('Stock: {}     Current Price: ${:.2f}     2x Strategy'.format(
                                    self.sym,
                                    self.last),
                     fontsize=16)
        plt.title('Options expiration: {}'.format(self.expiry),fontsize=14)
        plt.tight_layout()
        plt.savefig('images/{}_2X_payoff_plot.png'.format(self.sym))
        plt.show();

    def data_table(self):
        '''Creates strategy data table for use in presentations'''
        self.shares = input("Enter number of shares held:")
        self.shares = float(self.shares)
        self.contracts = self.shares / 100
        self.port_val = self.shares * self.last
        self.mgt_fee = self.mgt_fee_pct * self.port_val
        # Functions to calculate options payoffs at EXPIRY
        def call_payoff(stock_range,strike,premium):
            return np.where(stock_range>strike,stock_range-strike,0)-premium
        # Calculate payoffs for individual legs
        long_call=self.LCALL['STRIKE_PRC']
        long_call_prem=self.LCALL['MID']
        short_call=self.SCALL['STRIKE_PRC']
        short_call_prem=self.SCALL['MID']
        # Create DataFrame of Stock Prices every 5%
        pct_range=1+np.arange(-0.30,0.31,0.05)
        percent_range=pct_range*self.last
        payoff_long_call=call_payoff(percent_range,long_call,long_call_prem) * self.shares
        payoff_short_call = call_payoff(percent_range,short_call,short_call_prem) * self.shares * -2
        # Calculate Strategy Payoff
        stock_pl_5=(percent_range-self.last)*self.shares
        strat_pl_5=payoff_long_call + payoff_short_call + stock_pl_5
        strat_pl_5_net=strat_pl_5 - self.mgt_fee
        ret=pct_range-1
        df=pd.DataFrame({'return':ret,'2_x':strat_pl_5/self.port_val,'2_x_net':strat_pl_5_net/self.port_val,'stock':stock_pl_5/self.port_val})
        df.set_index('return',inplace=True)
        df=df.T
        display(df)
        ## UPDATE PRE-FORMATTED SPREADSHEET ###
        import openpyxl
        xl_file=r'C:\Users\mreis\OneDrive - Lido Advisors, LLC\Concentrated Stock\Lido CSP\Lido_CSP_strategy_workbook.xlsx'
        wb=openpyxl.load_workbook(xl_file) #load worksbbok
        ws3=wb.worksheets[2] #define worksheet
        ws3['A53'] = self.last * self.shares
        ws3['A55'] = self.net_opt_cost_dlrs * self.contracts
        ws3['A57'] = self.upside
        ws3['A59'] = '{:.0f} days'.format(self.adj_time_left) #DTE

        ws4=wb.worksheets[3] #define worksheet
        ws4['A21'] = f'{self.sym} with Lido CSP (gross)'
        ws4['A22'] = f'{self.sym} with Lido CSP (net)'
        ws4['A23'] = self.sym
        # #append covered call returns table
        for row in np.arange(0,len(df),1):
            for col in np.arange(0,len(df.columns),1):
                ws4.cell(row=row+21,column=col+2).value=df.iloc[row,col]

        wb.save(xl_file)#save workbook
        wb.close()

    def historical_chart(self,period='3y'):

        '''
        Creates historical daily price chart of underlying stock or ETF using non-adjusted close prices from YFinance.
        Intended for use in presentation decks

        :param period: <str> The lookback period of daily prices ('1y', '3y', '5y', '10y')

        '''
        import yfinance as yf
        df=yf.download('{}'.format(self.sym),period=period)['Close'].to_frame()
        plt.style.use('fivethirtyeight')
        plt.figure(figsize=(14,8))
        plt.plot(df['Close'],label=self.sym)
        plt.tick_params(labeltop=False,labelright=True)

        plt.hlines(y=self.SCALL['STRIKE_PRC'],xmin=df.index[0],xmax=df.index[-1],
                      linestyle='dashed',color='red',alpha=0.5,
                      label='Cap: ${} (+{:.1%}) OTM'.format(self.SCALL['STRIKE_PRC'],(self.SCALL['STRIKE_PRC']/self.last)-1))
        plt.hlines(y=self.last,xmin=df.index[0],xmax=df.index[-1],
                      linestyle='dashed',color='gray',alpha=0.5,
                      label='Current Price: ${}'.format(self.last))
        plt.hlines(y=self.LCALL['STRIKE_PRC'],xmin=df.index[0],xmax=df.index[-1],
                      linestyle='dashed',color='green',alpha=0.5,
                      label='2x Upside begins: ${} ({:.1%})'.format(self.LCALL['STRIKE_PRC'],(self.LCALL['STRIKE_PRC']/self.last)-1))
        plt.title('Stock: {}     |     Trade Duration: {:.1f}yrs     |     Last: ${:.2f}'.format(self.sym,self.adj_time_left/365,self.last),
                    fontsize=16,fontweight='bold')
        plt.legend(loc='upper left',fontsize=14)
        plt.xticks(fontsize=14)
        plt.yticks(fontsize=14)
        plt.gca().yaxis.set_major_formatter('${x:,.0f}')
        plt.xlabel('Date',fontsize=16,fontweight='bold')
        plt.ylabel(f'{self.sym} Price',fontsize=16,fontweight='bold')
        plt.tight_layout()
        plt.savefig('images/{}_{}r_price_chart_2X.png'.format(self.sym,period))
        plt.show();

#     def payoff_plot_plain(self):

#         '''Prints strategy payodd diagram at expiration for presentation decks'''

#         # Functions to calculate options payoffs at expiration
#         def call_payoff(stock_range,strike,premium):
#             return np.where(stock_range>strike,stock_range-strike,0)-premium

#         # Define stock price range at expiration
#         up_dn = 0.4
#         stock_range=np.arange((1-up_dn)*self.last,(1+up_dn)*self.last,1)

#         # Calculate payoffs for individual legs
#         long_put=self.LPUT['STRIKE_PRC']
#         long_put_prem=self.LPUT['MID']
#         short_put=self.SPUT['STRIKE_PRC']
#         short_put_prem=self.SPUT['MID']
#         short_call=self.SCALL['STRIKE_PRC']
#         short_call_prem=self.SCALL['MID']

#         payoff_long_put=put_payoff(stock_range,long_put,long_put_prem)
#         payoff_short_put=put_payoff(stock_range,short_put,short_put_prem)*-1
#         payoff_short_call=call_payoff(stock_range,short_call,short_call_prem)*-1

#         # Calculate Strategy Payoff
#         strategy=(((payoff_short_put+payoff_long_put+payoff_short_call+stock_range+self.plug)/self.last)-1)
#         strategy_net=(((payoff_short_put+payoff_long_put+payoff_short_call+stock_range+self.plug)/self.last)-1) - self.mgt_fee_pct
#         buy_hold_ret=((stock_range/(self.last))-1)

#         self.cap = (self.SCALL['STRIKE_PRC'] / self.last) - 1
#         protection = np.abs(((self.LPUT['STRIKE_PRC'] - self.SPUT['STRIKE_PRC'])  / self.last))  # protection calculation
#         gap = (self.LPUT['STRIKE_PRC'] / self.last) - 1 #calcualte gap

#         #Create Visualization
#         plt.style.use('fivethirtyeight')
#         fig,ax=plt.subplots(figsize=(8,8))

#         plt.plot(stock_range,buy_hold_ret,c='grey',lw=3,ls='dashed',label='{} Price Return'.format('SPY'))
#         plt.plot(stock_range,strategy,c='green',lw=3,label='Strategy Return')
#         plt.xlabel(F'{self.sym}',fontsize=10,fontweight='bold')
#         plt.ylabel('STRATEGY PAYOFF PROFILE',fontsize=10,fontweight='bold')
#         plt.vlines(x=self.last,ymin=buy_hold_ret.min(),ymax=buy_hold_ret.max(),color='grey',lw=2)
#         plt.hlines(y=0,xmin=stock_range.min(),xmax=stock_range.max(),color='grey',lw=2)
#         plt.legend(loc='lower right',fontsize=16)
#         plt.grid(False)
#         plt.axis('off')
#         plt.tight_layout()
#         plt.show();