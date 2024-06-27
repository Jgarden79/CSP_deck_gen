import eikon as ek
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt

pd.options.mode.chained_assignment = None  # default='warn'
import config
from plotly.subplots import make_subplots

ek.set_app_key(config.eikon_key)

import datetime
from IPython.display import display_html
import chart_studio
import plotly.graph_objects as go

chart_studio.tools.set_credentials_file(username='JGarden79', api_key='eWGoAmjzRp3GIVRTFfSR')
import plotly.io as pio
import options_dat as od
import warnings
import os
from PyPDF2 import PdfFileMerger

warnings.simplefilter(action='ignore', category=FutureWarning)

BOND_ETF_COLLATERAL = {
    "2024-12-20": "BSCO",
    "2025-12-19": "BSCP",
}

path = 'Images'
isExist = os.path.exists(path)  # check if images exists
if not isExist:
    # Create a new directory because it does not exist
    os.makedirs(path)
path_2 = 'html'
isExist = os.path.exists(path_2)  # check if images exists
if not isExist:
    # Create a new directory because it does not exist
    os.makedirs(path_2)
#
for f in os.listdir(path):
    os.remove(os.path.join(path, f))

for f in os.listdir(path_2):
    if f != ".git":
        os.remove(os.path.join(path_2, f))

path_3 = 'assets'
for f in os.listdir(path_3):
    os.remove(os.path.join(path_3, f))


def add_disc(file_name_in=list, file_name_out=str):
    pdfs = []
    for i in file_name_in:
        x = 'images/{}.pdf'.format(i)
        pdfs.append(x)
    pdfs.append('OptionsTemplateSpreadsheetDisclosures.pdf')
    merger = PdfFileMerger(strict=False)
    for pdf in pdfs:
        merger.append(open(pdf, "rb"))
    with open('images/{}.pdf'.format(file_name_out), "wb") as fout:
        merger.write(fout)


def dipsplay_side_by_side(*args):
    html_str = ''
    for df in args:
        html_str += df.to_html()
    display_html(html_str.replace('table', 'table style="display:inline"'), raw=True)


od.get_bonds()  # imports bond data in the morning


def plot_trade(
        chart_data: pd.DataFrame,
        trade_title: str,
        trade_subtitle: str,
        table_description: str,
        chart_description: str,
        trade_summary_data: list,
        fee: float,
        firm: str = 'L'
):
    """
    Creates, formats, and returns a standard plotly figure for a structure option trade given the
    trade data and details.

    Args:
        chart_data (pd.DataFrame): data used to populate the trade return table and chart
        trade_title (str): trade title displayed on the upper-left corner
        trade_subtitle (str): trade subtitle
                Currently used to display the ending date of the trade.
        table_description (str): description placed at the bottom of the trade return data table
                Currently used to describe the trade details and bounds.
        chart_description (str): description placed at the bottom of the trade graph
                Currently used to display the underlying price and timestamp at the time of the
                figure's creation.
        trade_summary_data (list): data displayed in the summary table below the trade title
                Currently displayes the trades structure, term, underlyer, cap, downside before
                protection, and protection level.
        firm (str, optional): firm the trade is being generated for and affects the formatting of
                the resulting figure. Limited to values of 'L' (Lido) or 'O' (Oakhurst) or 'C' (CFM). This
                parameter can be updated to an enum in future refactorings.
                Defaults to 'L' (Lido).

    Raises:
        ValueError: if the firm parameter is not one of the handled cases. Allowable values include:
                'L': Lido
                'O': Oakhurst

    Returns:
        plotly.graph_objects._figure.Figure: plotly figure with required formatting that can be
                saved, exported to html, included on a tear-sheet, etc.
    """
    # Set the brand details
    if firm.upper() == 'L':
        # brand settings for lido
        color_scheme = ['#2073C9', '#1A508C', '#041525', '#E8C245', 'white']
        fonts = ['Noe Display', 'Untitled Sans']
    elif firm.upper() == 'O':
        # brand settings for Oakhurst
        color_scheme = ['#6B8D73', '#4F5151', '#333333','black', 'white',]
        fonts = ['Open Sans Light', 'Open Sans Light']
    elif firm.upper() == 'C':
        color_scheme = ['#2073C9', '#1A508C', '#041525', '#E8C245', 'white']
        fonts = ['Noe Display', 'Untitled Sans']
    else:
        _err_str = f"Unhandled firm exception. '{firm}' is not a recognized value for firm."
        raise ValueError(_err_str)
    # Create the figure subplots
    fig = make_subplots(cols=3, rows=2, vertical_spacing=0.01, horizontal_spacing=0.05,
                        specs=[
                            [
                                {'type': 'table', 't': 0.2, 'l': 0.015, },
                                {'type': 'table', 'colspan': 2, 't': 0.08, 'l': 0, 'r': 0.02},
                                None,
                            ],
                            [
                                None,
                                {'type': 'scatter', 'colspan': 2, 'l': 0.08, 'r': 0.02},
                                None,
                            ]
                        ])

    # Add the Trade Return
    # fig.add_trace(go.Scatter(
    #     x=chart_data[chart_data.columns[0]],
    #     y=chart_data[chart_data.columns[1]],
    #     name='Trade Return',
    #     mode='lines',
    #     line=dict(color=color_scheme[0])), row=2, col=2)
    # add the market
    fig.add_trace(go.Scatter(
        x=chart_data[chart_data.columns[0]],
        y=chart_data[chart_data.columns[2]],
        name='Underlying Asset Price Return',
        mode='lines',
        line=dict(color="#1A508C")), row=2, col=2)
    fig.add_trace(go.Scatter(
        x=chart_data[chart_data.columns[0]],
        y=chart_data[chart_data.columns[3]],
        name='Net Return',
        mode='lines',
        # line=dict(color=color_scheme[3])), row=2, col=2)
        line=dict(color=color_scheme[0])), row=2, col=2)
    # add the breakeven line
    fig.add_trace(go.Scatter(
        x=chart_data[chart_data.columns[0]],
        y=(chart_data[chart_data.columns[2]] - chart_data[chart_data.columns[2]]),
        name='Break Even',
        mode='lines',
        line=dict(color='red')), row=2, col=2)

    # Format the axes of the graph
    fig.update_xaxes(showgrid=False, gridwidth=1, tickfont={'color': 'black'}, ticklabelposition='inside',
                     title_font={'color': 'black', 'family': fonts[1]})
    fig.update_yaxes(title_text='Return', showgrid=True, gridwidth=1, gridcolor='black', tickformat='.0%', dtick=0.1,
                     range=[chart_data[chart_data.columns[2]].min(), chart_data[chart_data.columns[2]].max()],
                     tickfont={'color': 'black'}, title_font={'color': 'black', 'family': fonts[1]})

    # Set the layout of the figure, titles, descriptions, etc.
    fig.update_layout(
        autosize=False,
        paper_bgcolor='rgba(0,0,0,0)',
        plot_bgcolor='rgba(0,0,0,0)',
        margin=dict(l=10, r=10, t=10, b=30, autoexpand=False),
        width=816,
        height=1050,
        legend=dict(font={'color': 'black', 'family': fonts[1], 'size': 15, }, x=0.357, y=0.0, orientation='h'),
        title=dict(text=trade_title, font=dict(family=fonts[0], size=31, color=color_scheme[-1]), x=0.03, y=0.92),
        annotations=[
            go.layout.Annotation(
                showarrow=False,
                text=trade_subtitle,
                font=dict(family=fonts[1], size=20, color=color_scheme[-1]),
                xref='paper',
                yref='paper',
                x=0.04,
                y=0.87),
            go.layout.Annotation(
                showarrow=False,
                text=chart_description,
                font=dict(family=fonts[1], size=12, color='black'),
                xref='paper',
                yref='paper',
                x=0.92,
                y=0.02),
            go.layout.Annotation(
                showarrow=False,
                text=table_description,
                align='left',
                xanchor='left',
                xref='paper',
                x=0.35,
                yanchor='top',
                yref='paper',
                y=0.54,
                font=dict(
                    family=fonts[1],
                    size=13,
                    color="Black")),
        ]
    )

    # Add the data table to the figure
    fig.add_trace(go.Table(
        columnorder=[1, 2, 3, 4],
        # columnwidth=[0.25, 0.25, 0.25, 0.25],
        columnwidth=[0.33, 0.34, 0.33],
        header=dict(
            # values=[f"<b>{chart_data.columns[i]}</b>" for i in [0, 2, 1, 3]],
            values=[f"<b>{chart_data.columns[i]}</b>" for i in [0, 2, 3]],
            line_color='Black',
            height=20,
            fill_color='white',
            font=dict(size=12, color='black', family=fonts[1])
        ),
        cells=dict(
            values=[
                [f"{i:.2f}" for i in chart_data[chart_data.columns[0]]][0::len(chart_data) // 15],
                chart_data[chart_data.columns[2]][0::len(chart_data) // 15],
                # chart_data[chart_data.columns[1]][0::len(chart_data) // 15],
                chart_data[chart_data.columns[3]][0::len(chart_data) // 15],
            ],
            fill_color=['white', 'white', 'white'],
            line_color='Black',
            font=dict(size=11, color=['Black', 'Black'], family=fonts[1]),
            height=20,
            align=['center', 'center', 'center'],
            format=[[None], ['.2%'], ['.2%'], ['.2%']]
        )
    ), row=1, col=2)

    # Add the trade summary table to the figure
    fig.add_trace(go.Table(
        columnorder=[1, 2],
        columnwidth=[0.33, 0.67],
        header=dict(
            values=['', '<b>Trade</b>'],
            fill_color='rgba(239, 236, 229, 0)',
            line_color=color_scheme[-1],
            height=14,
            align=['left', 'center'],
            font=dict(size=11, color=color_scheme[-1], family=fonts[1])
        ),
        cells=dict(
            values=[
                [f"<b>{i[0]}:</b>" for i in trade_summary_data],  # catagories
                [f"<b>{i[1]}</b>" for i in trade_summary_data],  # inputs
            ],
            fill_color=['rgba(239, 236, 229, 0)', 'rgba(239, 236, 229, 0)'],
            line_color=color_scheme[-1],
            font=dict(size=11, color=[color_scheme[-1], color_scheme[-1], ], family=fonts[1]),
            height=14,
            align=['left', 'left']
        )
    ), row=1, col=1)

    fig.update_layout(hovermode="x unified")

    return fig


def cap_cush(RIC, Date, get_data, output, firm, prot=None, fee=0.0125):
    """
    calculates a standard cap and cushion trade
    :param RIC: <str> the Reurters identification code for the security in question
    :param Date: <str> expiration date of the desired options chain.
    :param get_data: <str> y or n only - y will pull fresh data from tr, n will use stored data
    :param output: <str> "line" - will provide terms, "trades" will provide trades, "chart" will provide fact sheet,
    "execute" will generate csv for upload.
    :param firm: <str> "L" for Lido, "O" for Oakhurst
    :param prot: <flt> optional target protection
    :return: will return output based on the Output param of trade in question
    """
    if get_data == 'y':  # get data if necessary
        od.get_options([RIC])
    else:  # otherwise pass
        pass
    sym = pd.read_pickle('assets/{}_sym.pkl'.format(RIC)).iloc[0]['ticker']  # import the symbol
    und = pd.read_pickle('assets/{}_cached_und.pkl'.format(sym))  ##Gets Underlying Price###
    bnd_dat = pd.read_pickle('assets/bond_data.pkl')  # import bond data
    last = und.iloc[0]['CF_LAST']  # get most recent price

    # Set up Bond Data
    date = datetime.datetime.strptime(Date, '%Y-%m-%d')  ### Sets date into string
    bnd_lst = pd.read_pickle('assets/cached_bndlst.pkl')  ### reads list of UST bonds from pickles
    bnd_lst['MATUR_DATE'] = pd.to_datetime(bnd_lst['MATUR_DATE'])  # convert maturity dates to datetime
    bnd_dat['MATUR_DATE'] = pd.to_datetime(bnd_dat['MATUR_DATE'])  # convert maturity dates to datetime
    bnd_dat = bnd_dat.dropna()  # remove na's
    ind = abs(bnd_dat['MATUR_DATE'] - date).idxmin()  ## id the relevant bond
    if bnd_dat.loc[[ind]]['DIRTY_PRC'] is np.NaN:  ##this if ensures we skip over "when issued bonds"
        ind = ind + 1  # move up one
        bond = bnd_dat.loc[[ind]]  # isolate bond
    else:  # otherwise
        bond = bnd_dat.loc[[ind]]  # isolate bond
    bond_CU = bond.iloc[0]['Instrument']  # store cusip

    # begin options work
    # options chain set up
    chain = pd.read_pickle('assets/{}_cached_chain.pkl'.format(sym))  # read in the options chaing to pandas
    chain['MID'] = (chain['CF_ASK'] + chain['CF_BID']) / 2
    trade_chain = chain[chain['EXPIR_DATE'] == Date]  # isolate the chain we want
    sorted_trade_chain = trade_chain.sort_values(['PUTCALLIND', 'STRIKE_PRC'],
                                                 axis=0, )  # sort the chain by type and strike
    sorted_trade_chain_calls = sorted_trade_chain[sorted_trade_chain['PUTCALLIND'] == 'CALL']  # create a call df
    sorted_trade_chain_puts = sorted_trade_chain[sorted_trade_chain['PUTCALLIND'] == 'PUT ']  # create a put df

    # ATM CALL
    LCALL = sorted_trade_chain_calls.loc[
        sorted_trade_chain_calls['STRIKE_PRC'].sub(last).abs().idxmin()]  # get atm call

    # Solve for the other two legs
    date = datetime.datetime.strptime(Date, '%Y-%m-%d').date()  # convert date to datetime
    today = datetime.datetime.now().date()  # get today's date
    time_left = date - today  # days left
    adj_time_left = time_left / datetime.timedelta(days=1)  # convert to flt
    time_to_bnd = (adj_time_left / 365)  # convert to annual
    bnd = bnd_dat[bnd_dat['Instrument'] == bond_CU]  # get relevant bond info
    bnd['Minimum PAR'] = LCALL['STRIKE_PRC'] * 100
    bnd['1x Cost'] = bnd['Minimum PAR'] * bnd['DIRTY_PRC'] / 100
    rf = bnd['SEC_YLD_1'].iloc[0] * time_to_bnd  # calculate and store deanualzied yield
    notional = LCALL['STRIKE_PRC']  # trade's notional is the strike of the option
    bnd['TR'] = (notional * np.exp(rf / 100)) - notional  # calculate total expected return from bond

    # trades where no downside value is establised:
    if prot == None:  # if there is no target protection
        fv = LCALL['CF_ASK'] * np.exp(rf / 100)  # calc future value of call prem
        TGT = (fv / LCALL['STRIKE_PRC'])  # calculate market implied target
        top_x = np.floor(last * (1 + TGT))  # calculate target for short call
        SCALL = sorted_trade_chain_calls.loc[
            sorted_trade_chain_calls['STRIKE_PRC'].sub(top_x).abs().idxmin()]  # get short call
        inflows = SCALL['CF_BID'] + bnd['TR'].iloc[0]  # calculate all cash inflows
        tgt_put_val = LCALL['CF_ASK'] - inflows  # calculate target put value
        SPUT = sorted_trade_chain_puts.loc[
            sorted_trade_chain_puts['CF_BID'].sub(tgt_put_val).abs().idxmin()]  # get short put
    else:
        TGT = 1 - prot  # calc protection target
        sput_x = LCALL['STRIKE_PRC'] * TGT  # calculate protection target strike
        SPUT = sorted_trade_chain_puts.loc[
            sorted_trade_chain_puts['STRIKE_PRC'].sub(sput_x).abs().idxmin()]  # get short put
        inflows = SPUT['CF_BID'] + bnd['TR'].iloc[0]  # calculate all cash inflows
        tgt_call_val = LCALL['CF_ASK'] - inflows  # calculate target put value
        SCALL = sorted_trade_chain_calls.loc[
            sorted_trade_chain_calls['CF_BID'].sub(tgt_call_val).abs().idxmin()]  # get SCALL
    SPUT['Trans'] = 'SPO'  # add trade type
    LCALL['Trans'] = 'BCO'  # add trade type
    SCALL['Trans'] = 'SCO'  # add trade type
    Trade = pd.concat([pd.DataFrame(SCALL).T, pd.DataFrame(LCALL).T, pd.DataFrame(SPUT).T]).reset_index(drop=True)
    exp_date = pd.to_datetime(Trade['EXPIR_DATE'].iloc[0]).strftime('%y%m%d')  # get options formated date
    option_type = [s[0] for s in Trade["PUTCALLIND"].to_list()]  # isolate first letter of option type
    strikes = Trade['STRIKE_PRC'].to_list()  # # isolate strikes in list
    option_sym = ['{}{}{}{}'.format(sym, exp_date, option_type[i], int(strikes[i])) for i in
                  range(0, len(Trade))]  # create symbols
    Trade['Symbol'] = option_sym  # add to df
    trade_details = Trade.filter(['Trans', 'Symbol', 'STRIKE_PRC', "MID"])  # create trade_det df
    upside = (SCALL['STRIKE_PRC'] / LCALL['STRIKE_PRC']) - 1  # upside calculation
    annual_up = upside * (365 / adj_time_left)
    protection = np.abs((SPUT['STRIKE_PRC'] / LCALL['STRIKE_PRC']) - 1)  # protection calculation
    annual_p = protection * (365 / adj_time_left)
    months_left = np.round((adj_time_left / 365) * 12, 1)
    trade_line = pd.DataFrame({"Underlying Asset": sym, "Asset Price": last,
                               'Minimum ($)': round(LCALL['STRIKE_PRC'] * 100, -2),
                               "Expiration Date": Trade['EXPIR_DATE'].iloc[0], "Months Left": months_left,
                               'Potential Upside (%)': '{:.2%}'.format(upside),
                               'Annual Potential Upside (%)': '{:.2%}'.format(annual_up),
                               'Downside Before Protection (%)': '{:.2%}'.format(0),
                               'Protection (%)': '{:.2%}'.format(protection),
                               'Annual_Protection (%)': '{:.2%}'.format(annual_p),
                               'Spread (%)': '{:.2%}'.format(upside + protection), 'Collateral': 'UST'},
                              index=[0])  # creates line

    # begin work on plot data:
    chrt_rng = np.linspace(LCALL['STRIKE_PRC'] * 0.5, LCALL['STRIKE_PRC'] * 1.5, 50, dtype=float)  # set chart range
    chrt_rng = np.round(chrt_rng, 2)
    scall_ev = [np.maximum(p - SCALL['STRIKE_PRC'], 0) * -1 for p in chrt_rng]  # calc scall end val
    lcall_ev = [np.maximum(p - LCALL['STRIKE_PRC'], 0) for p in chrt_rng]  # calc lcall end val
    sput_ev = [np.maximum(SPUT['STRIKE_PRC'] - p, 0) * -1 for p in chrt_rng]  # calc sput end val
    perf_df = pd.DataFrame({"{} Price".format(sym): chrt_rng, "SCALL": scall_ev, "LCALL": lcall_ev, "SPUT": sput_ev,
                            "BOND": LCALL['STRIKE_PRC']})  # Create the df
    perf_df = perf_df.set_index("{} Price".format(sym))  # Set the mkt px as the index
    perf_df['Trade'] = perf_df.sum(axis=1)  # calculate total value
    cost = LCALL['CF_ASK'] - SCALL['CF_BID'] - SPUT['CF_BID'] + LCALL['STRIKE_PRC'] - bnd['TR'].iloc[0]  # trade cost
    perf_df['Trade Return'] = (perf_df['Trade'] / (cost)) - 1  # trade return
    perf_df = perf_df.sort_index(ascending=False)  # reorder in descending
    perf_df['{} Price Return'.format(sym)] = [(p / last) - 1 for p in perf_df.index]  # add underlying performance
    rets_tab = perf_df.filter(perf_df.columns[-2:]).reset_index()  # reset index
    rets_tab['Trade Return Net'] = [i - float(fee * time_to_bnd) for i in perf_df['Trade Return'].to_list()]
    fc_date = pd.to_datetime(Trade['EXPIR_DATE'].iloc[0]).strftime("%m-%d-%Y")
    upside = perf_df['Trade Return'].max()

    if output == 'line':  # line output
        return trade_line
    elif output == 'trades':  # trades output
        print(trade_details)
        print(bnd.filter(['Instrument', 'DIRTY_PRC', 'SEC_YLD_1', 'COUPN_RATE', 'MATUR_DATE',
                          'Minimum PAR', '1x Cost', ]))
    elif output == 'execute':
        op_exec = trade_details
        op_exec['units'] = [1 for i in range(0, len(op_exec))]
        bnd_exec = bnd.filter(['Instrument', 'Minimum PAR'])
        bnd_exec = bnd_exec.rename(columns={'Minimum PAR': 'units'})
        op_exec.to_csv('assets/temp_op.csv')
        bnd_exec.to_csv('assets/temp_bnd.csv')

    else:
        # Figure objects
        trade_title = f"{sym} Cap and<br>Cushion Trade"
        trade_subtitle = f"Ending {fc_date}"
        table_description = f"The investor gets about 100% of the price appreciation on {sym} beginning at a price<br>" \
                            f"of {int(LCALL['STRIKE_PRC'])}, capped at {int(SCALL['STRIKE_PRC'])} and full protection between " \
                            f"{int(LCALL['STRIKE_PRC'])} and {int(SPUT['STRIKE_PRC'])} ({-protection:.2%})"
        chart_description = f"{sym} Price (Created on: {today} when {sym} was at {last})"
        trade_summary_data = [
            ('Structure', 'Capped Upside with Protection'),
            ('Term', f"{months_left} Months"),
            ('Underlying Asset', sym),
            ('Cap', f"{upside:.2%}"),
            ('Downside Before Protection', 'No Gap'),
            ('Protection', f"{protection:.2%}"),
            ('Fee', f"{fee:.2%} Annually")
        ]

        # Create the figure
        fig = plot_trade(chart_data=rets_tab, trade_title=trade_title, trade_subtitle=trade_subtitle,
                         table_description=table_description, chart_description=chart_description,
                         trade_summary_data=trade_summary_data, firm=firm, fee=fee)

        pio.write_html(fig, file='html/{}_{}_{}_capp_cush.html'.format(Date, firm, sym),
                       auto_open=False)
        pio.write_image(fig, file='Images/{}_{}_{}_capp_cush.png'.format(Date, firm, sym), format='png', scale=6)
        if firm == "L":
            prs = Presentation('rebrand_test.pptx')
        elif firm == "O":
            prs = Presentation('rebrand_test_oak.pptx')
        else:
            prs = Presentation('rebrand_test_cfm.pptx')
        prs.save('Images/{}_{}_{}_capp_cush.pptx'.format(Date, firm, sym))
        shapes = prs.slides[0].shapes
        pic = shapes.add_picture('Images/{}_{}_{}_capp_cush.png'.format(Date, firm, sym), Inches(0.001), Inches(0.001),
                                 height=Inches(10.99), width=Inches(8.5))
        prs.save('Images/{}_{}_{}_capp_cush.pptx'.format(Date, firm, sym))

        return


def cap_cush_ETF(RIC, Date, get_data, output, firm, prot=None, exclude_div='n', custom_div='n', fee=0.0125):
    """
    ETF style cap and cushion
    :param RIC: <str> the Reurters identification code for the security in question
    :param Date: <str> expiration date of the desired options chain.
    :param get_data: <str> y or n only - y will pull fresh data from tr, n will use stored data
    :param output: <str> "line" - will provide terms, "trades" will provide trades, "chart" will provide fact sheet
    :param firm: <str> "L" for Lido, "O" for Oakhurst
    :param prot: <flt> target protection
    :param exclude_div:<str> y = exclude all dividends from calc
    :param custom_div: <flt> custom annual dividend to use
    :return:will return output based on the Output param of trade in question
    """
    if get_data == 'y':  # get data if necessary
        od.get_options([RIC])
    else:  # otherwise pass
        pass
    sym = pd.read_pickle('assets/{}_sym.pkl'.format(RIC)).iloc[0]['ticker']  # import the symbol
    und = pd.read_pickle('assets/{}_cached_und.pkl'.format(sym))  ##Gets Underlying Price###
    und = und.fillna(0)
    bnd_dat = pd.read_pickle('assets/bond_data.pkl')  # import bond data
    last = und.iloc[0]['CF_LAST']  # get most recent price

    # Set up Bond Data
    date = datetime.datetime.strptime(Date, '%Y-%m-%d')  ### Sets date into string
    bnd_lst = pd.read_pickle('assets/cached_bndlst.pkl')  ### reads list of UST bonds from pickles
    bnd_lst['MATUR_DATE'] = pd.to_datetime(bnd_lst['MATUR_DATE'])  # convert maturity dates to datetime
    bnd_dat['MATUR_DATE'] = pd.to_datetime(bnd_dat['MATUR_DATE'])  # convert maturity dates to datetime
    bnd_dat = bnd_dat.dropna()  # remove na's
    ind = abs(bnd_dat['MATUR_DATE'] - date).idxmin()  ## id the relevant bond
    if bnd_dat.loc[[ind]]['DIRTY_PRC'] is np.NaN:  ##this if ensures we skip over "when issued bonds"
        ind = ind + 1  # move up one
        bond = bnd_dat.loc[[ind]]  # isolate bond
    else:  # otherwise
        bond = bnd_dat.loc[[ind]]  # isolate bond
    bond_CU = bond.iloc[0]['Instrument']  # store cusip
    # get div data
    if custom_div == 'n':
        if exclude_div == 'n':
            QTR_div_yld = und.iloc[0]['YIELD'] / 400
        else:
            QTR_div_yld = 0
    else:
        QTR_div_yld = input('Enter Annual Dividend Yield:')
        QTR_div_yld = float(QTR_div_yld) / 4

    # begin options work
    # options chain set up
    chain = pd.read_pickle('assets/{}_cached_chain.pkl'.format(sym))  # read in the options chaing to pandas
    chain['MID'] = (chain['CF_ASK'] + chain['CF_BID']) / 2
    trade_chain = chain[chain['EXPIR_DATE'] == Date]  # isolate the chain we want
    sorted_trade_chain = trade_chain.sort_values(['PUTCALLIND', 'STRIKE_PRC'],
                                                 axis=0, )  # sort the chain by type and strike
    sorted_trade_chain_calls = sorted_trade_chain[sorted_trade_chain['PUTCALLIND'] == 'CALL']  # create a call df
    sorted_trade_chain_puts = sorted_trade_chain[sorted_trade_chain['PUTCALLIND'] == 'PUT ']  # create a put df

    LPUT = sorted_trade_chain_puts.loc[sorted_trade_chain_puts['STRIKE_PRC'].sub(last).abs().idxmin()]  # get atm put
    # Solve for the other two legs
    LCALL = sorted_trade_chain_calls.loc[
        sorted_trade_chain_calls['STRIKE_PRC'].sub(last).abs().idxmin()]  # get atm call
    date = datetime.datetime.strptime(Date, '%Y-%m-%d').date()  # convert date to datetime
    today = datetime.datetime.now().date()  # get today's date
    time_left = date - today  # days left
    adj_time_left = time_left / datetime.timedelta(days=1)  # convert to flt
    adj_time_left_div = np.floor(time_left / datetime.timedelta(days=1) / 365 * 4)
    time_to_bnd = (adj_time_left / 365)  # convert to annual
    bnd = bnd_dat[bnd_dat['Instrument'] == bond_CU]  # get relevant bond info
    rf = bnd['SEC_YLD_1'].iloc[0] * time_to_bnd  # calculate and store deanualzied yield
    notional = LCALL['STRIKE_PRC']  # trade's notional is the strike of the option
    plug = np.round(adj_time_left_div * QTR_div_yld * last, 2)  # total divs expected
    plug = np.nan_to_num(plug, 0)  # if no divs replace with 0
    # trades where no downside value is establised:
    if prot == None:  # if there is no target protection
        fv = LCALL['CF_ASK'] * np.exp(rf / 100)  # calc future value of call prem
        TGT = (fv / LCALL['STRIKE_PRC'])  # calculate market implied target
        top_x = np.floor(last * (1 + TGT))  # calculate target for short call
        SCALL = sorted_trade_chain_calls.loc[
            sorted_trade_chain_calls['STRIKE_PRC'].sub(top_x).abs().idxmin()]  # get short call
        inflows = SCALL['CF_BID'] + plug  # calculate all cash inflows
        tgt_put_val = LPUT['CF_ASK'] - inflows  # calculate target put value
        SPUT = sorted_trade_chain_puts.loc[
            sorted_trade_chain_puts['CF_BID'].sub(tgt_put_val).abs().idxmin()]  # get short put
    else:
        TGT = 1 - prot  # calc protection target
        sput_x = LPUT['STRIKE_PRC'] * TGT  # calculate protection target strike
        SPUT = sorted_trade_chain_puts.loc[
            sorted_trade_chain_puts['STRIKE_PRC'].sub(sput_x).abs().idxmin()]  # get short put
        inflows = SPUT['CF_BID'] + plug  # calculate all cash inflows
        tgt_call_val = LPUT['CF_ASK'] - inflows  # calculate target put value
        SCALL = sorted_trade_chain_calls.loc[
            sorted_trade_chain_calls['CF_BID'].sub(tgt_call_val).abs().idxmin()]  # get SCALL
    SPUT['Trans'] = 'SPO'  # add trade type
    LPUT['Trans'] = 'BPO'  # add trade type
    SCALL['Trans'] = 'SCO'  # add trade type
    Trade = pd.concat([SCALL.to_frame().T, LPUT.to_frame().T, SPUT.to_frame().T])  # create Trade df
    exp_date = pd.to_datetime(Trade['EXPIR_DATE'].iloc[0]).strftime('%y%m%d')  # get options formated date
    option_type = [s[0] for s in Trade["PUTCALLIND"].to_list()]  # isolate first letter of option type
    strikes = Trade['STRIKE_PRC'].to_list()  # # isolate strikes in list
    option_sym = ['{}{}{}{}'.format(sym, exp_date, option_type[i], int(strikes[i])) for i in
                  range(0, len(Trade))]  # create symbols
    Trade['Symbol'] = option_sym  # add to df
    trade_details = Trade.filter(['Trans', 'Symbol', 'STRIKE_PRC', "MID"])  # create trade_det df
    upside = (SCALL['STRIKE_PRC'] / last) - 1  # upside calculation
    annual_up = upside * (365 / adj_time_left)
    protection = np.abs((SPUT['STRIKE_PRC'] / LPUT['STRIKE_PRC']) - 1)  # protection calculation
    ds_before = (last / LPUT['STRIKE_PRC']) - 1
    annual_p = protection * (365 / adj_time_left)
    months_left = np.round((adj_time_left / 365) * 12, 1)
    trade_line = pd.DataFrame({"Underlying Asset": sym, "Asset Price": last,
                               'Minimum ($)': round(last * 100, -2),
                               "Expiration Date": Trade['EXPIR_DATE'].iloc[0], "Months Left": months_left,
                               'Potential Upside (%)': '{:.2%}'.format(upside),
                               'Annual Potential Upside (%)': '{:.2%}'.format(annual_up),
                               'Downside Before Protection (%)': '{:.2%}'.format(ds_before),
                               'Protection (%)': '{:.2%}'.format(protection),
                               'Annual_Protection (%)': '{:.2%}'.format(annual_p),
                               'Spread (%)': '{:.2%}'.format(upside + protection), 'Collateral': '{}'.format(sym)},
                              index=[0])  # creates line

    # begin work on plot data:
    chrt_rng = np.linspace(LPUT['STRIKE_PRC'] * 0.5, LPUT['STRIKE_PRC'] * 1.5, 50, dtype=float)  # set chart range
    chrt_rng = np.round(chrt_rng, 2)
    scall_ev = [np.maximum(p - SCALL['STRIKE_PRC'], 0) * -1 for p in chrt_rng]  # calc scall end val
    lput_ev = [np.maximum(LPUT['STRIKE_PRC'] - p, 0) for p in chrt_rng]  # calc lcall end val
    sput_ev = [np.maximum(SPUT['STRIKE_PRC'] - p, 0) * -1 for p in chrt_rng]  # calc sput end val
    perf_df = pd.DataFrame({"{} Price".format(sym): chrt_rng, "SCALL": scall_ev, "LPUT": lput_ev, "SPUT": sput_ev,
                            "UND": chrt_rng, "DIVS": [plug for i in range(0, len(chrt_rng))]})  # Create the df
    perf_df = perf_df.set_index("{} Price".format(sym))  # Set the mkt px as the index
    perf_df['Trade'] = perf_df.sum(axis=1)  # calculate total value
    cost = LPUT['CF_ASK'] - SCALL['CF_BID'] - SPUT['CF_BID'] + last  # trade cost
    perf_df['Trade Return'] = (perf_df['Trade'] / (cost)) - 1  # trade return
    perf_df = perf_df.sort_index(ascending=False)  # reorder in descending
    perf_df['{} Price Return'.format(sym)] = [(p / last) - 1 for p in perf_df.index]  # add underlying performance
    rets_tab = perf_df.filter(perf_df.columns[-2:]).reset_index()  # reset index
    rets_tab['Trade Return Net'] = [i - float(fee * time_to_bnd) for i in perf_df['Trade Return'].to_list()]
    fc_date = pd.to_datetime(Trade['EXPIR_DATE'].iloc[0]).strftime("%m-%d-%Y")
    if exclude_div == 'n':  # use div data for chart
        div_dat = " "
    else:
        div_dat = " Excluding Dividends "  # insert if no divs
    upside = perf_df['Trade Return'].max()

    if output == 'line':  # line output
        return trade_line
    elif output == 'trades':  # trades output
        und['Min_Shares'] = 100
        und['Min_Cost'] = und['Min_Shares'] * last
        print(trade_details)
        print(und)
    elif output == 'execute':
        op_exec = trade_details
        op_exec['units'] = [1 for i in range(0, len(op_exec))]
        und['Min_Shares'] = 100
        und_exec = und.filter(['Instrument', 'CF_LAST'])
        op_exec.to_csv('assets/temp_op.csv')
        und_exec.to_csv('assets/temp_und.csv')
    else:
        # Figure objects
        trade_title = f"{sym} Cap and<br>Cushion Trade"
        trade_subtitle = f"Ending {fc_date}"
        table_description = f"The investor gets about 100% of the price appreciation on {sym} beginning at a price<br>" \
                            f"of {int(LCALL['STRIKE_PRC'])}, capped at {int(SCALL['STRIKE_PRC'])} and full protection between " \
                            f"{int(LCALL['STRIKE_PRC'])} and {int(SPUT['STRIKE_PRC'])} ({-protection:.2%})"
        chart_description = f"{sym} Price (Created on: {today} when {sym} was at {last})"
        trade_summary_data = [
            ('Structure', 'Capped Upside with Protection'),
            ('Term', f"{months_left} Months"),
            ('Underlying Asset', sym),
            ('Cap', f"{upside:.2%}"),
            ('Downside Before Protection', 'No Gap'),
            ('Protection', f"{protection:.2%}"),
            ('Fee', f"{fee:.2%} Annually")
        ]

        # Create the figure
        fig = plot_trade(chart_data=rets_tab, trade_title=trade_title, trade_subtitle=trade_subtitle,
                         table_description=table_description, chart_description=chart_description,
                         trade_summary_data=trade_summary_data, firm=firm, fee=fee)
        pio.write_html(fig, file='html/{}_{}_{}_capp_cush.html'.format(Date, firm, sym),
                       auto_open=True)
        pio.write_image(fig, file='Images/{}_{}_{}_capp_cush.png'.format(Date, firm, sym), format='png', scale=6)
        if firm == "L":
            prs = Presentation('rebrand_test.pptx')
        elif firm == "O":
            prs = Presentation('rebrand_test_oak.pptx')
        else:
            prs = Presentation('rebrand_test_cfm.pptx')
        prs.save('Images/{}_{}_{}_capp_cush.pptx'.format(Date, firm, sym))
        shapes = prs.slides[0].shapes
        pic = shapes.add_picture('Images/{}_{}_{}_capp_cush.png'.format(Date, firm, sym), Inches(0.001), Inches(0.001),
                                 height=Inches(10.99), width=Inches(8.5))
        prs.save('Images/{}_{}_{}_capp_cush.pptx'.format(Date, firm, sym))

    return


def powerlift(get_data, output, firm, version=1, fee=0.0125):
    """
    special trade with more protections using corps
    :param get_data: <str> y or n only - y will pull fresh data from tr, n will use stored data
    :param output: <str> "line" - will provide terms, "trades" will provide trades, "chart" will provide fact sheet
    :param firm: <str> "L" for Lido, "O" for Oakhurst
    :param version: <flt> 1 for 50% protection 2 for 36% protection, 3 for longer trade.
    :return: based on output param
    """
    if get_data == 'y':  # get data if necessary
        od.get_options(["SPY"])
        if version == 3 or version == 4 or version == 2:
            ytw = od.calc_credit('BSCO', get_data='y')
            Date = '2024-12-20'
            bnd_tick = "BSCO.O"

        else:
            ytw = od.calc_credit('BSCN', get_data='y')
            Date = '2023-12-15'
            bnd_tick = "BSCN.O"
    else:  # otherwise pass
        if version == 3 or version == 4 or version == 2:
            ytw = od.calc_credit('BSCO', get_data='y')
            Date = '2024-12-20'
            bnd_tick = "BSCO.O"

        else:
            ytw = od.calc_credit('BSCN', get_data='y')
            Date = '2023-12-15'
            bnd_tick = "BSCN.O"

    sym = pd.read_pickle('assets/{}_sym.pkl'.format("SPY")).iloc[0]['ticker']  # import the symbol
    und = pd.read_pickle('assets/{}_cached_und.pkl'.format(sym))  ##Gets Underlying Price###
    bnd_dat = pd.read_pickle('assets/bond_data.pkl')  # import bond data
    last = und.iloc[0]['CF_LAST']  # get most recent price

    # Set up Bond Data

    # begin options work
    # options chain set up
    chain = pd.read_pickle('assets/{}_cached_chain.pkl'.format(sym))  # read in the options chaing to pandas
    chain['MID'] = (chain['CF_ASK'] + chain['CF_BID']) / 2
    trade_chain = chain[chain['EXPIR_DATE'] == Date]  # isolate the chain we want
    sorted_trade_chain = trade_chain.sort_values(['PUTCALLIND', 'STRIKE_PRC'],
                                                 axis=0, )  # sort the chain by type and strike
    sorted_trade_chain_calls = sorted_trade_chain[sorted_trade_chain['PUTCALLIND'] == 'CALL']  # create a call df
    sorted_trade_chain_puts = sorted_trade_chain[sorted_trade_chain['PUTCALLIND'] == 'PUT ']  # create a put df

    # ATM CALL
    LCALL = sorted_trade_chain_calls.loc[
        sorted_trade_chain_calls['STRIKE_PRC'].sub(last).abs().idxmin()]  # get atm call

    # Solve for the other two legs
    date = datetime.datetime.strptime(Date, '%Y-%m-%d').date()  # convert date to datetime
    today = datetime.datetime.now().date()  # get today's date
    time_left = date - today  # days left
    adj_time_left = time_left / datetime.timedelta(days=1)  # convert to flt
    time_to_bnd = (adj_time_left / 365)  # convert to annual
    bnd = ek.get_data(bnd_tick, fields=['CF_LAST'])[0]  ###Gets BSCN Price###
    rf = (1 + np.round(ytw, 4)) ** time_to_bnd
    bnd['1x Cost'] = (LCALL['STRIKE_PRC'] * 100) / rf
    bnd['Shares'] = (bnd['1x Cost'] / bnd['CF_LAST']).round(0)
    notional = LCALL['STRIKE_PRC'] * 1  # trade's notional is the strike of the option
    bnd['TR'] = (bnd['1x Cost'] * rf) / 100 - bnd['1x Cost'] / 100  # calculate total expected return from bond

    # trades where no downside value is establised:
    if version == 1 or version == 3:  # if this is the 50% trade
        TGT = 0.6  # calc protection target
        sput_x = LCALL['STRIKE_PRC'] * TGT  # calculate protection target strike
        SPUT = sorted_trade_chain_puts.loc[
            sorted_trade_chain_puts['STRIKE_PRC'].sub(sput_x).abs().idxmin()]  # get short put
        inflows = SPUT['CF_BID'] + bnd['TR'].iloc[0]  # calculate all cash inflows
        tgt_call_val = LCALL['CF_ASK'] - inflows  # calculate target put value
        SCALL = sorted_trade_chain_calls.loc[
            sorted_trade_chain_calls['CF_BID'].sub(tgt_call_val).abs().idxmin()]  # get SCALL
    elif version == 2:
        TGT = 0.85  # calc protection target
        sput_x = LCALL['STRIKE_PRC'] * TGT  # calculate protection target strike
        SPUT = sorted_trade_chain_puts.loc[
            sorted_trade_chain_puts['STRIKE_PRC'].sub(sput_x).abs().idxmin()]  # get short put
        inflows = SPUT['CF_BID'] + bnd['TR'].iloc[0]  # calculate all cash inflows
        tgt_call_val = LCALL['CF_ASK'] - inflows  # calculate target put value
        SCALL = sorted_trade_chain_calls.loc[
            sorted_trade_chain_calls['CF_BID'].sub(tgt_call_val).abs().idxmin()]  # get SCALL
    else:
        TGT = 0.75  # calc protection target
        sput_x = LCALL['STRIKE_PRC'] * TGT  # calculate protection target strike
        SPUT = sorted_trade_chain_puts.loc[
            sorted_trade_chain_puts['STRIKE_PRC'].sub(sput_x).abs().idxmin()]  # get short put
        inflows = SPUT['CF_BID'] + bnd['TR'].iloc[0]  # calculate all cash inflows
        tgt_call_val = LCALL['CF_ASK'] - inflows  # calculate target put value
        SCALL = sorted_trade_chain_calls.loc[
            sorted_trade_chain_calls['CF_BID'].sub(tgt_call_val).abs().idxmin()]  # get SCALL
    SPUT['Trans'] = 'SPO'  # add trade type
    LCALL['Trans'] = 'BCO'  # add trade type
    SCALL['Trans'] = 'SCO'  # add trade type
    Trade = pd.concat([pd.DataFrame(SCALL).T, pd.DataFrame(LCALL).T, pd.DataFrame(SPUT).T]).reset_index(drop=True)
    exp_date = pd.to_datetime(Trade['EXPIR_DATE'].iloc[0]).strftime('%y%m%d')  # get options formated date
    option_type = [s[0] for s in Trade["PUTCALLIND"].to_list()]  # isolate first letter of option type
    strikes = Trade['STRIKE_PRC'].to_list()  # # isolate strikes in list
    option_sym = ['{}{}{}{}'.format(sym, exp_date, option_type[i], int(strikes[i])) for i in
                  range(0, len(Trade))]  # create symbols
    Trade['Symbol'] = option_sym  # add to df
    trade_details = Trade.filter(['Trans', 'Symbol', 'STRIKE_PRC', "MID"])  # create trade_det df
    upside = (SCALL['STRIKE_PRC'] / LCALL['STRIKE_PRC']) - 1  # upside calculation
    annual_up = upside * (365 / adj_time_left)
    protection = np.abs((SPUT['STRIKE_PRC'] / LCALL['STRIKE_PRC']) - 1)  # protection calculation
    annual_p = protection * (365 / adj_time_left)
    months_left = np.round((adj_time_left / 365) * 12, 1)
    trade_line = pd.DataFrame({"Underlying Asset": sym, "Asset Price": last,
                               'Minimum ($)': round(LCALL['STRIKE_PRC'] * 100, -2),
                               "Expiration Date": Trade['EXPIR_DATE'].iloc[0], "Months Left": months_left,
                               'Potential Upside (%)': '{:.2%}'.format(upside),
                               'Annual Potential Upside (%)': '{:.2%}'.format(annual_up),
                               'Downside Before Protection (%)': '{:.2%}'.format(0),
                               'Protection (%)': '{:.2%}'.format(protection),
                               'Annual_Protection (%)': '{:.2%}'.format(annual_p),
                               'Spread (%)': '{:.2%}'.format(upside + protection), 'Collateral': 'BSCN'},
                              index=[0])  # creates line

    # begin work on plot data:
    chrt_rng = np.linspace(LCALL['STRIKE_PRC'] * 0.3, LCALL['STRIKE_PRC'] * 1.3, 50, dtype=float)  # set chart range
    chrt_rng = np.round(chrt_rng, 2)
    scall_ev = [np.maximum(p - SCALL['STRIKE_PRC'], 0) * -1 for p in chrt_rng]  # calc scall end val
    lcall_ev = [np.maximum(p - LCALL['STRIKE_PRC'], 0) for p in chrt_rng]  # calc lcall end val
    sput_ev = [np.maximum(SPUT['STRIKE_PRC'] - p, 0) * -1 for p in chrt_rng]  # calc sput end val
    perf_df = pd.DataFrame({"{} Price".format(sym): chrt_rng, "SCALL": scall_ev, "LCALL": lcall_ev, "SPUT": sput_ev,
                            "BOND": LCALL['STRIKE_PRC']})  # Create the df
    perf_df = perf_df.set_index("{} Price".format(sym))  # Set the mkt px as the index
    perf_df['Trade'] = perf_df.sum(axis=1)  # calculate total value
    cost = LCALL['CF_ASK'] - SCALL['CF_BID'] - SPUT['CF_BID'] + LCALL['STRIKE_PRC'] - bnd['TR'].iloc[0]  # trade cost
    perf_df['Trade Return'] = (perf_df['Trade'] / (cost)) - 1  # trade return
    perf_df = perf_df.sort_index(ascending=False)  # reorder in descending
    perf_df['{} Price Return'.format(sym)] = [(p / last) - 1 for p in perf_df.index]  # add underlying performance
    rets_tab = perf_df.filter(perf_df.columns[-2:]).reset_index()  # reset index
    rets_tab['Trade Return Net'] = [i - float(fee * time_to_bnd) for i in perf_df['Trade Return'].to_list()]
    fc_date = pd.to_datetime(Trade['EXPIR_DATE'].iloc[0]).strftime("%m-%d-%Y")
    upside = perf_df['Trade Return'].max()

    if output == 'line':  # line output
        return trade_line
    elif output == 'trades':  # trades output
        print(trade_details)
        print(bnd.filter(['Instrument', "CF_LAST", "Shares", '1x Cost', ]))
    elif output == 'execute':
        op_exec = trade_details
        op_exec['units'] = [1 for i in range(0, len(op_exec))]
        bnd_exec = bnd.filter(['Instrument', 'Shares', 'CF_LAST'])
        bnd_exec = bnd_exec.rename(columns={'Shares': 'units'})
        op_exec.to_csv('assets/temp_op.csv')
        bnd_exec.to_csv('assets/temp_bnd.csv')
    else:
        # Figure objects
        trade_title = "Lido Powerlift<br>Trade"
        trade_subtitle = f"Ending {fc_date}"
        table_description = f"The investor gets about 100% of the price appreciation on {sym} beginning at a price<br>" \
                            f"of {int(LCALL['STRIKE_PRC'])}, capped at {int(SCALL['STRIKE_PRC'])} and full protection between " \
                            f"{int(LCALL['STRIKE_PRC'])} and {int(SPUT['STRIKE_PRC'])} ({-protection:.2%})"
        chart_description = f"{sym} Price (Created on: {today} when {sym} was at {last})"
        trade_summary_data = [
            ('Structure', 'High Protection'),
            ('Term', f"{months_left} Months"),
            ('Underlying Asset', sym),
            ('Cap', f"Capped at {upside:.2%}"),
            ('Downside Before Protection', 'No Gap'),
            ('Protection', f"{protection:.2%}"),
            ('Fee', f"{fee:.2%} Annually"),
        ]

        # Create the figure
        fig = plot_trade(chart_data=rets_tab, trade_title=trade_title, trade_subtitle=trade_subtitle,
                         table_description=table_description, chart_description=chart_description,
                         trade_summary_data=trade_summary_data, firm=firm, fee=fee)

        pio.write_html(fig, file='html/{}_{}_{}_powerlift_v{}.html'.format(Date, firm, sym, version),
                       auto_open=False)
        pio.write_image(fig, file='Images/{}_{}_{}_powerlift_v{}.png'.format(Date, firm, sym, version), format='png',
                        scale=6)
        if firm == "L":
            prs = Presentation('rebrand_test.pptx')
        elif firm == "O":
            prs = Presentation('rebrand_test_oak.pptx')
        else:
            prs = Presentation('rebrand_test_cfm.pptx')
        prs.save('Images/{}_{}_{}_powerlift_v{}.pptx'.format(Date, firm, sym, version))
        shapes = prs.slides[0].shapes
        pic = shapes.add_picture('Images/{}_{}_{}_powerlift_v{}.png'.format(Date, firm, sym, version), Inches(0.001),
                                 Inches(0.001),
                                 height=Inches(10.99), width=Inches(8.5))
        prs.save('Images/{}_{}_{}_powerlift_v{}.pptx'.format(Date, firm, sym, version))
        return


def gap_trade(RIC, Date, get_data, output, firm, gap, tgt, exclude_div, custom_div='n', fee=0.0125):
    """

    :param RIC: <str> the Reurters identification code for the security in question
    :param Date: <str> expiration date of the desired options chain.
    :param get_data: <str> y or n only - y will pull fresh data from tr, n will use stored data
    :param output: <str> "line" - will provide terms, "trades" will provide trades, "chart" will provide fact sheet
    :param firm: <str> "L" for Lido, "O" for Oakhurst
    :param gap: <flt> the target exposure before protection begins
    :param tgt: <flt> the target ANNUALIZED upside for the trade
    :param exclude_div: <str> "y" = calculate trade without dividends
    :param custom_div: <str> "y" = use a custom dividend yield (specified when prompted)
    :return: based on output
    """
    if get_data == 'y':  # get data if necessary
        od.get_options([RIC])
    else:  # otherwise pass
        pass
    sym = pd.read_pickle('assets/{}_sym.pkl'.format(RIC)).iloc[0]['ticker']  # import the symbol
    und = pd.read_pickle('assets/{}_cached_und.pkl'.format(sym))  ##Gets Underlying Price###
    bnd_dat = pd.read_pickle('assets/bond_data.pkl')  # import bond data
    last = und.iloc[0]['CF_LAST']  # get most recent price

    # Set up Bond Data
    date = datetime.datetime.strptime(Date, '%Y-%m-%d')  ### Sets date into string
    bnd_lst = pd.read_pickle('assets/cached_bndlst.pkl')  ### reads list of UST bonds from pickles
    bnd_lst['MATUR_DATE'] = pd.to_datetime(bnd_lst['MATUR_DATE'])  # convert maturity dates to datetime
    bnd_dat['MATUR_DATE'] = pd.to_datetime(bnd_dat['MATUR_DATE'])  # convert maturity dates to datetime
    bnd_dat = bnd_dat.dropna()  # remove na's
    ind = abs(bnd_dat['MATUR_DATE'] - date).idxmin()  ## id the relevant bond
    if bnd_dat.loc[[ind]]['DIRTY_PRC'] is np.NaN:  ##this if ensures we skip over "when issued bonds"
        ind = ind + 1  # move up one
        bond = bnd_dat.loc[[ind]]  # isolate bond
    else:  # otherwise
        bond = bnd_dat.loc[[ind]]  # isolate bond
    bond_CU = bond.iloc[0]['Instrument']  # store cusip
    # get div data
    if custom_div == 'n':
        if exclude_div == 'n':
            QTR_div_yld = und.iloc[0]['YIELD'] / 400
        else:
            QTR_div_yld = 0
    else:
        QTR_div_yld = input('Enter Annual Dividend Yield:')
        QTR_div_yld = float(QTR_div_yld) / 4

    # begin options work
    # options chain set up
    chain = pd.read_pickle('assets/{}_cached_chain.pkl'.format(sym))  # read in the options chaing to pandas
    chain['MID'] = (chain['CF_BID'] + chain["CF_ASK"]) / 2
    trade_chain = chain[chain['EXPIR_DATE'] == Date]  # isolate the chain we want
    sorted_trade_chain = trade_chain.sort_values(['PUTCALLIND', 'STRIKE_PRC'],
                                                 axis=0, )  # sort the chain by type and strike
    sorted_trade_chain_calls = sorted_trade_chain[sorted_trade_chain['PUTCALLIND'] == 'CALL']  # create a call df
    sorted_trade_chain_puts = sorted_trade_chain[sorted_trade_chain['PUTCALLIND'] == 'PUT ']  # create a put df

    LPUT = sorted_trade_chain_puts.loc[
        sorted_trade_chain_puts['STRIKE_PRC'].sub(last * (1 - gap)).abs().idxmin()]  # get atm put
    # Solve for the other two legs
    LCALL = sorted_trade_chain_calls.loc[
        sorted_trade_chain_calls['STRIKE_PRC'].sub(last).abs().idxmin()]  # get atm call
    date = datetime.datetime.strptime(Date, '%Y-%m-%d').date()  # convert date to datetime
    today = datetime.datetime.now().date()  # get today's date
    time_left = date - today  # days left
    adj_time_left = time_left / datetime.timedelta(days=1)  # convert to flt
    adj_time_left_div = np.floor(time_left / datetime.timedelta(days=1) / 365 * 4)
    time_to_bnd = (adj_time_left / 365)  # convert to annual
    bnd = bnd_dat[bnd_dat['Instrument'] == bond_CU]  # get relevant bond info
    # rf = bnd['SEC_YLD_1'].iloc[0] * time_to_bnd  # calculate and store deanualzied yield
    notional = LCALL['STRIKE_PRC']  # trade's notional is the strike of the option
    plug = np.round(adj_time_left_div * QTR_div_yld * last, 2)  # total divs expected
    plug = np.nan_to_num(plug, 0)  # if no divs replace with 0
    TGT = np.exp(tgt * time_to_bnd)  # calculate market implied target
    top_x = np.floor(last * TGT)  # calculate target for short call
    SCALL = sorted_trade_chain_calls.loc[
        sorted_trade_chain_calls['STRIKE_PRC'].sub(top_x).abs().idxmin()]  # get short call
    inflows = SCALL['CF_BID'] + plug  # calculate all cash inflows
    tgt_put_val = LPUT['CF_ASK'] - inflows  # calculate target put value
    SPUT = sorted_trade_chain_puts.loc[
        sorted_trade_chain_puts['CF_BID'].sub(tgt_put_val).abs().idxmin()]  # get short put
    SPUT['Trans'] = 'SPO'  # add trade type
    LPUT['Trans'] = 'BPO'  # add trade type
    SCALL['Trans'] = 'SCO'  # add trade type
    Trade = pd.concat([SCALL.to_frame().T, LPUT.to_frame().T, SPUT.to_frame().T])  # create Trade df
    exp_date = pd.to_datetime(Trade['EXPIR_DATE'].iloc[0]).strftime('%y%m%d')  # get options formated date
    option_type = [s[0] for s in Trade["PUTCALLIND"].to_list()]  # isolate first letter of option type
    strikes = Trade['STRIKE_PRC'].to_list()  # # isolate strikes in list
    option_sym = ['{}{}{}{}'.format(sym, exp_date, option_type[i], int(strikes[i])) for i in
                  range(0, len(Trade))]  # create symbols
    Trade['Symbol'] = option_sym  # add to df
    trade_details = Trade.filter(['Trans', 'Symbol', 'STRIKE_PRC', "MID"])  # create trade_det df
    upside = (SCALL['STRIKE_PRC'] / last) - 1  # upside calculation
    annual_up = upside * (365 / adj_time_left)
    protection = np.abs(
        (((SPUT['STRIKE_PRC'] / last) - 1) - ((LPUT['STRIKE_PRC'] / last) - 1)))  # protection calculation
    ds_before = (last / LPUT['STRIKE_PRC']) - 1
    annual_p = protection * (365 / adj_time_left)
    months_left = np.round((adj_time_left / 365) * 12, 1)
    trade_line = pd.DataFrame({"Underlying Asset": sym, "Asset Price": last,
                               'Minimum ($)': round(last * 100, -2),
                               "Expiration Date": Trade['EXPIR_DATE'].iloc[0], "Months Left": months_left,
                               'Potential Upside (%)': '{:.2%}'.format(upside),
                               'Annual Potential Upside (%)': '{:.2%}'.format(annual_up),
                               'Downside Before Protection (%)': '{:.2%}'.format(ds_before),
                               'Protection (%)': '{:.2%}'.format(protection),
                               'Annual_Protection (%)': '{:.2%}'.format(annual_p),
                               'Spread (%)': '{:.2%}'.format(upside + protection), 'Collateral': '{}'.format(sym)},
                              index=[0])  # creates line

    # begin work on plot data:
    chrt_rng = np.linspace(LCALL['STRIKE_PRC'] * 0.3, LCALL['STRIKE_PRC'] * 1.3, 50, dtype=float)  # set chart range
    chrt_rng = np.round(chrt_rng, 2)
    scall_ev = [np.maximum(p - SCALL['STRIKE_PRC'], 0) * -1 for p in chrt_rng]  # calc scall end val
    lput_ev = [np.maximum(LPUT['STRIKE_PRC'] - p, 0) for p in chrt_rng]  # calc lcall end val
    sput_ev = [np.maximum(SPUT['STRIKE_PRC'] - p, 0) * -1 for p in chrt_rng]  # calc sput end val
    perf_df = pd.DataFrame({"{} Price".format(sym): chrt_rng, "SCALL": scall_ev, "LPUT": lput_ev, "SPUT": sput_ev,
                            "UND": chrt_rng, "DIVS": [plug for i in range(0, len(chrt_rng))]})  # Create the df
    perf_df = perf_df.set_index("{} Price".format(sym))  # Set the mkt px as the index
    perf_df['Trade'] = perf_df.sum(axis=1)  # calculate total value
    cost = LPUT['CF_ASK'] - SCALL['CF_BID'] - SPUT['CF_BID'] + last  # trade cost
    perf_df['Trade Return'] = (perf_df['Trade'] / (cost)) - 1  # trade return
    perf_df = perf_df.sort_index(ascending=False)  # reorder in descending
    perf_df['{} Price Return'.format(sym)] = [(p / last) - 1 for p in perf_df.index]  # add underlying performance
    rets_tab = perf_df.filter(perf_df.columns[-2:]).reset_index()  # reset index
    rets_tab['Trade Return Net'] = [i - float(fee * time_to_bnd) for i in perf_df['Trade Return'].to_list()]
    fc_date = pd.to_datetime(Trade['EXPIR_DATE'].iloc[0]).strftime("%m-%d-%Y")
    upside = perf_df['Trade Return'].max()
    if exclude_div == 'n':  # use div data for chart
        div_dat = " "
    else:
        div_dat = " Excluding Dividends "  # insert if no divs

    if output == 'line':  # line output
        return trade_line
    elif output == 'trades':  # trades output
        und['Min_Shares'] = 100
        und['Min_Cost'] = und['Min_Shares'] * last
        print(trade_details)
        print(und)
    elif output == 'execute':
        op_exec = trade_details
        op_exec['units'] = [1 for i in range(0, len(op_exec))]
        und['Min_Shares'] = 100
        und_exec = und.filter(['Instrument', 'CF_LAST'])
        op_exec.to_csv('assets/temp_op.csv')
        und_exec.to_csv('assets/temp_und.csv')
    else:
        # Figure objects
        trade_title = f"{sym} Cap and<br>Gap Trade"
        trade_subtitle = f"Ending {fc_date}"
        table_description = f"The investor gets about 100% of the price appreciation on {sym} beginning at a price<br>" \
                            f"of {int(LCALL['STRIKE_PRC'])}, capped at {int(SCALL['STRIKE_PRC'])} and full protection between " \
                            f"{int(LCALL['STRIKE_PRC'])} and {int(SPUT['STRIKE_PRC'])} ({-protection:.2%})"
        chart_description = f"{sym} Price (Created on: {today} when {sym} was at {last})"
        trade_summary_data = [
            ('Structure', 'Capped Upside with deferred protection'),
            ('Term', f"{months_left} Months"),
            ('Underlying Asset', sym),
            ('Cap', f"Capped at {upside:.2%}"),
            ('Downside Before Protection', f"{ds_before:.2%}"),
            ('Protection', f"{protection:.2%}"),
            ('Fee', f"{fee:.2%} Annually"),
        ]

        # Create the figure
        fig = plot_trade(chart_data=rets_tab, trade_title=trade_title, trade_subtitle=trade_subtitle,
                         table_description=table_description, chart_description=chart_description,
                         trade_summary_data=trade_summary_data, firm=firm, fee=fee)

        pio.write_html(fig, file='html/{}_{}_{}_gap.html'.format(Date, firm, sym),
                       auto_open=False)
        pio.write_image(fig, file='Images/{}_{}_{}_gap.png'.format(Date, firm, sym), format='png', scale=6)
        if firm == "L":
            prs = Presentation('rebrand_test.pptx')
        elif firm == "O":
            prs = Presentation('rebrand_test_oak.pptx')
        else:
            prs = Presentation('rebrand_test_cfm.pptx')
        prs.save('Images/{}_{}_{}_gap.pptx'.format(Date, firm, sym))
        shapes = prs.slides[0].shapes
        pic = shapes.add_picture('Images/{}_{}_{}_gap.png'.format(Date, firm, sym), Inches(0.001), Inches(0.001),
                                 height=Inches(10.99), width=Inches(8.5))
        prs.save('Images/{}_{}_{}_gap.pptx'.format(Date, firm, sym))


def two_x_plus(RIC, Date, get_data, firm, output, fee=0.0125):
    """
    function for 2x plus protection trade
    :param RIC: <str> the Reurters identification code for the security in question
    :param Date: <str> expiration date of the desired options chain.
    :param get_data: <str> y or n only - y will pull fresh data from tr, n will use stored data
    :param output: <str> "line" - will provide terms, "trades" will provide trades, "chart" will provide fact sheet
    :param firm: <str> "L" for Lido, "O" for Oakhurst
    :return: based on output param
    """
    if get_data == 'y':  # get data if necessary
        od.get_options([RIC])
    else:  # otherwise pass
        pass
    sym = pd.read_pickle('assets/{}_sym.pkl'.format(RIC)).iloc[0]['ticker']  # import the symbol
    und = pd.read_pickle('assets/{}_cached_und.pkl'.format(sym))  ##Gets Underlying Price###
    last = und.iloc[0]['CF_LAST']  # get most recent price
    und['YIELD'] = und['YIELD'].fillna(0)  # make sure no nas exist in the div number
    QTR_div_yld = und.iloc[0]['YIELD'] / 400
    date = datetime.datetime.strptime(Date, '%Y-%m-%d').date()  # convert date to datetime
    today = datetime.datetime.now().date()  # get today's date
    time_left = date - today  # days left
    adj_time_left = time_left / datetime.timedelta(days=1)  # convert to flt
    adj_time_left_div = np.floor(time_left / datetime.timedelta(days=1) / 365 * 4)
    time_to_bnd = (adj_time_left / 365)  # convert to annual
    plug = np.round(adj_time_left_div * QTR_div_yld * last, 2)  # total divs expected
    plug = np.nan_to_num(plug, 0)  # if no divs replace with 0

    # begin options work
    # options chain set up
    chain = pd.read_pickle('assets/{}_cached_chain.pkl'.format(sym))  # read in the options chaing to pandas
    trade_chain = chain[chain['EXPIR_DATE'] == Date]  # isolate the chain we want
    trade_chain['MID'] = (trade_chain["CF_BID"] + trade_chain["CF_ASK"]) * 0.5  # calculate mid px
    sorted_trade_chain = trade_chain.sort_values(['PUTCALLIND', 'STRIKE_PRC'],
                                                 axis=0, )  # sort the chain by type and strike
    sorted_trade_chain_calls = sorted_trade_chain[sorted_trade_chain['PUTCALLIND'] == 'CALL']  # create a call df
    sorted_trade_chain_puts = sorted_trade_chain[sorted_trade_chain['PUTCALLIND'] == 'PUT ']  # create a put df

    # ATM CALL
    LCALL = sorted_trade_chain_calls.loc[
        sorted_trade_chain_calls['STRIKE_PRC'].sub(last).abs().idxmin()]  # get atm call
    # set up put spread:
    LPUT = sorted_trade_chain_puts.loc[
        sorted_trade_chain_puts['STRIKE_PRC'].sub(last * 0.9).abs().idxmin()]  # get long put
    SPUT = sorted_trade_chain_puts.loc[
        sorted_trade_chain_puts['STRIKE_PRC'].sub(last * 0.75).abs().idxmin()]  # get short Put
    spread_cst = LPUT['MID'] - SPUT["MID"] - plug  # cost of spread net of any dividend
    tot_prem = spread_cst + LCALL["MID"]
    scall_tgt = tot_prem / 2
    SCALL = sorted_trade_chain_calls.loc[
        sorted_trade_chain_calls['MID'].sub(scall_tgt).abs().idxmin()]  # get Short call

    SPUT['Trans'] = 'SPO'  # add trade type
    LPUT['Trans'] = 'BPO'  # add trade type
    LCALL['Trans'] = 'BCO'
    SCALL['Trans'] = 'SCO'  # add trade type
    Trade = pd.concat([SCALL.to_frame().T, LCALL.to_frame().T, LPUT.to_frame().T, SPUT.to_frame().T])  # create Trade df
    exp_date = pd.to_datetime(Trade['EXPIR_DATE'].iloc[0]).strftime('%y%m%d')  # get options formated date
    option_type = [s[0] for s in Trade["PUTCALLIND"].to_list()]  # isolate first letter of option type
    strikes = Trade['STRIKE_PRC'].to_list()  # # isolate strikes in list
    option_sym = ['{}{}{}{}'.format(sym, exp_date, option_type[i], int(strikes[i])) for i in
                  range(0, len(Trade))]  # create symbols
    Trade['Symbol'] = option_sym  # add to df
    trade_details = Trade.filter(['Trans', 'Symbol', 'STRIKE_PRC', 'MID'])  # create trade_det df
    upside = ((SCALL['STRIKE_PRC'] / LCALL['STRIKE_PRC']) - 1) * 2  # upside calculation
    annual_up = upside * (365 / adj_time_left)
    protection = np.abs(
        (((SPUT['STRIKE_PRC'] / last) - 1) - ((LPUT['STRIKE_PRC'] / last) - 1)))  # protection calculation
    ds_before = (last / LPUT['STRIKE_PRC']) - 1
    annual_p = protection * (365 / adj_time_left)
    months_left = np.round((adj_time_left / 365) * 12, 1)
    trade_line = pd.DataFrame({"Underlying Asset": sym, "Asset Price": last,
                               'Minimum ($)': round(last * 100, -2),
                               "Expiration Date": Trade['EXPIR_DATE'].iloc[0], "Months Left": months_left,
                               'Potential Upside (%)': '{:.2%}'.format(upside),
                               'Annual Potential Upside (%)': '{:.2%}'.format(annual_up),
                               'Downside Before Protection (%)': '{:.2%}'.format(ds_before),
                               'Protection (%)': '{:.2%}'.format(protection),
                               'Annual_Protection (%)': '{:.2%}'.format(annual_p),
                               'Spread (%)': '{:.2%}'.format(upside + protection), 'Collateral': '{}'.format(sym)},
                              index=[0])  # creates line

    chrt_rng = np.linspace(LCALL['STRIKE_PRC'] * 0.5, LCALL['STRIKE_PRC'] * 1.75, 50, dtype=float)  # set chart range
    chrt_rng = np.round(chrt_rng, 2)
    scall_ev = [np.maximum(p - SCALL['STRIKE_PRC'], 0) * -2 for p in chrt_rng]  # calc scall end val
    lcall_ev = [np.maximum(p - LCALL['STRIKE_PRC'], 0) for p in chrt_rng]
    lput_ev = [np.maximum(LPUT['STRIKE_PRC'] - p, 0) for p in chrt_rng]  # calc lcall end val
    sput_ev = [np.maximum(SPUT['STRIKE_PRC'] - p, 0) * -1 for p in chrt_rng]  # calc sput end val
    perf_df = pd.DataFrame({"{} Price".format(sym): chrt_rng, "SCALL": scall_ev, "LCALL": lcall_ev,
                            "LPUT": lput_ev, "SPUT": sput_ev,
                            "UND": chrt_rng, "DIVS": [plug for i in range(0, len(chrt_rng))]})  # Create the df
    perf_df = perf_df.set_index("{} Price".format(sym))  # Set the mkt px as the index
    perf_df['Trade'] = perf_df.sum(axis=1)  # calculate total value
    cost = (LPUT['MID'] + LCALL['MID'] + last) - (2 * SCALL['MID']) - SPUT['MID']  # trade cost
    perf_df['Trade Return'] = (perf_df['Trade'] / (cost)) - 1  # trade return
    perf_df = perf_df.sort_index(ascending=False)  # reorder in descending
    perf_df['{} Price Return'.format(sym)] = [(p / last) - 1 for p in perf_df.index]  # add underlying performance
    rets_tab = perf_df.filter(perf_df.columns[-2:]).reset_index()  # reset index
    rets_tab['Trade Return Net'] = [i - float(fee * time_to_bnd) for i in perf_df['Trade Return'].to_list()]
    fc_date = pd.to_datetime(Trade['EXPIR_DATE'].iloc[0]).strftime("%m-%d-%Y")
    upside = perf_df['Trade Return'].max()  # update upside
    trade_line['Potential Upside (%)'] = '{:.2%}'.format(upside)  # Update Upside

    if output == 'line':  # line output
        return trade_line
    elif output == 'trades':  # trades output
        und['Min_Shares'] = 100
        und['Min_Cost'] = und['Min_Shares'] * last
        print(trade_details)
        print(und)
    elif output == 'execute':
        op_exec = trade_details
        op_exec['units'] = [2, 1, 1, 1]
        und['Min_Shares'] = 100
        und_exec = und.filter(['Instrument', 'CF_LAST'])
        op_exec.to_csv('assets/temp_op.csv')
        und_exec.to_csv('assets/temp_und.csv')
    else:
        # Figure objects
        trade_title = f"{sym} 2x<br>Plus Protection"
        trade_subtitle = f"Ending {fc_date}"
        table_description = f"The investor gets up to 200% of the price appreciation on {sym} beginning at a price of" \
                            f"<br>{LCALL['STRIKE_PRC']}, capped at {SCALL['STRIKE_PRC']} and full protection between " \
                            f"{LPUT['STRIKE_PRC']} and {SPUT['STRIKE_PRC']} ({-protection:.2%})"
        chart_description = f"{sym} Price (Created on: {today} when {sym} was at {last})"
        trade_summary_data = [
            ('Structure', 'Limited 2x Upside + Protection'),
            ('Term', f"{months_left} Months"),
            ('Underlying Asset', sym),
            ('Cap', f"Capped at {upside:.2%}"),
            ('Downside Before Protection', f"{ds_before:.2%}"),
            ('Protection', f"{protection:.2%}"),
            ('Fee', f"{fee:.2%} Annually"),
        ]

        # Create the figure
        fig = plot_trade(chart_data=rets_tab, trade_title=trade_title, trade_subtitle=trade_subtitle,
                         table_description=table_description, chart_description=chart_description,
                         trade_summary_data=trade_summary_data, firm=firm, fee=fee)

        pio.write_html(fig, file='html/{}_{}_{}_2xprot.html'.format(Date, firm, sym),
                       auto_open=False)
        pio.write_image(fig, file='Images/{}_{}_{}_2xprot.png'.format(Date, firm, sym), format='png', scale=6)
        if firm == "L":
            prs = Presentation('rebrand_test.pptx')
        elif firm == "O":
            prs = Presentation('rebrand_test_oak.pptx')
        else:
            prs = Presentation('rebrand_test_cfm.pptx')
        prs.save('Images/{}_{}_{}_2xprot.pptx'.format(Date, firm, sym))
        shapes = prs.slides[0].shapes
        pic = shapes.add_picture('Images/{}_{}_{}_2xprot.png'.format(Date, firm, sym), Inches(0.001), Inches(0.001),
                                 height=Inches(10.99), width=Inches(8.5))
        prs.save('Images/{}_{}_{}_2xprot.pptx'.format(Date, firm, sym))
        return


def two_x(RIC, Date, get_data, firm, output, fee=0.0125):
    """
    function for 2x trade
    :param RIC: <str> the Reurters identification code for the security in question
    :param Date: <str> expiration date of the desired options chain.
    :param get_data: <str> y or n only - y will pull fresh data from tr, n will use stored data
    :param output: <str> "line" - will provide terms, "trades" will provide trades, "chart" will provide fact sheet
    :param firm: <str> "L" for Lido, "O" for Oakhurst
    :return: based on output param
    """
    if get_data == 'y':  # get data if necessary
        od.get_options([RIC])
    else:  # otherwise pass
        pass
    sym = pd.read_pickle('assets/{}_sym.pkl'.format(RIC)).iloc[0]['ticker']  # import the symbol
    und = pd.read_pickle('assets/{}_cached_und.pkl'.format(sym))  ##Gets Underlying Price###
    last = und.iloc[0]['CF_LAST']  # get most recent price
    und['YIELD'] = und['YIELD'].fillna(0)  # make sure no nas exist in the div number
    QTR_div_yld = und.iloc[0]['YIELD'] / 400
    date = datetime.datetime.strptime(Date, '%Y-%m-%d').date()  # convert date to datetime
    today = datetime.datetime.now().date()  # get today's date
    time_left = date - today  # days left
    adj_time_left = time_left / datetime.timedelta(days=1)  # convert to flt
    adj_time_left_div = np.floor(time_left / datetime.timedelta(days=1) / 365 * 4)
    time_to_bnd = (adj_time_left / 365)  # convert to annual
    plug = np.round(adj_time_left_div * QTR_div_yld * last, 2)  # total divs expected
    plug = np.nan_to_num(plug, 0)  # if no divs replace with 0

    # begin options work
    # options chain set up
    chain = pd.read_pickle('assets/{}_cached_chain.pkl'.format(sym))  # read in the options chaing to pandas
    trade_chain = chain[chain['EXPIR_DATE'] == Date]  # isolate the chain we want
    trade_chain['MID'] = (trade_chain["CF_BID"] + trade_chain["CF_ASK"]) * 0.5  # calculate mid px
    sorted_trade_chain = trade_chain.sort_values(['PUTCALLIND', 'STRIKE_PRC'],
                                                 axis=0, )  # sort the chain by type and strike
    sorted_trade_chain_calls = sorted_trade_chain[sorted_trade_chain['PUTCALLIND'] == 'CALL']  # create a call df

    # ATM CALL
    LCALL = sorted_trade_chain_calls.loc[
        sorted_trade_chain_calls['STRIKE_PRC'].sub(last).abs().idxmin()]  # get atm call
    # set up put spread:
    scall_tgt = LCALL["MID"] / 2
    SCALL = sorted_trade_chain_calls.loc[
        sorted_trade_chain_calls['MID'].sub(scall_tgt).abs().idxmin()]  # get Short call
    LCALL['Trans'] = 'BCO'
    SCALL['Trans'] = 'SCO'  # add trade type
    Trade = pd.concat([SCALL.to_frame().T, LCALL.to_frame().T])  # create Trade df
    exp_date = pd.to_datetime(Trade['EXPIR_DATE'].iloc[0]).strftime('%y%m%d')  # get options formated date
    option_type = [s[0] for s in Trade["PUTCALLIND"].to_list()]  # isolate first letter of option type
    strikes = Trade['STRIKE_PRC'].to_list()  # # isolate strikes in list
    option_sym = ['{}{}{}{}'.format(sym, exp_date, option_type[i], int(strikes[i])) for i in
                  range(0, len(Trade))]  # create symbols
    Trade['Symbol'] = option_sym  # add to df
    trade_details = Trade.filter(['Trans', 'Symbol', 'STRIKE_PRC', "MID"])  # create trade_det df
    upside = ((SCALL['STRIKE_PRC'] / LCALL['STRIKE_PRC']) - 1) * 2  # upside calculation
    annual_up = upside * (365 / adj_time_left)
    protection = 0  # protection calculation
    ds_before = 1
    annual_p = 1
    months_left = np.round((adj_time_left / 365) * 12, 1)
    trade_line = pd.DataFrame({"Underlying Asset": sym, "Asset Price": last,
                               'Minimum ($)': round(last * 100, -2),
                               "Expiration Date": Trade['EXPIR_DATE'].iloc[0], "Months Left": months_left,
                               'Potential Upside (%)': '{:.2%}'.format(upside),
                               'Annual Potential Upside (%)': '{:.2%}'.format(annual_up),
                               'Downside Before Protection (%)': '{:.2%}'.format(ds_before),
                               'Protection (%)': '{:.2%}'.format(protection),
                               'Annual_Protection (%)': '{:.2%}'.format(annual_p),
                               'Spread (%)': '{:.2%}'.format(upside + protection), 'Collateral': '{}'.format(sym)},
                              index=[0])  # creates line

    chrt_rng = np.linspace(LCALL['STRIKE_PRC'] * 0.5, LCALL['STRIKE_PRC'] * 1.75, 50, dtype=float)  # set chart range
    chrt_rng = np.round(chrt_rng, 2)
    scall_ev = [np.maximum(p - SCALL['STRIKE_PRC'], 0) * -2 for p in chrt_rng]  # calc scall end val
    lcall_ev = [np.maximum(p - LCALL['STRIKE_PRC'], 0) for p in chrt_rng]
    perf_df = pd.DataFrame({"{} Price".format(sym): chrt_rng, "SCALL": scall_ev, "LCALL": lcall_ev,
                            "UND": chrt_rng, "DIVS": [plug for i in range(0, len(chrt_rng))]})  # Create the df
    perf_df = perf_df.set_index("{} Price".format(sym))  # Set the mkt px as the index
    perf_df['Trade'] = perf_df.sum(axis=1)  # calculate total value
    cost = (LCALL['MID'] + last) - (2 * SCALL['MID'])  # trade cost
    perf_df['Trade Return'] = (perf_df['Trade'] / (cost)) - 1  # trade return
    perf_df = perf_df.sort_index(ascending=False)  # reorder in descending
    perf_df['{} Price Return'.format(sym)] = [(p / last) - 1 for p in perf_df.index]  # add underlying performance
    rets_tab = perf_df.filter(perf_df.columns[-2:]).reset_index()  # reset index
    rets_tab['Trade Return Net'] = [i - float(fee * time_to_bnd) for i in perf_df['Trade Return'].to_list()]
    fc_date = pd.to_datetime(Trade['EXPIR_DATE'].iloc[0]).strftime("%m-%d-%Y")
    upside = perf_df['Trade Return'].max()  # update upside
    trade_line['Potential Upside (%)'] = '{:.2%}'.format(upside)  # Update Upside

    if output == 'line':  # line output
        return trade_line
    elif output == 'trades':  # trades output
        und['Min_Shares'] = 100
        und['Min_Cost'] = und['Min_Shares'] * last
        print(trade_details)
        print(und)
    elif output == 'execute':
        op_exec = trade_details
        op_exec['units'] = [2, 1]
        und['Min_Shares'] = 100
        und_exec = und.filter(['Instrument', 'CF_LAST'])
        op_exec.to_csv('assets/temp_op.csv')
        und_exec.to_csv('assets/temp_und.csv')
    else:
        # Figure objects
        trade_title = f"{sym} 2x Trade"
        trade_subtitle = f"Ending {fc_date}"
        table_description = f"The investor gets up to 200% of the price appreciation on {sym} beginning at a price of " \
                            f"{LCALL['STRIKE_PRC']},<br>capped at {SCALL['STRIKE_PRC']}."
        chart_description = f"{sym} Price (Created on: {today} when {sym} was at {last})"
        trade_summary_data = [
            ('Structure', 'Limited 2x Upside'),
            ('Term', f"{months_left} Months"),
            ('Underlying Asset', sym),
            ('Cap', f"Capped at {upside:.2%}"),
            ('Downside Before Protection', 'No Gap'),
            ('Protection', 'No Protection'),
            ('Fee', f"{fee:.2%} Annually"),
        ]

        # Create the figure
        fig = plot_trade(chart_data=rets_tab, trade_title=trade_title, trade_subtitle=trade_subtitle,
                         table_description=table_description, chart_description=chart_description,
                         trade_summary_data=trade_summary_data, firm=firm, fee=fee)

        pio.write_html(fig, file='html/{}_{}_{}_2x.html'.format(Date, firm, sym),
                       auto_open=False)
        pio.write_image(fig, file='Images/{}_{}_{}_2x.png'.format(Date, firm, sym), format='png', scale=6)
        if firm == "L":
            prs = Presentation('rebrand_test.pptx')
        elif firm == "O":
            prs = Presentation('rebrand_test_oak.pptx')
        else:
            prs = Presentation('rebrand_test_cfm.pptx')
        prs.save('Images/{}_{}_{}_2x.pptx'.format(Date, firm, sym))
        shapes = prs.slides[0].shapes
        pic = shapes.add_picture('Images/{}_{}_{}_2x.png'.format(Date, firm, sym), Inches(0.001), Inches(0.001),
                                 height=Inches(10.99), width=Inches(8.5))
        prs.save('Images/{}_{}_{}_2x.pptx'.format(Date, firm, sym))
        return


def uncapped_plus(get_data, firm, output, fee=0.0125):
    """
    uncapped trade using corporate bond etf
    :param get_data: <str> y or n only - y will pull fresh data from tr, n will use stored data
    :param output: <str> "line" - will provide terms, "trades" will provide trades, "chart" will provide fact sheet
    :param firm: <str> "L" for Lido, "O" for Oakhurst
    :return: based on output param
    """
    if get_data == 'y':  # get data if necessary
        od.get_options(["SPY"])
        ytw = od.calc_credit('BSCP', get_data='y')

    else:  # otherwise pass
        ytw = od.calc_credit('BSCP', get_data='n')

    Date = '2025-12-19'
    sym = pd.read_pickle('assets/{}_sym.pkl'.format("SPY")).iloc[0]['ticker']  # import the symbol
    und = pd.read_pickle('assets/{}_cached_und.pkl'.format(sym))  ##Gets Underlying Price###
    bnd_dat = pd.read_pickle('assets/bond_data.pkl')  # import bond data
    last = und.iloc[0]['CF_LAST']  # get most recent price

    # Set up Bond Data

    # begin options work
    # options chain set up
    chain = pd.read_pickle('assets/{}_cached_chain.pkl'.format(sym))  # read in the options chaing to pandas
    chain['MID'] = (chain['CF_BID'] + chain['CF_ASK']) / 2
    trade_chain = chain[chain['EXPIR_DATE'] == Date]  # isolate the chain we want
    sorted_trade_chain = trade_chain.sort_values(['PUTCALLIND', 'STRIKE_PRC'],
                                                 axis=0, )  # sort the chain by type and strike
    sorted_trade_chain_calls = sorted_trade_chain[sorted_trade_chain['PUTCALLIND'] == 'CALL']  # create a call df
    sorted_trade_chain_puts = sorted_trade_chain[sorted_trade_chain['PUTCALLIND'] == 'PUT ']  # create a put df

    # ATM CALL
    LCALL = sorted_trade_chain_calls.loc[
        sorted_trade_chain_calls['STRIKE_PRC'].sub(last).abs().idxmin()]  # get atm call

    # Solve for the other two legs
    date = datetime.datetime.strptime(Date, '%Y-%m-%d').date()  # convert date to datetime
    today = datetime.datetime.now().date()  # get today's date
    time_left = date - today  # days left
    adj_time_left = time_left / datetime.timedelta(days=1)  # convert to flt
    time_to_bnd = (adj_time_left / 365)  # convert to annual
    bnd = ek.get_data("BSCP.O", fields=['CF_LAST'])[0]  ###Gets BSCN Price###
    rf = (1 + np.round(ytw, 4)) ** time_to_bnd
    bnd['1x Cost'] = (LCALL['STRIKE_PRC'] * 100) / rf
    bnd['Shares'] = (bnd['1x Cost'] / bnd['CF_LAST']).round(0)
    notional = LCALL['STRIKE_PRC']  # trade's notional is the strike of the option
    bnd['TR'] = (bnd['1x Cost'] / 100 * rf) - bnd['1x Cost'] / 100  # calculate total expected return from bond
    # begin work on put
    tgt = LCALL["MID"] - bnd['TR'].iloc[0]  # tgt put value
    SPUT = sorted_trade_chain_puts.loc[sorted_trade_chain_puts['MID'].sub(tgt).abs().idxmin()]  # get atm call
    LCALL['Trans'] = 'BCO'
    SPUT['Trans'] = 'SPO'  # add trade type
    Trade = pd.concat([LCALL.to_frame().T, SPUT.to_frame().T, ])  # create Trade df
    exp_date = pd.to_datetime(Trade['EXPIR_DATE'].iloc[0]).strftime('%y%m%d')  # get options formated date
    option_type = [s[0] for s in Trade["PUTCALLIND"].to_list()]  # isolate first letter of option type
    strikes = Trade['STRIKE_PRC'].to_list()  # # isolate strikes in list
    option_sym = ['{}{}{}{}'.format(sym, exp_date, option_type[i], int(strikes[i])) for i in
                  range(0, len(Trade))]  # create symbols
    Trade['Symbol'] = option_sym  # add to df
    trade_details = Trade.filter(['Trans', 'Symbol', 'STRIKE_PRC', 'CF_BID', 'CF_ASK'])  # create trade_det df
    upside = "Uncapped"  # upside calculation
    protection = np.abs((SPUT['STRIKE_PRC'] / LCALL['STRIKE_PRC']) - 1)  # protection calculation
    annual_p = protection * (365 / adj_time_left)
    months_left = np.round((adj_time_left / 365) * 12, 1)
    trade_line = pd.DataFrame({"Underlying Asset": sym, "Asset Price": last,
                               'Minimum ($)': round(LCALL['STRIKE_PRC'] * 100, -2),
                               "Expiration Date": Trade['EXPIR_DATE'].iloc[0], "Months Left": months_left,
                               'Potential Upside (%)': upside,
                               'Annual Potential Upside (%)': upside,
                               'Downside Before Protection (%)': '{:.2%}'.format(0),
                               'Protection (%)': '{:.2%}'.format(protection),
                               'Annual_Protection (%)': '{:.2%}'.format(annual_p),
                               'Spread (%)': 'NA', 'Collateral': 'BSCP'}, index=[0])  # creates line

    # begin work on plot data:
    chrt_rng = np.linspace(LCALL['STRIKE_PRC'] * 0.5, LCALL['STRIKE_PRC'] * 1.5, 50, dtype=float)  # set chart range
    chrt_rng = np.round(chrt_rng, 2)
    lcall_ev = [np.maximum(p - LCALL['STRIKE_PRC'], 0) for p in chrt_rng]  # calc lcall end val
    sput_ev = [np.maximum(SPUT['STRIKE_PRC'] - p, 0) * -1 for p in chrt_rng]  # calc sput end val
    perf_df = pd.DataFrame({"{} Price".format(sym): chrt_rng, "LCALL": lcall_ev, "SPUT": sput_ev,
                            "BOND": LCALL['STRIKE_PRC']})  # Create the df
    perf_df = perf_df.set_index("{} Price".format(sym))  # Set the mkt px as the index
    perf_df['Trade'] = perf_df.sum(axis=1)  # calculate total value
    cost = LCALL['MID'] - SPUT['MID'] + LCALL['STRIKE_PRC'] - bnd['TR'].iloc[0]  # trade cost
    perf_df['Trade Return'] = (perf_df['Trade'] / (cost)) - 1  # trade return
    perf_df = perf_df.sort_index(ascending=False)  # reorder in descending
    perf_df['{} Price Return'.format(sym)] = [(p / last) - 1 for p in perf_df.index]  # add underlying performance
    rets_tab = perf_df.filter(perf_df.columns[-2:]).reset_index()  # reset index
    rets_tab['Trade Return Net'] = [i - float(fee * time_to_bnd) for i in perf_df['Trade Return'].to_list()]
    fc_date = pd.to_datetime(Trade['EXPIR_DATE'].iloc[0]).strftime("%m-%d-%Y")

    if output == 'line':  # line output
        return trade_line
    elif output == 'trades':  # trades output
        print(trade_details)
        print(bnd.filter(['Instrument', "CF_LAST", "Shares", '1x Cost', ]))
    elif output == 'execute':
        op_exec = trade_details
        op_exec['MID'] = (op_exec['CF_BID'] + op_exec['CF_ASK']) / 2
        op_exec = op_exec.drop(['CF_BID', 'CF_ASK'], axis=1)
        op_exec['units'] = [1 for i in range(0, len(op_exec))]
        bnd_exec = bnd.filter(['Instrument', 'Shares', 'CF_LAST'])
        bnd_exec = bnd_exec.rename(columns={'Shares': 'units'})
        op_exec.to_csv('assets/temp_op.csv')
        bnd_exec.to_csv('assets/temp_bnd.csv')
    else:
        # Figure objects
        trade_title = "Uncapped Plus<br>Trade"
        trade_subtitle = f"Ending {fc_date}"
        table_description = f"The investor gets about 100% of the price appreciation on {sym} beginning at a<br>" \
                            f"price of {int(LCALL['STRIKE_PRC'])} and full protection between " \
                            f"{int(LCALL['STRIKE_PRC'])} and {int(SPUT['STRIKE_PRC'])} ({-protection:.2%})"
        chart_description = f"{sym} Price (Created on: {today} when {sym} was at {last})"
        trade_summary_data = [
            ('Structure', 'Unlimited Upside + Protection'),
            ('Term', f"{months_left} Months"),
            ('Underlying Asset', sym),
            ('Cap', 'Uncapped'),
            ('Downside Before Protection', 'No Gap'),
            ('Protection', f"{protection:.2%}"),
            ('Fee', f"{fee:.2%} Annually"),
        ]

        # Create the figure
        fig = plot_trade(chart_data=rets_tab, trade_title=trade_title, trade_subtitle=trade_subtitle,
                         table_description=table_description, chart_description=chart_description,
                         trade_summary_data=trade_summary_data, firm=firm, fee=fee)

        pio.write_html(fig, file='html/{}_{}_{}_unc.html'.format(Date, firm, sym),
                       auto_open=False)
        pio.write_image(fig, file='Images/{}_{}_{}_uncapped.png'.format(Date, firm, sym), format='png', scale=6)
        if firm == "L":
            prs = Presentation('rebrand_test.pptx')
        elif firm == "O":
            prs = Presentation('rebrand_test_oak.pptx')
        else:
            prs = Presentation('rebrand_test_cfm.pptx')
        prs.save('Images/{}_{}_{}_uncapped.pptx'.format(Date, firm, sym))
        shapes = prs.slides[0].shapes
        pic = shapes.add_picture('Images/{}_{}_{}_uncapped.png'.format(Date, firm, sym), Inches(0.001), Inches(0.001),
                                 height=Inches(10.99), width=Inches(8.5))
        prs.save('Images/{}_{}_{}_uncapped.pptx'.format(Date, firm, sym))
        return


def uncapped1_5x(get_data, firm, output, fee=0.0125):
    """
    uncapped trade using corporate bond etf
    :param get_data: <str> y or n only - y will pull fresh data from tr, n will use stored data
    :param output: <str> "line" - will provide terms, "trades" will provide trades, "chart" will provide fact sheet
    :param firm: <str> "L" for Lido, "O" for Oakhurst
    :return: based on output param
    """
    if get_data == 'y':  # get data if necessary
        od.get_options(["SPY"])
        ytw = od.calc_credit('BSCP', get_data='y')

    else:  # otherwise pass
        ytw = od.calc_credit('BSCP', get_data='n')

    Date = '2025-12-19'
    sym = pd.read_pickle('assets/{}_sym.pkl'.format("SPY")).iloc[0]['ticker']  # import the symbol
    und = pd.read_pickle('assets/{}_cached_und.pkl'.format(sym))  ##Gets Underlying Price###
    bnd_dat = pd.read_pickle('assets/bond_data.pkl')  # import bond data
    last = und.iloc[0]['CF_LAST']  # get most recent price

    # begin options work
    # options chain set up
    chain = pd.read_pickle('assets/{}_cached_chain.pkl'.format(sym))  # read in the options chaing to pandas
    chain['MID'] = (chain['CF_BID'] + chain['CF_ASK']) / 2
    trade_chain = chain[chain['EXPIR_DATE'] == Date]  # isolate the chain we want
    sorted_trade_chain = trade_chain.sort_values(['PUTCALLIND', 'STRIKE_PRC'],
                                                 axis=0, )  # sort the chain by type and strike
    sorted_trade_chain_calls = sorted_trade_chain[sorted_trade_chain['PUTCALLIND'] == 'CALL']  # create a call df
    sorted_trade_chain_puts = sorted_trade_chain[sorted_trade_chain['PUTCALLIND'] == 'PUT ']  # create a put df

    # ATM CALL
    LCALL = sorted_trade_chain_calls.loc[
        sorted_trade_chain_calls['STRIKE_PRC'].sub(last).abs().idxmin()]  # get atm call
    SPUT = sorted_trade_chain_puts.loc[
        sorted_trade_chain_puts['STRIKE_PRC'].sub(LCALL['STRIKE_PRC']).abs().idxmin()]  # get atm PUTS
    # Solve for the other two legs
    date = datetime.datetime.strptime(Date, '%Y-%m-%d').date()  # convert date to datetime
    today = datetime.datetime.now().date()  # get today's date
    time_left = date - today  # days left
    adj_time_left = time_left / datetime.timedelta(days=1)  # convert to flt
    time_to_bnd = (adj_time_left / 365)  # convert to annual
    bnd = ek.get_data("BSCP.O", fields=['CF_LAST'])[0]  ###Gets BSCN Price###
    rf = (1 + np.round(ytw, 4)) ** time_to_bnd
    bnd['1x Cost'] = (LCALL['STRIKE_PRC'] * 200) / rf
    bnd['Shares'] = (bnd['1x Cost'] / bnd['CF_LAST']).round(0)
    notional = LCALL['STRIKE_PRC'] * 2  # trade's notional is the strike of the option
    bnd['TR'] = (bnd['1x Cost'] / 100 * rf) - bnd['1x Cost'] / 100  # calculate total expected return from bond
    LCALL['Trans'] = 'BCO'
    SPUT['Trans'] = 'SPO'  # add trade type
    Trade = pd.concat([LCALL.to_frame().T, SPUT.to_frame().T, ])  # create Trade df
    exp_date = pd.to_datetime(Trade['EXPIR_DATE'].iloc[0]).strftime('%y%m%d')  # get options formated date
    option_type = [s[0] for s in Trade["PUTCALLIND"].to_list()]  # isolate first letter of option type
    strikes = Trade['STRIKE_PRC'].to_list()  # # isolate strikes in list
    option_sym = ['{}{}{}{}'.format(sym, exp_date, option_type[i], int(strikes[i])) for i in
                  range(0, len(Trade))]  # create symbols
    Trade['Symbol'] = option_sym  # add to df
    trade_details = Trade.filter(['Trans', 'Symbol', 'STRIKE_PRC', "MID"])  # create trade_det df
    trade_details['No. Contracts'] = [3, 2]
    upside = "Uncapped"  # upside calculation
    protection = 0  # protection calculation
    annual_p = protection * (365 / adj_time_left)
    months_left = np.round((adj_time_left / 365) * 12, 1)
    trade_line = pd.DataFrame({"Underlying Asset": sym, "Asset Price": last,
                               'Minimum ($)': round(LCALL['STRIKE_PRC'] * 200, -2),
                               "Expiration Date": Trade['EXPIR_DATE'].iloc[0], "Months Left": months_left,
                               'Potential Upside (%)': upside,
                               'Annual Potential Upside (%)': upside,
                               'Downside Before Protection (%)': '{:.2%}'.format(0),
                               'Protection (%)': '{:.2%}'.format(protection),
                               'Annual_Protection (%)': '{:.2%}'.format(annual_p),
                               'Spread (%)': 'NA', 'Collateral': 'BSCP'}, index=[0])  # creates line

    # begin work on plot data:
    chrt_rng = np.linspace(LCALL['STRIKE_PRC'] * 0.5, LCALL['STRIKE_PRC'] * 1.5, 50, dtype=float)  # set chart range
    chrt_rng = np.round(chrt_rng, 2)
    lcall_ev = [np.maximum(p - LCALL['STRIKE_PRC'], 0) * 3 for p in chrt_rng]  # calc lcall end val
    sput_ev = [np.maximum(SPUT['STRIKE_PRC'] - p, 0) * -2 for p in chrt_rng]  # calc sput end val
    perf_df = pd.DataFrame({"{} Price".format(sym): chrt_rng, "LCALL": lcall_ev, "SPUT": sput_ev,
                            "BOND": LCALL['STRIKE_PRC'] * 2})  # Create the df
    perf_df = perf_df.set_index("{} Price".format(sym))  # Set the mkt px as the index
    perf_df['Trade'] = perf_df.sum(axis=1)  # calculate total value
    cost = LCALL['MID'] * 3 - SPUT['MID'] * 2 + LCALL['STRIKE_PRC'] * 2 - bnd['TR'].iloc[0]  # trade cost
    perf_df['Trade Return'] = (perf_df['Trade'] / (cost)) - 1  # trade return
    perf_df = perf_df.sort_index(ascending=False)  # reorder in descending
    perf_df['{} Price Return'.format(sym)] = [(p / last) - 1 for p in perf_df.index]  # add underlying performance
    rets_tab = perf_df.filter(perf_df.columns[-2:]).reset_index()  # reset index
    rets_tab['Trade Return Net'] = [i - float(fee * time_to_bnd) for i in perf_df['Trade Return'].to_list()]
    fc_date = pd.to_datetime(Trade['EXPIR_DATE'].iloc[0]).strftime("%m-%d-%Y")

    if output == 'line':  # line output
        return trade_line
    elif output == 'trades':  # trades output
        print(trade_details)
        print(bnd.filter(['Instrument', "CF_LAST", "Shares", '1x Cost', ]))
    elif output == 'execute':
        op_exec = trade_details
        op_exec['units'] = [1.5, 1]
        bnd_exec = bnd.filter(['Instrument', 'Shares', 'CF_LAST'])
        bnd_exec = bnd_exec.rename(columns={'Shares': 'units'})
        op_exec.to_csv('assets/temp_op.csv')
        bnd_exec.to_csv('assets/temp_bnd.csv')
    else:
        # Figure objects
        trade_title = "Uncapped 1.5x<br>Trade"
        trade_subtitle = f"Ending {fc_date}"
        table_description = f"The investor gets <b>UP TO 150%</b> of the price appreciation on {sym} beginning at a " \
                            f"price<br>of {int(LCALL['STRIKE_PRC'])}"
        chart_description = f"{sym} Price (Created on: {today} when {sym} was at {last})"
        trade_summary_data = [
            ('Structure', '1.5x levered upside with unlevered downside'),
            ('Term', f"{months_left} Months"),
            ('Underlying Asset', sym),
            ('Cap', 'Uncapped'),
            ('Downside Before Protection', 'No Gap'),
            ('Protection', 'No Protection'),
            ('Fee', f"{fee:.2%} Annually"),
        ]

        # Create the figure
        fig = plot_trade(chart_data=rets_tab, trade_title=trade_title, trade_subtitle=trade_subtitle,
                         table_description=table_description, chart_description=chart_description,
                         trade_summary_data=trade_summary_data, firm=firm, fee=fee)

        pio.write_html(fig, file='html/{}_{}_{}_unc15.html'.format(Date, firm, sym),
                       auto_open=False)
        pio.write_image(fig, file='Images/{}_{}_{}_uncapped15.png'.format(Date, firm, sym), format='png', scale=6)
        if firm == "L":
            prs = Presentation('rebrand_test.pptx')
        elif firm == "O":
            prs = Presentation('rebrand_test_oak.pptx')
        else:
            prs = Presentation('rebrand_test_cfm.pptx')
        prs.save('Images/{}_{}_{}_uncapped15.pptx'.format(Date, firm, sym))
        shapes = prs.slides[0].shapes
        pic = shapes.add_picture('Images/{}_{}_{}_uncapped15.png'.format(Date, firm, sym), Inches(0.001), Inches(0.001),
                                 height=Inches(10.99), width=Inches(8.5))
        prs.save('Images/{}_{}_{}_uncapped15.pptx'.format(Date, firm, sym))
        return


def uncapped1_4x(get_data, firm, output, fee=0.0125):
    """
    uncapped trade using corporate bond etf
    :param get_data: <str> y or n only - y will pull fresh data from tr, n will use stored data
    :param output: <str> "line" - will provide terms, "trades" will provide trades, "chart" will provide fact sheet
    :param firm: <str> "L" for Lido, "O" for Oakhurst
    :return: based on output param
    """
    if get_data == 'y':  # get data if necessary
        od.get_options(["SPY"])
        ytw = od.calc_credit('BSCP', get_data='y')

    else:  # otherwise pass
        ytw = od.calc_credit('BSCP', get_data='n')

    Date = '2025-12-19'
    sym = pd.read_pickle('assets/{}_sym.pkl'.format("SPY")).iloc[0]['ticker']  # import the symbol
    und = pd.read_pickle('assets/{}_cached_und.pkl'.format(sym))  ##Gets Underlying Price###
    bnd_dat = pd.read_pickle('assets/bond_data.pkl')  # import bond data
    last = und.iloc[0]['CF_LAST']  # get most recent price

    # begin options work
    # options chain set up
    chain = pd.read_pickle('assets/{}_cached_chain.pkl'.format(sym))  # read in the options chaing to pandas
    chain['MID'] = (chain['CF_BID'] + chain['CF_ASK']) / 2
    trade_chain = chain[chain['EXPIR_DATE'] == Date]  # isolate the chain we want
    sorted_trade_chain = trade_chain.sort_values(['PUTCALLIND', 'STRIKE_PRC'],
                                                 axis=0, )  # sort the chain by type and strike
    sorted_trade_chain_calls = sorted_trade_chain[sorted_trade_chain['PUTCALLIND'] == 'CALL']  # create a call df
    sorted_trade_chain_puts = sorted_trade_chain[sorted_trade_chain['PUTCALLIND'] == 'PUT ']  # create a put df

    # ATM CALL
    LCALL = sorted_trade_chain_calls.loc[
        sorted_trade_chain_calls['STRIKE_PRC'].sub(last).abs().idxmin()]  # get atm call
    SPUT = sorted_trade_chain_puts.loc[
        sorted_trade_chain_puts['STRIKE_PRC'].sub(LCALL['STRIKE_PRC']).abs().idxmin()]  # get atm PUTS
    # Solve for the other two legs
    date = datetime.datetime.strptime(Date, '%Y-%m-%d').date()  # convert date to datetime
    today = datetime.datetime.now().date()  # get today's date
    time_left = date - today  # days left
    adj_time_left = time_left / datetime.timedelta(days=1)  # convert to flt
    time_to_bnd = (adj_time_left / 365)  # convert to annual
    bnd = ek.get_data("BSCP.O", fields=['CF_LAST'])[0]  ###Gets BSCN Price###
    rf = (1 + np.round(ytw, 4)) ** time_to_bnd
    bnd['1x Cost'] = (LCALL['STRIKE_PRC'] * 500) / rf
    bnd['Shares'] = (bnd['1x Cost'] / bnd['CF_LAST']).round(0)
    notional = LCALL['STRIKE_PRC'] * 5  # trade's notional is the strike of the option
    bnd['TR'] = (bnd['1x Cost'] / 100 * rf) - bnd['1x Cost'] / 100  # calculate total expected return from bond
    LCALL['Trans'] = 'BCO'
    SPUT['Trans'] = 'SPO'  # add trade type
    Trade = pd.concat([LCALL.to_frame().T, SPUT.to_frame().T, ])  # create Trade df
    exp_date = pd.to_datetime(Trade['EXPIR_DATE'].iloc[0]).strftime('%y%m%d')  # get options formated date
    option_type = [s[0] for s in Trade["PUTCALLIND"].to_list()]  # isolate first letter of option type
    strikes = Trade['STRIKE_PRC'].to_list()  # # isolate strikes in list
    option_sym = ['{}{}{}{}'.format(sym, exp_date, option_type[i], int(strikes[i])) for i in
                  range(0, len(Trade))]  # create symbols
    Trade['Symbol'] = option_sym  # add to df
    trade_details = Trade.filter(['Trans', 'Symbol', 'STRIKE_PRC', "MID"])  # create trade_det df
    trade_details['No. Contracts'] = [7, 5]
    upside = "Uncapped"  # upside calculation
    protection = 0  # protection calculation
    annual_p = protection * (365 / adj_time_left)
    months_left = np.round((adj_time_left / 365) * 12, 1)
    trade_line = pd.DataFrame({"Underlying Asset": sym, "Asset Price": last,
                               'Minimum ($)': round(LCALL['STRIKE_PRC'] * 500, -2),
                               "Expiration Date": Trade['EXPIR_DATE'].iloc[0], "Months Left": months_left,
                               'Potential Upside (%)': upside,
                               'Annual Potential Upside (%)': upside,
                               'Downside Before Protection (%)': '{:.2%}'.format(0),
                               'Protection (%)': '{:.2%}'.format(protection),
                               'Annual_Protection (%)': '{:.2%}'.format(annual_p),
                               'Spread (%)': 'NA', 'Collateral': 'BSCP'}, index=[0])  # creates line

    # begin work on plot data:
    chrt_rng = np.linspace(LCALL['STRIKE_PRC'] * 0.5, LCALL['STRIKE_PRC'] * 1.5, 50, dtype=float)  # set chart range
    chrt_rng = np.round(chrt_rng, 2)
    lcall_ev = [np.maximum(p - LCALL['STRIKE_PRC'], 0) * 7 for p in chrt_rng]  # calc lcall end val
    sput_ev = [np.maximum(SPUT['STRIKE_PRC'] - p, 0) * -5 for p in chrt_rng]  # calc sput end val
    perf_df = pd.DataFrame({"{} Price".format(sym): chrt_rng, "LCALL": lcall_ev, "SPUT": sput_ev,
                            "BOND": LCALL['STRIKE_PRC'] * 5})  # Create the df
    perf_df = perf_df.set_index("{} Price".format(sym))  # Set the mkt px as the index
    perf_df['Trade'] = perf_df.sum(axis=1)  # calculate total value
    cost = LCALL['MID'] * 7 - SPUT['MID'] * 5 + LCALL['STRIKE_PRC'] * 5 - bnd['TR'].iloc[0]  # trade cost
    perf_df['Trade Return'] = (perf_df['Trade'] / (cost)) - 1  # trade return
    perf_df = perf_df.sort_index(ascending=False)  # reorder in descending
    perf_df['{} Price Return'.format(sym)] = [(p / last) - 1 for p in perf_df.index]  # add underlying performance
    rets_tab = perf_df.filter(perf_df.columns[-2:]).reset_index()  # reset index
    rets_tab['Trade Return Net'] = [i - float(fee * time_to_bnd) for i in perf_df['Trade Return'].to_list()]
    fc_date = pd.to_datetime(Trade['EXPIR_DATE'].iloc[0]).strftime("%m-%d-%Y")

    if output == 'line':  # line output
        return trade_line
    elif output == 'trades':  # trades output
        print(trade_details)
        print(bnd.filter(['Instrument', "CF_LAST", "Shares", '1x Cost', ]))
    elif output == 'execute':
        op_exec = trade_details
        op_exec['units'] = [1.4, 1]
        bnd_exec = bnd.filter(['Instrument', 'Shares', 'CF_LAST'])
        bnd_exec = bnd_exec.rename(columns={'Shares': 'units'})
        op_exec.to_csv('assets/temp_op.csv')
        bnd_exec.to_csv('assets/temp_bnd.csv')
    else:
        # Figure objects
        trade_title = "Uncapped 1.4x<br>Trade"
        trade_subtitle = f"Ending {fc_date}"
        table_description = f"The investor gets <b>UP TO 140%</b> of the price appreciation on {sym} beginning at a " \
                            f"price<br>of {int(LCALL['STRIKE_PRC'])}"
        chart_description = f"{sym} Price (Created on: {today} when {sym} was at {last})"
        trade_summary_data = [
            ('Structure', '1.4x levered upside with unlevered downside'),
            ('Term', f"{months_left} Months"),
            ('Underlying Asset', sym),
            ('Cap', 'Uncapped'),
            ('Downside Before Protection', 'No Gap'),
            ('Protection', 'No Protection'),
            ('Fee', f"{fee:.2%} Annually"),
        ]

        # Create the figure
        fig = plot_trade(chart_data=rets_tab, trade_title=trade_title, trade_subtitle=trade_subtitle,
                         table_description=table_description, chart_description=chart_description,
                         trade_summary_data=trade_summary_data, firm=firm, fee=fee)

        pio.write_html(fig, file='html/{}_{}_{}_unc14.html'.format(Date, firm, sym),
                       auto_open=False)
        pio.write_image(fig, file='Images/{}_{}_{}_uncapped14.png'.format(Date, firm, sym), format='png', scale=6)
        if firm == "L":
            prs = Presentation('rebrand_test.pptx')
        elif firm == "O":
            prs = Presentation('rebrand_test_oak.pptx')
        else:
            prs = Presentation('rebrand_test_cfm.pptx')
        prs.save('Images/{}_{}_{}_uncapped14.pptx'.format(Date, firm, sym))
        shapes = prs.slides[0].shapes
        pic = shapes.add_picture('Images/{}_{}_{}_uncapped14.png'.format(Date, firm, sym), Inches(0.001), Inches(0.001),
                                 height=Inches(10.99), width=Inches(8.5))
        prs.save('Images/{}_{}_{}_uncapped14.pptx'.format(Date, firm, sym))
        return


def spy_killer120(get_data, firm, output, fee=0.0125):
    """
    uncapped trade using corporate bond etf
    :param get_data: <str> y or n only - y will pull fresh data from tr, n will use stored data
    :param output: <str> "line" - will provide terms, "trades" will provide trades, "chart" will provide fact sheet
    :param firm: <str> "L" for Lido, "O" for Oakhurst
    :return: based on output param
    """
    if get_data == 'y':  # get data if necessary
        od.get_options(["SPY"])
        ytw = od.calc_credit('BSCP', get_data='y')
    else:  # otherwise pass
        ytw = od.calc_credit('BSCP', get_data='n')

    Date = '2025-12-19'
    sym = pd.read_pickle('assets/{}_sym.pkl'.format("SPY")).iloc[0]['ticker']  # import the symbol
    und = pd.read_pickle('assets/{}_cached_und.pkl'.format(sym))  ##Gets Underlying Price###
    bnd_dat = pd.read_pickle('assets/bond_data.pkl')  # import bond data
    last = und.iloc[0]['CF_LAST']  # get most recent price

    # options chain set up
    chain = pd.read_pickle('assets/{}_cached_chain.pkl'.format(sym))  # read in the options chaing to pandas
    chain['MID'] = (chain['CF_BID'] + chain['CF_ASK']) / 2
    trade_chain = chain[chain['EXPIR_DATE'] == Date]  # isolate the chain we want
    trade_chain['MID'] = (trade_chain['CF_BID'] + trade_chain['CF_ASK']) / 2
    sorted_trade_chain = trade_chain.sort_values(['PUTCALLIND', 'STRIKE_PRC'],
                                                 axis=0, )  # sort the chain by type and strike
    sorted_trade_chain_calls = sorted_trade_chain[sorted_trade_chain['PUTCALLIND'] == 'CALL']  # create a call df
    sorted_trade_chain_puts = sorted_trade_chain[sorted_trade_chain['PUTCALLIND'] == 'PUT ']  # create a put df

    # ATM CALL
    LCALL = sorted_trade_chain_calls.loc[
        sorted_trade_chain_calls['STRIKE_PRC'].sub(last).abs().idxmin()]  # get atm call
    SPUT1 = sorted_trade_chain_puts.loc[
        sorted_trade_chain_puts['STRIKE_PRC'].sub(LCALL['STRIKE_PRC']).abs().idxmin()]  # get atm put
    LPUT1 = sorted_trade_chain_puts.loc[
        sorted_trade_chain_puts['STRIKE_PRC'].sub(LCALL['STRIKE_PRC'] * 0.95).abs().idxmin()]  # get 7% OTM Put
    date = datetime.datetime.strptime(Date, '%Y-%m-%d').date()  # convert date to datetime
    today = datetime.datetime.now().date()  # get today's date
    time_left = date - today  # days left
    adj_time_left = time_left / datetime.timedelta(days=1)  # convert to flt
    time_to_bnd = (adj_time_left / 365)  # convert to annual
    bnd = ek.get_data("BSCP.O", fields=['CF_LAST'])[0]  ###Gets BSCN Price###
    rf = (1 + np.round(ytw, 4)) ** time_to_bnd
    bnd['1x Cost'] = (LCALL['STRIKE_PRC'] * 500) / rf
    bnd['Shares'] = (bnd['1x Cost'] / bnd['CF_LAST']).round(0)
    notional = LCALL['STRIKE_PRC'] * 5  # trade's notional is the strike of the option
    bnd['TR'] = (bnd['1x Cost'] * rf) / 100 - bnd['1x Cost'] / 100  # calculate total expected return from bond

    # solve for long_put2
    linsp = np.arange(start=(int(LCALL['STRIKE_PRC'] * 1.5)), stop=0, step=-1)
    lcall_atm = [np.maximum(i - LCALL['STRIKE_PRC'], 0) * 6 for i in linsp]  # ATM CALL PL CALC
    sput_atm1 = [np.maximum(SPUT1['STRIKE_PRC'] - i, 0) * -5 for i in linsp]  # ATM PUT PL CALC
    lput_otm1 = [np.maximum(LPUT1['STRIKE_PRC'] - i, 0) * 5 for i in linsp]  # OTM PUT PL CALC
    op_df = pd.DataFrame({"LCALL_ATM": lcall_atm, "SPUT_ATM": sput_atm1,
                          "LPUT_OTM1": lput_otm1,
                          "END_FI": [LCALL['STRIKE_PRC'] * 5 for i in
                                     linsp]})  # this is the df to use in solving for otm put
    op_df['Temp_val'] = op_df.sum(axis=1)
    op_df['mkt'] = linsp
    temp_cost = LCALL['MID'] * 6 - SPUT1["MID"] * 5 + LPUT1["MID"] * 5 + (
            notional - bnd['TR'].iloc[0])  # this is the current cost
    test_range = np.arange(LPUT1['STRIKE_PRC'] - 5, LPUT1['STRIKE_PRC'] * 0.5, -5,
                           dtype=int)  # this is the range of options to test
    put_spread_df = pd.DataFrame()  # store df
    for p in test_range:
        sp = sorted_trade_chain_puts.loc[sorted_trade_chain_puts['STRIKE_PRC'].sub(p).abs().idxmin()]  # get long put
        for q in test_range:  # loop through again
            if q >= p:  # ignore any strike above or equal to the one we picked
                pass
            else:
                lp = sorted_trade_chain_puts.loc[
                    sorted_trade_chain_puts['STRIKE_PRC'].sub(q).abs().idxmin()]  # get short put
                shrt_sr = sp['STRIKE_PRC']  # short Strike
                long_str = lp['STRIKE_PRC']  # long Strike
                spr_cst = (sp['MID'] * -20) + (lp['MID'] * 15)
                spr = pd.DataFrame({"SPUT": shrt_sr, "LPUT": long_str, "tot_cst": spr_cst}, index=[0])
                put_spread_df = pd.concat([put_spread_df, spr])  # put spread
                tgt_val = SPUT1['STRIKE_PRC'] * 5  # bond val
    put_spread_df['tgt_val'] = tgt_val  # add to df
    put_spread_df['temp_cst'] = temp_cost
    put_spread_df['adj_cost'] = put_spread_df['temp_cst'] + put_spread_df['tot_cst']
    put_spread_df['unit_delta'] = np.abs(put_spread_df['adj_cost'] - put_spread_df['tgt_val'])
    put_spread_df = put_spread_df.sort_values('unit_delta')
    put_spread_df['prem'] = put_spread_df['unit_delta'] / put_spread_df['tgt_val']
    # put_spread_df = put_spread_df[put_spread_df['prem'] <= 0.005]
    put_spread_df['downside_prot'] = \
        np.abs(((put_spread_df["SPUT"] / LCALL['STRIKE_PRC']) - 1) - ((LPUT1['STRIKE_PRC'] / LCALL['STRIKE_PRC']) - 1))
    put_spread_df = put_spread_df[put_spread_df['downside_prot'] >= 0.15]
    put_spread_df = put_spread_df[put_spread_df['downside_prot'] <= 0.22]
    best = 0
    metric = None
    op_df['mkt_perf'] = (op_df['mkt'] / LCALL['STRIKE_PRC']) - 1
    for i in range(0, len(put_spread_df)):
        test_df = op_df.copy()
        spread = put_spread_df.iloc[i]
        test_df['SPUT_OTM'] = [np.maximum(spread["SPUT"] - j, 0) * -20 for j in linsp]
        test_df['LPUT_OTM2'] = [np.maximum(spread["LPUT"] - j, 0) * 15 for j in linsp]
        test_df['Temp_val'] = test_df['Temp_val'] + test_df['SPUT_OTM'] + test_df['LPUT_OTM2']
        inv_cap = temp_cost + spread['tot_cst']
        test_df['perf_delta'] = (test_df['Temp_val'] / inv_cap) - 1
        if metric == None:
            metric = test_df['Temp_val'].iloc[-1]
            if metric >= 0:
                best = i
        elif test_df['Temp_val'].iloc[-1] >= 0:
            if test_df['Temp_val'].iloc[-1] < metric:
                metric = test_df['Temp_val'].iloc[-1]
                best = i
    spr = put_spread_df.iloc[best]
    op_df["LCALL_ATM"] = [np.maximum(i - LCALL['STRIKE_PRC'], 0) * 6 for i in op_df['mkt']]
    op_df["SPUT_ATM"] = [np.maximum(SPUT1['STRIKE_PRC'] - i, 0) * -5 for i in op_df['mkt']]  # ATM PUT PL CALC
    op_df["LPUT_OTM1"] = [np.maximum(LPUT1['STRIKE_PRC'] - i, 0) * 5 for i in op_df['mkt']]  # OTM PUT PL CALC
    op_df['SPUT_OTM'] = [np.maximum(spr["SPUT"] - j, 0) * -20 for j in op_df['mkt']]
    op_df['LPUT_OTM2'] = [np.maximum(spr["LPUT"] - j, 0) * 15 for j in op_df['mkt']]
    op_df['Temp_val'] = op_df["LCALL_ATM"] + op_df["SPUT_ATM"] + op_df["LPUT_OTM1"] + op_df['SPUT_OTM'] + op_df[
        'LPUT_OTM2'] + op_df['END_FI']
    inv_cap = temp_cost + spr['tot_cst']
    op_df['perfomance'] = (op_df['Temp_val'] / inv_cap) - 1
    perf_df = op_df.filter(['perfomance', 'mkt_perf', 'mkt'])
    SPUT_OTM = sorted_trade_chain_puts.loc[sorted_trade_chain_puts['STRIKE_PRC'].sub(spr['SPUT']).abs().idxmin()]
    LPUT2_OTM = sorted_trade_chain_puts.loc[sorted_trade_chain_puts['STRIKE_PRC'].sub(spr['LPUT']).abs().idxmin()]
    LCALL['Trans'] = 'BCO'
    SPUT1['Trans'] = 'SPO'  # add trade type
    LPUT1['Trans'] = "BPO"
    SPUT_OTM['Trans'] = "SPO"
    LPUT2_OTM['Trans'] = "BPO"
    Trade = pd.concat([LCALL.to_frame().T, SPUT1.to_frame().T, LPUT1.to_frame().T,
                       SPUT_OTM.to_frame().T, LPUT2_OTM.to_frame().T])  # create Trade df
    exp_date = pd.to_datetime(Trade['EXPIR_DATE'].iloc[0]).strftime('%y%m%d')  # get options formated date
    option_type = [s[0] for s in Trade["PUTCALLIND"].to_list()]  # isolate first letter of option type
    strikes = Trade['STRIKE_PRC'].to_list()  # # isolate strikes in list
    option_sym = ['{}{}{}{}'.format(sym, exp_date, option_type[i], int(strikes[i])) for i in
                  range(0, len(Trade))]  # create symbols
    Trade['Symbol'] = option_sym  # add to df
    trade_details = Trade.filter(['Trans', 'Symbol', 'STRIKE_PRC', "MID"])  # create trade_det df
    trade_details['No. Contracts'] = [6, 5, 5, 20, 15]
    upside = "Uncapped - Levered Upside"  # upside calculation
    gap = (LPUT1['STRIKE_PRC'] / LCALL['STRIKE_PRC']) - 1
    protection = ((SPUT_OTM['STRIKE_PRC'] / LCALL['STRIKE_PRC']) - 1) - gap  # protection calculation
    annual_p = protection * (365 / adj_time_left)
    months_left = np.round((adj_time_left / 365) * 12, 1)
    trade_line = pd.DataFrame({"Underlying Asset": sym, "Asset Price": last,
                               'Minimum ($)': round(LCALL['STRIKE_PRC'] * 500, -2),
                               "Expiration Date": Trade['EXPIR_DATE'].iloc[0], "Months Left": months_left,
                               'Potential Upside (%)': upside,
                               'Annual Potential Upside (%)': upside,
                               'Downside Before Protection (%)': '{:.2%}'.format(gap),
                               'Protection (%)': '{:.2%}'.format(protection),
                               'Annual_Protection (%)': '{:.2%}'.format(annual_p),
                               'Spread (%)': 'NA', 'Collateral': 'BSCP'}, index=[0])  # creates line
    perf_df = perf_df.rename(columns={"perfomance": "Trade Return",
                                      "mkt_perf": '{} Price Return'.format(sym)})
    rets_tab = perf_df.copy()
    rets_tab = rets_tab.filter(['mkt', 'Trade Return', '{} Price Return'.format(sym)])
    rets_tab['Trade Return Net'] = [i - float(fee * time_to_bnd) for i in perf_df['Trade Return'].to_list()]
    fc_date = pd.to_datetime(Trade['EXPIR_DATE'].iloc[0]).strftime("%m-%d-%Y")

    if output == 'line':  # line output
        return trade_line
    elif output == 'trades':  # trades output
        print(trade_details)
        print(bnd.filter(['Instrument', "CF_LAST", "Shares", '1x Cost', ]))
    elif output == 'execute':
        op_exec = trade_details.filter(['Trans', 'Symbol', 'STRIKE_PRC', 'No. Contracts', 'MID'])
        op_exec = op_exec.rename(columns={'No. Contracts': 'units'})
        bnd_exec = bnd.filter(['Instrument', 'Shares', 'CF_LAST'])
        bnd_exec = bnd_exec.rename(columns={'Shares': 'units'})
        op_exec.to_csv('assets/temp_op.csv')
        bnd_exec.to_csv('assets/temp_bnd.csv')
    else:
        # Figure objects
        trade_title = "1.2x SPY Killer<br> Trade"
        trade_subtitle = f"Ending {fc_date}"
        table_description = f"The investor gets UP TO 120% of the price appreciation on {sym} beginning at a price " \
                            f"of<br>{int(LCALL['STRIKE_PRC'])}"
        chart_description = f"{sym} Price (Created on: {today} when {sym} was at {last})"
        trade_summary_data = [
            ('Structure', 'Unlimited 1.2x levered upside with a deferred soft buffer'),
            ('Term', f"{months_left} Months"),
            ('Underlying Asset', sym),
            ('Cap', 'Uncapped 1.2x'),
            ('Downside Before Protection', f"{gap:.2%}"),
            ('Protection', f"{protection:.2%} (Soft Buffer)"),
            ('Fee', f"{fee:.2%} Annually"),
        ]

        # Create the figure
        fig = plot_trade(chart_data=rets_tab, trade_title=trade_title, trade_subtitle=trade_subtitle,
                         table_description=table_description, chart_description=chart_description,
                         trade_summary_data=trade_summary_data, firm=firm, fee=fee)

        pio.write_html(fig, file='html/{}_{}_{}_spykiller12.html'.format(Date, firm, sym),
                       auto_open=False)
        pio.write_image(fig, file='Images/{}_{}_{}_spykiller12.png'.format(Date, firm, sym), format='png', scale=6)
        if firm == "L":
            prs = Presentation('rebrand_test.pptx')
        elif firm == "O":
            prs = Presentation('rebrand_test_oak.pptx')
        else:
            prs = Presentation('rebrand_test_cfm.pptx')
        prs.save('Images/{}_{}_{}_spykiller12.pptx'.format(Date, firm, sym))
        shapes = prs.slides[0].shapes
        pic = shapes.add_picture('Images/{}_{}_{}_spykiller12.png'.format(Date, firm, sym), Inches(0.001),
                                 Inches(0.001),
                                 height=Inches(10.99), width=Inches(8.5))
        prs.save('Images/{}_{}_{}_spykiller12.pptx'.format(Date, firm, sym))
        return


def uncapped(get_data, firm, output, fee=0.0125):
    """
    uncapped trade using UST
    :param get_data: <str> y or n only - y will pull fresh data from tr, n will use stored data
    :param output: <str> "line" - will provide terms, "trades" will provide trades, "chart" will provide fact sheet
    :param firm: <str> "L" for Lido, "O" for Oakhurst
    :return: based on output param
    """
    if get_data == 'y':  # get data if necessary
        od.get_options(['SPY'])
    else:  # otherwise pass
        pass
    sym = pd.read_pickle('assets/{}_sym.pkl'.format("SPY")).iloc[0]['ticker']  # import the symbol
    und = pd.read_pickle('assets/{}_cached_und.pkl'.format(sym))  ##Gets Underlying Price###
    bnd_dat = pd.read_pickle('assets/bond_data.pkl')  # import bond data
    last = und.iloc[0]['CF_LAST']  # get most recent price

    # Set up Bond Data
    date = "2025-12-19"  ### Sets date into string
    Date = date
    bnd_lst = pd.read_pickle('assets/cached_bndlst.pkl')  ### reads list of UST bonds from pickles
    bnd_lst['MATUR_DATE'] = pd.to_datetime(bnd_lst['MATUR_DATE'])  # convert maturity dates to datetime
    bnd_dat['MATUR_DATE'] = pd.to_datetime(bnd_dat['MATUR_DATE'])  # convert maturity dates to datetime
    bnd_dat = bnd_dat.dropna()  # remove na's
    ind = abs(bnd_dat['MATUR_DATE'] - pd.to_datetime(date)).idxmin()  ## id the relevant bond
    if bnd_dat.loc[[ind]]['DIRTY_PRC'] is np.NaN:  ##this if ensures we skip over "when issued bonds"
        ind = ind + 1  # move up one
        bond = bnd_dat.loc[[ind]]  # isolate bond
    else:  # otherwise
        bond = bnd_dat.loc[[ind]]  # isolate bond
    bond_CU = bond.iloc[0]['Instrument']  # store cusip

    # begin options work
    # options chain set up
    chain = pd.read_pickle('assets/{}_cached_chain.pkl'.format(sym))  # read in the options chaing to pandas
    chain['MID'] = (chain['CF_ASK'] + chain['CF_BID']) / 2
    trade_chain = chain[chain['EXPIR_DATE'] == Date]  # isolate the chain we want
    sorted_trade_chain = trade_chain.sort_values(['PUTCALLIND', 'STRIKE_PRC'],
                                                 axis=0, )  # sort the chain by type and strike
    sorted_trade_chain_calls = sorted_trade_chain[sorted_trade_chain['PUTCALLIND'] == 'CALL']  # create a call df
    sorted_trade_chain_puts = sorted_trade_chain[sorted_trade_chain['PUTCALLIND'] == 'PUT ']  # create a put df

    # ATM CALL
    LCALL = sorted_trade_chain_calls.loc[
        sorted_trade_chain_calls['STRIKE_PRC'].sub(last).abs().idxmin()]  # get atm call

    # Solve for the other two legs
    date = datetime.datetime.strptime(Date, '%Y-%m-%d').date()  # convert date to datetime
    today = datetime.datetime.now().date()  # get today's date
    time_left = date - today  # days left
    adj_time_left = time_left / datetime.timedelta(days=1)  # convert to flt
    time_to_bnd = (adj_time_left / 365)  # convert to annual
    bnd = bnd_dat[bnd_dat['Instrument'] == bond_CU]  # get relevant bond info
    bnd['Minimum PAR'] = LCALL['STRIKE_PRC'] * 100
    bnd['1x Cost'] = bnd['Minimum PAR'] * bnd['DIRTY_PRC'] / 100
    rf = (1 + (bnd['SEC_YLD_1'].iloc[0] / 100)) ** time_to_bnd  # calculate and store deanualzied yield
    notional = LCALL['STRIKE_PRC']  # trade's notional is the strike of the option
    bnd['TR'] = notional - (notional / rf)  # calculate total expected return from bond

    # begin work on put
    tgt = LCALL["MID"] - bnd['TR'].iloc[0]  # tgt put value
    SPUT = sorted_trade_chain_puts.loc[sorted_trade_chain_puts['MID'].sub(tgt).abs().idxmin()]  # get atm call
    LCALL['Trans'] = 'BCO'
    SPUT['Trans'] = 'SPO'  # add trade type
    Trade = pd.concat([LCALL.to_frame().T, SPUT.to_frame().T, ])  # create Trade df
    exp_date = pd.to_datetime(Trade['EXPIR_DATE'].iloc[0]).strftime('%y%m%d')  # get options formated date
    option_type = [s[0] for s in Trade["PUTCALLIND"].to_list()]  # isolate first letter of option type
    strikes = Trade['STRIKE_PRC'].to_list()  # # isolate strikes in list
    option_sym = ['{}{}{}{}'.format(sym, exp_date, option_type[i], int(strikes[i])) for i in
                  range(0, len(Trade))]  # create symbols
    Trade['Symbol'] = option_sym  # add to df
    trade_details = Trade.filter(['Trans', 'Symbol', 'STRIKE_PRC', 'CF_BID', 'CF_ASK'])  # create trade_det df
    upside = "Uncapped"  # upside calculation
    protection = np.abs((SPUT['STRIKE_PRC'] / LCALL['STRIKE_PRC']) - 1)  # protection calculation
    annual_p = protection * (365 / adj_time_left)
    months_left = np.round((adj_time_left / 365) * 12, 1)
    trade_line = pd.DataFrame({"Underlying Asset": sym, "Asset Price": last,
                               'Minimum ($)': round(LCALL['STRIKE_PRC'] * 100, -2),
                               "Expiration Date": Trade['EXPIR_DATE'].iloc[0], "Months Left": months_left,
                               'Potential Upside (%)': upside,
                               'Annual Potential Upside (%)': upside,
                               'Downside Before Protection (%)': '{:.2%}'.format(0),
                               'Protection (%)': '{:.2%}'.format(protection),
                               'Annual_Protection (%)': '{:.2%}'.format(annual_p),
                               'Spread (%)': 'NA', 'Collateral': 'UST'}, index=[0])  # creates line

    # begin work on plot data:
    chrt_rng = np.linspace(LCALL['STRIKE_PRC'] * 0.5, LCALL['STRIKE_PRC'] * 1.5, 50, dtype=float)  # set chart range
    chrt_rng = np.round(chrt_rng, 2)
    lcall_ev = [np.maximum(p - LCALL['STRIKE_PRC'], 0) for p in chrt_rng]  # calc lcall end val
    sput_ev = [np.maximum(SPUT['STRIKE_PRC'] - p, 0) * -1 for p in chrt_rng]  # calc sput end val
    perf_df = pd.DataFrame({"{} Price".format(sym): chrt_rng, "LCALL": lcall_ev, "SPUT": sput_ev,
                            "BOND": LCALL['STRIKE_PRC']})  # Create the df
    perf_df = perf_df.set_index("{} Price".format(sym))  # Set the mkt px as the index
    perf_df['Trade'] = perf_df.sum(axis=1)  # calculate total value
    cost = LCALL['MID'] - SPUT['MID'] + LCALL['STRIKE_PRC'] - bnd['TR'].iloc[0]  # trade cost
    perf_df['Trade Return'] = (perf_df['Trade'] / (cost)) - 1  # trade return
    perf_df = perf_df.sort_index(ascending=False)  # reorder in descending
    perf_df['{} Price Return'.format(sym)] = [(p / last) - 1 for p in perf_df.index]  # add underlying performance
    rets_tab = perf_df.filter(perf_df.columns[-2:]).reset_index()  # reset index
    rets_tab['Trade Return Net'] = [i - float(fee * time_to_bnd) for i in perf_df['Trade Return'].to_list()]
    fc_date = pd.to_datetime(Trade['EXPIR_DATE'].iloc[0]).strftime("%m-%d-%Y")

    if output == 'line':  # line output
        return trade_line
    elif output == 'trades':  # trades output
        print(trade_details)
        print(bnd.filter(['Instrument', 'DIRTY_PRC', 'SEC_YLD_1', 'COUPN_RATE', 'MATUR_DATE',
                          'Minimum PAR', '1x Cost', ]))
    elif output == 'execute':
        op_exec = trade_details
        op_exec['MID'] = (op_exec['CF_BID'] + op_exec['CF_ASK']) / 2
        op_exec['units'] = [1 for i in range(0, len(op_exec))]
        bnd_exec = bnd.filter(['Instrument', 'Minimum PAR', 'DIRTY_PRC', ])
        bnd_exec = bnd_exec.rename(columns={'Minimum PAR': 'units'})
        op_exec.to_csv('assets/temp_op.csv')
        bnd_exec.to_csv('assets/temp_bnd.csv')
    else:
        # Figure objects
        trade_title = "Uncapped Plus<br>Trade"
        trade_subtitle = f"Ending {fc_date}"
        table_description = f"The investor gets about 100% of the price appreciation on {sym} beginning at a<br>" \
                            f"price of {int(LCALL['STRIKE_PRC'])} and full protection between " \
                            f"{int(LCALL['STRIKE_PRC'])} and {int(SPUT['STRIKE_PRC'])} ({-protection:.2%})"
        chart_description = f"{sym} Price (Created on: {today} when {sym} was at {last})"
        trade_summary_data = [
            ('Structure', 'Unlimited Upside + Protection'),
            ('Term', f"{months_left} Months"),
            ('Underlying Asset', sym),
            ('Cap', 'Uncapped'),
            ('Downside Before Protection', 'No Gap'),
            ('Protection', f"{protection:.2%}"),
            ('Fee', f"{fee:.2%} Annually"),
        ]

        # Create the figure
        fig = plot_trade(chart_data=rets_tab, trade_title=trade_title, trade_subtitle=trade_subtitle,
                         table_description=table_description, chart_description=chart_description,
                         trade_summary_data=trade_summary_data, firm=firm, fee=fee)

        pio.write_html(fig, file='html/{}_{}_{}_uncust.html'.format(Date, firm, sym),
                       auto_open=False)
        pio.write_image(fig, file='Images/{}_{}_{}_uncappedust.png'.format(Date, firm, sym), format='png', scale=6)
        if firm == "L":
            prs = Presentation('rebrand_test.pptx')
        elif firm == "O":
            prs = Presentation('rebrand_test_oak.pptx')
        else:
            prs = Presentation('rebrand_test_cfm.pptx')
        prs.save('Images/{}_{}_{}_uncappedust.pptx'.format(Date, firm, sym))
        shapes = prs.slides[0].shapes
        pic = shapes.add_picture('Images/{}_{}_{}_uncappedust.png'.format(Date, firm, sym), Inches(0.001),
                                 Inches(0.001),
                                 height=Inches(10.99), width=Inches(8.5))
        prs.save('Images/{}_{}_{}_uncappedust.pptx'.format(Date, firm, sym))
        return


def uncapped_gap(get_data, firm, output, Date="2025-12-19", fee=0.0125):
    """
    uncapped_gapp trade using corporate bond etf
    :param Date: <str> expiration date of the desired options chain.
    :param get_data: <str> y or n only - y will pull fresh data from tr, n will use stored data
    :param output: <str> "line" - will provide terms, "trades" will provide trades, "chart" will provide fact sheet
    :param firm: <str> "L" for Lido, "O" for Oakhurst
    :return: based on output param
    """
    # Get the Bond ETF Fund that will be used for collateral for this expiry date
    collateral_etf = BOND_ETF_COLLATERAL.get(Date)
    if collateral_etf is None:
        raise ValueError("Invalid date. Provided expiry date must be a valid Dec expiry date.")

    if get_data == 'y':  # get data if necessary
        od.get_options(["SPY"])
        ytw = od.calc_credit(collateral_etf, get_data='y')
    else:  # otherwise pass
        ytw = od.calc_credit(collateral_etf, get_data='n')

    sym = pd.read_pickle('assets/{}_sym.pkl'.format("SPY")).iloc[0]['ticker']  # import the symbol
    und = pd.read_pickle('assets/{}_cached_und.pkl'.format(sym))  ##Gets Underlying Price###
    bnd_dat = pd.read_pickle('assets/bond_data.pkl')  # import bond data
    last = und.iloc[0]['CF_LAST']  # get most recent price

    # Set up Bond Data

    # begin options work
    # options chain set up
    chain = pd.read_pickle('assets/{}_cached_chain.pkl'.format(sym))  # read in the options chaing to pandas
    chain['MID'] = (chain['CF_BID'] + chain['CF_ASK']) / 2
    trade_chain = chain[chain['EXPIR_DATE'] == Date]  # isolate the chain we want
    sorted_trade_chain = trade_chain.sort_values(['PUTCALLIND', 'STRIKE_PRC'],
                                                 axis=0, )  # sort the chain by type and strike
    sorted_trade_chain_calls = sorted_trade_chain[sorted_trade_chain['PUTCALLIND'] == 'CALL']  # create a call df
    sorted_trade_chain_puts = sorted_trade_chain[sorted_trade_chain['PUTCALLIND'] == 'PUT ']  # create a put df

    # ATM CALL
    LCALL = sorted_trade_chain_calls.loc[
        sorted_trade_chain_calls['STRIKE_PRC'].sub(last).abs().idxmin()]  # get atm call
    SPUT_ATM = sorted_trade_chain_puts.loc[
        sorted_trade_chain_puts['STRIKE_PRC'].sub(last).abs().idxmin()]  # get atm put
    LPUT_OTM = sorted_trade_chain_puts.loc[sorted_trade_chain_puts['STRIKE_PRC'].sub(last * 0.93).abs().idxmin()]
    # Solve for the other two legs
    date = datetime.datetime.strptime(Date, '%Y-%m-%d').date()  # convert date to datetime
    today = datetime.datetime.now().date()  # get today's date
    time_left = date - today  # days left
    adj_time_left = time_left / datetime.timedelta(days=1)  # convert to flt
    time_to_bnd = (adj_time_left / 365)  # convert to annual
    bnd = ek.get_data(f"{collateral_etf}.O", fields=['CF_LAST'])[0]  ###Gets ETF collateral Price###
    rf = (1 + np.round(ytw, 4)) ** time_to_bnd
    bnd['1x Cost'] = (LCALL['STRIKE_PRC'] * 100) / rf
    bnd['Shares'] = (bnd['1x Cost'] / bnd['CF_LAST']).round(0)
    notional = LCALL['STRIKE_PRC'] * 1  # trade's notional is the strike of the option
    bnd['TR'] = (bnd['1x Cost'] * rf) / 100 - bnd['1x Cost'] / 100  # calculate total expected return from bond
    # begin work on put

    tgt = LCALL["MID"] - SPUT_ATM["MID"] + LPUT_OTM["MID"] - bnd['TR'].iloc[0]  # tgt put value
    SPUT_OTM = sorted_trade_chain_puts.loc[
        sorted_trade_chain_puts['MID'].sub(tgt).abs().idxmin()]  # solve for short put
    LCALL['Trans'] = 'BCO'  # add trade type
    SPUT_ATM['Trans'] = 'SPO'  # add trade type
    LPUT_OTM['Trans'] = 'BPO'  # add trade type
    SPUT_OTM['Trans'] = 'SPO'  # add trade type
    Trade = pd.concat([LCALL.to_frame().T, SPUT_ATM.to_frame().T,
                       LPUT_OTM.to_frame().T, SPUT_OTM.to_frame().T])  # create Trade df
    exp_date = pd.to_datetime(Trade['EXPIR_DATE'].iloc[0]).strftime('%y%m%d')  # get options formated date
    option_type = [s[0] for s in Trade["PUTCALLIND"].to_list()]  # isolate first letter of option type
    strikes = Trade['STRIKE_PRC'].to_list()  # # isolate strikes in list
    option_sym = ['{}{}{}{}'.format(sym, exp_date, option_type[i], int(strikes[i])) for i in
                  range(0, len(Trade))]  # create symbols
    Trade['Symbol'] = option_sym  # add to df
    trade_details = Trade.filter(['Trans', 'Symbol', 'STRIKE_PRC', "MID"])  # create trade_det df
    upside = "Uncapped"  # upside calculation
    protection = np.abs((((SPUT_OTM['STRIKE_PRC'] / LCALL['STRIKE_PRC']) - 1) - (
            (LPUT_OTM['STRIKE_PRC'] / LCALL['STRIKE_PRC']) - 1)))  # protection calculation
    annual_p = protection * (365 / adj_time_left)
    ds_before = (LCALL['STRIKE_PRC'] / LPUT_OTM['STRIKE_PRC']) - 1
    months_left = np.round((adj_time_left / 365) * 12, 1)
    trade_line = pd.DataFrame({"Underlying Asset": sym, "Asset Price": last,
                               'Minimum ($)': round(LCALL['STRIKE_PRC'] * 100, -2),
                               "Expiration Date": Trade['EXPIR_DATE'].iloc[0], "Months Left": months_left,
                               'Potential Upside (%)': upside,
                               'Annual Potential Upside (%)': upside,
                               'Downside Before Protection (%)': '{:.2%}'.format(ds_before),
                               'Protection (%)': '{:.2%}'.format(protection),
                               'Annual_Protection (%)': '{:.2%}'.format(annual_p),
                               'Spread (%)': 'NA', 'Collateral': collateral_etf}, index=[0])  # creates line
    # begin work on plot data:
    chrt_rng = np.linspace(LCALL['STRIKE_PRC'] * 0.5, LCALL['STRIKE_PRC'] * 1.5, 50, dtype=float)  # set chart range
    chrt_rng = np.round(chrt_rng, 2)
    lcall_ev = [np.maximum(p - LCALL['STRIKE_PRC'], 0) for p in chrt_rng]  # calc lcall end val
    sput_atm_ev = [np.maximum(SPUT_ATM['STRIKE_PRC'] - p, 0) * -1 for p in chrt_rng]  # calc sput end val
    lput_otm_ev = [np.maximum(LPUT_OTM['STRIKE_PRC'] - p, 0) * 1 for p in chrt_rng]  # calc sput end val
    sput_otm_ev = [np.maximum(SPUT_OTM['STRIKE_PRC'] - p, 0) * -1 for p in chrt_rng]  # calc sput end val
    perf_df = pd.DataFrame({"{} Price".format(sym): chrt_rng, "LCALL": lcall_ev, "SPUT_ATM": sput_atm_ev,
                            "LPUT_OTM": lput_otm_ev, "SPUT_OTM": sput_otm_ev,
                            "BOND": LCALL['STRIKE_PRC']})  # Create the df
    perf_df = perf_df.set_index("{} Price".format(sym))  # Set the mkt px as the index
    perf_df['Trade'] = perf_df.sum(axis=1)  # calculate total value
    cost = LCALL['MID'] - SPUT_ATM['MID'] + LPUT_OTM["MID"] - SPUT_OTM["MID"] + \
           (LCALL['STRIKE_PRC'] - bnd['TR'].iloc[0])  # trade cost
    perf_df['Trade Return'] = (perf_df['Trade'] / (cost)) - 1  # trade return
    perf_df = perf_df.sort_index(ascending=False)  # reorder in descending
    perf_df['{} Price Return'.format(sym)] = [(p / last) - 1 for p in perf_df.index]  # add underlying performance
    rets_tab = perf_df.filter(perf_df.columns[-2:]).reset_index()  # reset index
    rets_tab['Trade Return Net'] = [i - float(fee * time_to_bnd) for i in perf_df['Trade Return'].to_list()]
    fc_date = pd.to_datetime(Trade['EXPIR_DATE'].iloc[0]).strftime("%m-%d-%Y")

    if output == 'line':  # line output
        return trade_line
    elif output == 'trades':  # trades output
        print(trade_details)
        print(bnd.filter(['Instrument', "CF_ASK", "Shares", '1x Cost', ]))
    elif output == 'execute':
        op_exec = trade_details
        op_exec['units'] = [1 for i in range(0, len(op_exec))]
        bnd_exec = bnd.filter(['Instrument', 'Shares'])
        bnd_exec = bnd_exec.rename(columns={'Shares': 'units'})
        op_exec.to_csv('assets/temp_op.csv')
        bnd_exec.to_csv('assets/temp_bnd.csv')
    else:
        # Figure objects
        trade_title = "Uncapped and<br>Gap Trade"
        trade_subtitle = f"Ending {fc_date}"
        table_description = f"The investor gets about 100% of the price appreciation on {sym} beginning at a price<br>" \
                            f"of {LCALL['STRIKE_PRC']} and full protection between " \
                            f"{LPUT_OTM['STRIKE_PRC']} and {SPUT_OTM['STRIKE_PRC']} ({-protection:.2%})"
        chart_description = f"{sym} Price (Created on: {today} when {sym} was at {last})"
        trade_summary_data = [
            ('Structure', 'Unlimited upside with deferred downside protection'),
            ('Term', f"{months_left} Months"),
            ('Underlying Asset', sym),
            ('Cap', 'Uncapped'),
            ('Downside Before Protection', f"{ds_before:.2%} Gap"),
            ('Protection', f"{protection:.2%}"),
            ('Fee', f"{fee:.2%} Annually"),
        ]

        # Create the figure
        fig = plot_trade(chart_data=rets_tab, trade_title=trade_title, trade_subtitle=trade_subtitle,
                         table_description=table_description, chart_description=chart_description,
                         trade_summary_data=trade_summary_data, firm=firm, fee=fee)

        pio.write_html(fig, file='html/{}_{}_{}_uncapped_gap.html'.format(Date, firm, sym),
                       auto_open=False)
        pio.write_image(fig, file='Images/{}_{}_{}_uncapped_gap.png'.format(Date, firm, sym), format='png', scale=6)
        if firm == "L":
            prs = Presentation('rebrand_test.pptx')
        elif firm == "O":
            prs = Presentation('rebrand_test_oak.pptx')
        else:
            prs = Presentation('rebrand_test_cfm.pptx')
        prs.save('Images/{}_{}_{}_uncapped_gap.pptx'.format(Date, firm, sym))
        shapes = prs.slides[0].shapes
        pic = shapes.add_picture('Images/{}_{}_{}_uncapped_gap.png'.format(Date, firm, sym), Inches(0.001),
                                 Inches(0.001), height=Inches(10.99), width=Inches(8.5))
        prs.save('Images/{}_{}_{}_uncapped_gap.pptx'.format(Date, firm, sym))
        return


def no_downside(get_data, firm, output, fee=0.0125):
    """
    uncapped_gapp trade using corporate bond etf
    :param get_data: <str> y or n only - y will pull fresh data from tr, n will use stored data
    :param output: <str> "line" - will provide terms, "trades" will provide trades, "chart" will provide fact sheet
    :param firm: <str> "L" for Lido, "O" for Oakhurst
    :return: based on output param
    """
    if get_data == 'y':  # get data if necessary
        od.get_options(["SPY"])
        ytw = od.calc_credit('BSCP', get_data='y')
    else:  # otherwise pass
        ytw = od.calc_credit('BSCP', get_data='n')

    Date = '2025-12-19'
    sym = pd.read_pickle('assets/{}_sym.pkl'.format("SPY")).iloc[0]['ticker']  # import the symbol
    und = pd.read_pickle('assets/{}_cached_und.pkl'.format(sym))  ##Gets Underlying Price###
    bnd_dat = pd.read_pickle('assets/bond_data.pkl')  # import bond data
    last = und.iloc[0]['CF_LAST']  # get most recent price

    # Set up Bond Data

    # begin options work
    # options chain set up
    chain = pd.read_pickle('assets/{}_cached_chain.pkl'.format(sym))  # read in the options chaing to pandas
    chain['MID'] = (chain['CF_BID'] + chain['CF_ASK']) / 2
    trade_chain = chain[chain['EXPIR_DATE'] == Date]  # isolate the chain we want
    sorted_trade_chain = trade_chain.sort_values(['PUTCALLIND', 'STRIKE_PRC'],
                                                 axis=0, )  # sort the chain by type and strike
    sorted_trade_chain_calls = sorted_trade_chain[sorted_trade_chain['PUTCALLIND'] == 'CALL']  # create a call df

    # ATM CALL
    LCALL = sorted_trade_chain_calls.loc[
        sorted_trade_chain_calls['STRIKE_PRC'].sub(last).abs().idxmin()]  # get atm call
    date = datetime.datetime.strptime(Date, '%Y-%m-%d').date()  # convert date to datetime
    today = datetime.datetime.now().date()  # get today's date
    time_left = date - today  # days left
    adj_time_left = time_left / datetime.timedelta(days=1)  # convert to flt
    time_to_bnd = (adj_time_left / 365)  # convert to annual
    bnd = ek.get_data("BSCP.O", fields=['CF_LAST'])[0]  ###Gets BSCN Price###
    rf = (1 + np.round(ytw, 4)) ** time_to_bnd
    bnd['1x Cost'] = (LCALL['STRIKE_PRC'] * 100) / rf
    bnd['Shares'] = (bnd['1x Cost'] / bnd['CF_LAST']).round(0)
    notional = LCALL['STRIKE_PRC'] * 1  # trade's notional is the strike of the option
    bnd['TR'] = (bnd['1x Cost'] * rf) / 100 - bnd['1x Cost'] / 100  # calculate total expected return from bond
    scall_tgt = LCALL['MID'] - bnd['TR'].iloc[0]
    scall_tab = np.abs((sorted_trade_chain_calls['MID'] - scall_tgt)).idxmin()
    SCALL = sorted_trade_chain_calls.loc[scall_tab]
    LCALL['Trans'] = 'BCO'  # add trade type
    SCALL['Trans'] = 'SCO'  # add trade type
    Trade = pd.concat([pd.DataFrame(SCALL).T, pd.DataFrame(LCALL).T]).reset_index(drop=True)
    exp_date = pd.to_datetime(Trade['EXPIR_DATE'].iloc[0]).strftime('%y%m%d')  # get options formated date
    option_type = [s[0] for s in Trade["PUTCALLIND"].to_list()]  # isolate first letter of option type
    strikes = Trade['STRIKE_PRC'].to_list()  # # isolate strikes in list
    option_sym = ['{}{}{}{}'.format(sym, exp_date, option_type[i], int(strikes[i])) for i in
                  range(0, len(Trade))]  # create symbols
    Trade['Symbol'] = option_sym  # add to df
    trade_details = Trade.filter(['Trans', 'Symbol', 'STRIKE_PRC', "MID"])  # create trade_det df
    upside = (SCALL['STRIKE_PRC'] / LCALL['STRIKE_PRC']) - 1  # upside calculation
    annual_up = upside * (365 / adj_time_left)
    protection = 1
    annual_p = 1
    months_left = np.round((adj_time_left / 365) * 12, 1)
    trade_line = pd.DataFrame({"Underlying Asset": sym, "Asset Price": last,
                               'Minimum ($)': round(LCALL['STRIKE_PRC'] * 100, -2),
                               "Expiration Date": Trade['EXPIR_DATE'].iloc[0], "Months Left": months_left,
                               'Potential Upside (%)': '{:.2%}'.format(upside),
                               'Annual Potential Upside (%)': '{:.2%}'.format(annual_up),
                               'Downside Before Protection (%)': '{:.2%}'.format(0),
                               'Protection (%)': '{:.2%}'.format(protection),
                               'Annual_Protection (%)': '{:.2%}'.format(annual_p),
                               'Spread (%)': '{:.2%}'.format(upside + protection), 'Collateral': 'BSCP'},
                              index=[0])  # creates line

    # begin work on plot data:
    chrt_rng = np.linspace(LCALL['STRIKE_PRC'] * 0.5, LCALL['STRIKE_PRC'] * 1.5, 50, dtype=float)  # set chart range
    chrt_rng = np.round(chrt_rng, 2)
    scall_ev = [np.maximum(p - SCALL['STRIKE_PRC'], 0) * -1 for p in chrt_rng]  # calc scall end val
    lcall_ev = [np.maximum(p - LCALL['STRIKE_PRC'], 0) for p in chrt_rng]  # calc lcall end val
    perf_df = pd.DataFrame({"{} Price".format(sym): chrt_rng, "SCALL": scall_ev, "LCALL": lcall_ev,
                            "BOND": LCALL['STRIKE_PRC']})  # Create the df
    perf_df = perf_df.set_index("{} Price".format(sym))  # Set the mkt px as the index
    perf_df['Trade'] = perf_df.sum(axis=1)  # calculate total value
    cost = LCALL['MID'] - SCALL['MID'] + LCALL['STRIKE_PRC'] - bnd['TR'].iloc[0]  # trade cost
    perf_df['Trade Return'] = (perf_df['Trade'] / (cost)) - 1  # trade return
    perf_df = perf_df.sort_index(ascending=False)  # reorder in descending
    perf_df['{} Price Return'.format(sym)] = [(p / last) - 1 for p in perf_df.index]  # add underlying performance
    rets_tab = perf_df.filter(perf_df.columns[-2:]).reset_index()  # reset index
    rets_tab['Trade Return Net'] = [i - float(fee * time_to_bnd) for i in perf_df['Trade Return'].to_list()]
    fc_date = pd.to_datetime(Trade['EXPIR_DATE'].iloc[0]).strftime("%m-%d-%Y")
    upside = perf_df['Trade Return'].max()

    if output == 'line':  # line output
        return trade_line
    elif output == 'trades':  # trades output
        print(trade_details)
        print(bnd.filter(['Instrument', "CF_ASK", "Shares", '1x Cost', ]))
    elif output == 'execute':
        op_exec = trade_details
        op_exec['units'] = [1 for i in range(0, len(op_exec))]
        bnd_exec = bnd.filter(['Instrument', 'Shares'])
        bnd_exec = bnd_exec.rename(columns={'Shares': 'units'})
        op_exec.to_csv('assets/temp_op.csv')
        bnd_exec.to_csv('assets/temp_bnd.csv')

    else:
        # Figure objects
        trade_title = "Cap and Go<br>Trade"
        trade_subtitle = f"Ending {fc_date}"
        table_description = f"The investor gets about 100% of the price appreciation on {sym} beginning at a price " \
                            f"of {LCALL['STRIKE_PRC']}, capped at {SCALL['STRIKE_PRC']}<br>and full protection to the downside."
        chart_description = f"{sym} Price (Created on: {today} when {sym} was at {last})"
        trade_summary_data = [
            ('Structure', 'Cap and Go'),
            ('Term', f"{months_left} Months"),
            ('Underlying Asset', sym),
            ('Cap', f"Capped at {upside:.2%}"),
            ('Downside Before Protection', 'No Gap'),
            ('Protection', f"{protection:.2%}"),
            ('Fee', f"{fee:.2%} Annually"),
        ]

        # Create the figure
        fig = plot_trade(chart_data=rets_tab, trade_title=trade_title, trade_subtitle=trade_subtitle,
                         table_description=table_description, chart_description=chart_description,
                         trade_summary_data=trade_summary_data, firm=firm, fee=fee)

        pio.write_html(fig, file='html/{}_{}_{}_cgo.html'.format(Date, firm, sym),
                       auto_open=False)
        pio.write_image(fig, file='Images/{}_{}_{}_cgo.png'.format(Date, firm, sym), format='png', scale=6)
        if firm == "L":
            prs = Presentation('rebrand_test.pptx')
        elif firm == "O":
            prs = Presentation('rebrand_test_oak.pptx')
        else:
            prs = Presentation('rebrand_test_cfm.pptx')
        prs.save('Images/{}_{}_{}_cgo.pptx'.format(Date, firm, sym))
        shapes = prs.slides[0].shapes
        pic = shapes.add_picture('Images/{}_{}_{}_cgo.png'.format(Date, firm, sym), Inches(0.001), Inches(0.001),
                                 height=Inches(10.99), width=Inches(8.5))
        prs.save('Images/{}_{}_{}_cgo.pptx'.format(Date, firm, sym))
        return


def eighty(get_data, firm, output, fee=0.0125):
    """
    uncapped_gapp trade using corporate bond etf
    :param get_data: <str> y or n only - y will pull fresh data from tr, n will use stored data
    :param output: <str> "line" - will provide terms, "trades" will provide trades, "chart" will provide fact sheet
    :param firm: <str> "L" for Lido, "O" for Oakhurst
    :return: based on output param
    """
    if get_data == 'y':  # get data if necessary
        od.get_options(["SPY"])
        ytw = od.calc_credit('BSCP', get_data='y')
    else:  # otherwise pass
        ytw = od.calc_credit('BSCP', get_data='n')

    Date = '2025-12-19'
    sym = pd.read_pickle('assets/{}_sym.pkl'.format("SPY")).iloc[0]['ticker']  # import the symbol
    und = pd.read_pickle('assets/{}_cached_und.pkl'.format(sym))  ##Gets Underlying Price###
    bnd_dat = pd.read_pickle('assets/bond_data.pkl')  # import bond data
    last = und.iloc[0]['CF_LAST']  # get most recent price

    # Set up Bond Data

    # begin options work
    # options chain set up
    chain = pd.read_pickle('assets/{}_cached_chain.pkl'.format(sym))  # read in the options chaing to pandas
    chain['MID'] = (chain['CF_BID'] + chain['CF_ASK']) / 2
    trade_chain = chain[chain['EXPIR_DATE'] == Date]  # isolate the chain we want
    sorted_trade_chain = trade_chain.sort_values(['PUTCALLIND', 'STRIKE_PRC'],
                                                 axis=0, )  # sort the chain by type and strike
    sorted_trade_chain_calls = sorted_trade_chain[sorted_trade_chain['PUTCALLIND'] == 'CALL']  # create a call df

    # ATM CALL
    LCALL = sorted_trade_chain_calls.loc[
        sorted_trade_chain_calls['STRIKE_PRC'].sub(last).abs().idxmin()]  # get atm call
    date = datetime.datetime.strptime(Date, '%Y-%m-%d').date()  # convert date to datetime
    today = datetime.datetime.now().date()  # get today's date
    time_left = date - today  # days left
    adj_time_left = time_left / datetime.timedelta(days=1)  # convert to flt
    time_to_bnd = (adj_time_left / 365)  # convert to annual
    bnd = ek.get_data("BSCP.O", fields=['CF_LAST'])[0]  ###Gets BSCN Price###
    rf = (1 + np.round(ytw, 4)) ** time_to_bnd
    bnd['1x Cost'] = (LCALL['STRIKE_PRC'] * 500) / rf
    bnd['Shares'] = (bnd['1x Cost'] / bnd['CF_LAST']).round(0)
    notional = LCALL['STRIKE_PRC'] * 1  # trade's notional is the strike of the option
    bnd['TR'] = (bnd['1x Cost'] * rf) - bnd['1x Cost']  # calculate total expected return from bond
    mult = np.round(bnd['TR'] / (LCALL['MID'] * 100), 0) * 100
    mult = mult[0]
    LCALL['Trans'] = 'BCO'  # add trade type
    Trade = pd.concat([pd.DataFrame(LCALL).T]).reset_index(drop=True)
    exp_date = pd.to_datetime(Trade['EXPIR_DATE'].iloc[0]).strftime('%y%m%d')  # get options formated date
    option_type = [s[0] for s in Trade["PUTCALLIND"].to_list()]  # isolate first letter of option type
    strikes = Trade['STRIKE_PRC'].to_list()  # # isolate strikes in list
    option_sym = ['{}{}{}{}'.format(sym, exp_date, option_type[i], int(strikes[i])) for i in
                  range(0, len(Trade))]  # create symbols
    Trade['Symbol'] = option_sym  # add to df
    trade_details = Trade.filter(['Trans', 'Symbol', 'STRIKE_PRC', "MID"])  # create trade_det df
    upside = 'Uncapped'  # upside calculation
    annual_up = 'Uncapped'
    protection = 1
    annual_p = 1
    months_left = np.round((adj_time_left / 365) * 12, 1)
    trade_line = pd.DataFrame({"Underlying Asset": sym, "Asset Price": last,
                               'Minimum ($)': round(LCALL['STRIKE_PRC'] * 100, -2),
                               "Expiration Date": Trade['EXPIR_DATE'].iloc[0], "Months Left": months_left,
                               'Potential Upside (%)': '{}'.format(upside),
                               'Annual Potential Upside (%)': '{}'.format(annual_up),
                               'Downside Before Protection (%)': '{:.2%}'.format(0),
                               'Protection (%)': '{:.2%}'.format(protection),
                               'Annual_Protection (%)': '{:.2%}'.format(annual_p),
                               'Spread (%)': '{}'.format('N/A'), 'Collateral': 'BSCP'},
                              index=[0])  # creates line

    # begin work on plot data:
    chrt_rng = np.linspace(LCALL['STRIKE_PRC'] * 0.5, LCALL['STRIKE_PRC'] * 1.5, 50, dtype=float)  # set chart range
    chrt_rng = np.round(chrt_rng, 2)
    lcall_ev = [np.maximum(p - LCALL['STRIKE_PRC'], 0) * mult for p in chrt_rng]  # calc lcall end val
    perf_df = pd.DataFrame({"{} Price".format(sym): chrt_rng, "LCALL": lcall_ev,
                            "BOND": LCALL['STRIKE_PRC'] * 500})  # Create the df
    perf_df = perf_df.set_index("{} Price".format(sym))  # Set the mkt px as the index
    perf_df['Trade'] = perf_df.sum(axis=1)  # calculate total value
    cost = LCALL['MID'] * mult + LCALL['STRIKE_PRC'] * 500 - bnd['TR'].iloc[0]  # trade cost
    perf_df['Trade Return'] = (perf_df['Trade'] / (cost)) - 1  # trade return
    perf_df = perf_df.sort_index(ascending=False)  # reorder in descending
    perf_df['{} Price Return'.format(sym)] = [(p / last) - 1 for p in perf_df.index]  # add underlying performance
    rets_tab = perf_df.filter(perf_df.columns[-2:]).reset_index()  # reset index
    rets_tab['Trade Return Net'] = [i - float(fee * time_to_bnd) for i in perf_df['Trade Return'].to_list()]
    fc_date = pd.to_datetime(Trade['EXPIR_DATE'].iloc[0]).strftime("%m-%d-%Y")

    if output == 'line':  # line output
        return trade_line
    elif output == 'trades':  # trades output
        print(trade_details)
        print(bnd.filter(['Instrument', "CF_ASK", "Shares", '1x Cost', ]))
    elif output == 'execute':
        op_exec = trade_details
        op_exec['units'] = [4 for i in range(0, len(op_exec))]
        bnd_exec = bnd.filter(['Instrument', 'Shares', 'CF_LAST'])
        bnd_exec = bnd_exec.rename(columns={'Shares': 'units'})
        op_exec.to_csv('assets/temp_op.csv')
        bnd_exec.to_csv('assets/temp_bnd.csv')
    else:
        # Figure objects
        trade_title = f"{sym} 80 - Up<br>Trade"
        trade_subtitle = f"Ending {fc_date}"
        table_description = f"The investor gets about 70% - 80% of the price appreciation on {sym} beginning<br>" \
                            f"at a price of {LCALL['STRIKE_PRC']} and full protection to the downside."
        chart_description = f"{sym} Price (Created on: {today} when {sym} was at {last})"
        trade_summary_data = [
            ('Structure', '80% Up'),
            ('Term', f"{months_left} Months"),
            ('Underlying Asset', sym),
            ('Cap', 'Uncapped (70%-80% Capture)'),
            ('Downside Before Protection', 'No Gap'),
            ('Protection', f"{protection:.2%}"),
            ('Fee', f"{fee:.2%} Annually"),
        ]

        # Create the figure
        fig = plot_trade(chart_data=rets_tab, trade_title=trade_title, trade_subtitle=trade_subtitle,
                         table_description=table_description, chart_description=chart_description,
                         trade_summary_data=trade_summary_data, firm=firm, fee=fee)

        pio.write_html(fig, file='html/{}_{}_{}_eightyup.html'.format(Date, firm, sym),
                       auto_open=False)
        pio.write_image(fig, file='Images/{}_{}_{}_eightyup.png'.format(Date, firm, sym), format='png', scale=6)
        if firm == "L":
            prs = Presentation('rebrand_test.pptx')
        elif firm == "O":
            prs = Presentation('rebrand_test_oak.pptx')
        else:
            prs = Presentation('rebrand_test_cfm.pptx')
        prs.save('Images/{}_{}_{}_eightyup.pptx'.format(Date, firm, sym))
        shapes = prs.slides[0].shapes
        pic = shapes.add_picture('Images/{}_{}_{}_eightyup.png'.format(Date, firm, sym), Inches(0.001), Inches(0.001),
                                 height=Inches(10.99), width=Inches(8.5))
        prs.save('Images/{}_{}_{}_eightyup.pptx'.format(Date, firm, sym))

        return


def uncapped1_33x(get_data, firm, output, fee=0.0125):
    """
    uncapped 1.33x trade using corporate bond etf
    :param get_data: <str> y or n only - y will pull fresh data from tr, n will use stored data
    :param output: <str> "line" - will provide terms, "trades" will provide trades, "chart" will provide fact sheet
    :param firm: <str> "L" for Lido, "O" for Oakhurst
    :return: based on output param
    """
    if get_data == 'y':  # get data if necessary
        od.get_options(["SPY"])
        ytw = od.calc_credit('BSCP', get_data='y')

    else:  # otherwise pass
        ytw = od.calc_credit('BSCP', get_data='n')

    Date = '2025-12-19'
    sym = pd.read_pickle('assets/{}_sym.pkl'.format("SPY")).iloc[0]['ticker']  # import the symbol
    und = pd.read_pickle('assets/{}_cached_und.pkl'.format(sym))  ##Gets Underlying Price###
    bnd_dat = pd.read_pickle('assets/bond_data.pkl')  # import bond data
    last = und.iloc[0]['CF_LAST']  # get most recent price

    # begin options work
    # options chain set up
    chain = pd.read_pickle('assets/{}_cached_chain.pkl'.format(sym))  # read in the options chaing to pandas
    chain['MID'] = (chain['CF_BID'] + chain['CF_ASK']) / 2
    trade_chain = chain[chain['EXPIR_DATE'] == Date]  # isolate the chain we want
    sorted_trade_chain = trade_chain.sort_values(['PUTCALLIND', 'STRIKE_PRC'],
                                                 axis=0, )  # sort the chain by type and strike
    sorted_trade_chain_calls = sorted_trade_chain[sorted_trade_chain['PUTCALLIND'] == 'CALL']  # create a call df
    sorted_trade_chain_puts = sorted_trade_chain[sorted_trade_chain['PUTCALLIND'] == 'PUT ']  # create a put df

    # ATM CALL
    LCALL = sorted_trade_chain_calls.loc[
        sorted_trade_chain_calls['STRIKE_PRC'].sub(last).abs().idxmin()]  # get atm call
    SPUT = sorted_trade_chain_puts.loc[
        sorted_trade_chain_puts['STRIKE_PRC'].sub(LCALL['STRIKE_PRC']).abs().idxmin()]  # get atm PUTS
    # Solve for the other two legs
    date = datetime.datetime.strptime(Date, '%Y-%m-%d').date()  # convert date to datetime
    today = datetime.datetime.now().date()  # get today's date
    time_left = date - today  # days left
    adj_time_left = time_left / datetime.timedelta(days=1)  # convert to flt
    time_to_bnd = (adj_time_left / 365)  # convert to annual
    bnd = ek.get_data("BSCP.O", fields=['CF_LAST'])[0]  ###Gets BSCN Price###
    rf = (1 + np.round(ytw, 4)) ** time_to_bnd
    bnd['1x Cost'] = (LCALL['STRIKE_PRC'] * 300) / rf
    bnd['Shares'] = (bnd['1x Cost'] / bnd['CF_LAST']).round(0)
    notional = LCALL['STRIKE_PRC'] * 3  # trade's notional is the strike of the option
    bnd['TR'] = (bnd['1x Cost'] / 100 * rf) - bnd['1x Cost'] / 100  # calculate total expected return from bond
    LCALL['Trans'] = 'BCO'
    SPUT['Trans'] = 'SPO'  # add trade type
    Trade = pd.concat([LCALL.to_frame().T, SPUT.to_frame().T, ])  # create Trade df
    exp_date = pd.to_datetime(Trade['EXPIR_DATE'].iloc[0]).strftime('%y%m%d')  # get options formated date
    option_type = [s[0] for s in Trade["PUTCALLIND"].to_list()]  # isolate first letter of option type
    strikes = Trade['STRIKE_PRC'].to_list()  # # isolate strikes in list
    option_sym = ['{}{}{}{}'.format(sym, exp_date, option_type[i], int(strikes[i])) for i in
                  range(0, len(Trade))]  # create symbols
    Trade['Symbol'] = option_sym  # add to df
    trade_details = Trade.filter(['Trans', 'Symbol', 'STRIKE_PRC', "MID"])  # create trade_det df
    trade_details['No. Contracts'] = [4, 3]
    upside = "Uncapped"  # upside calculation
    protection = 0  # protection calculation
    annual_p = protection * (365 / adj_time_left)
    months_left = np.round((adj_time_left / 365) * 12, 1)
    trade_line = pd.DataFrame({"Underlying Asset": sym, "Asset Price": last,
                               'Minimum ($)': round(LCALL['STRIKE_PRC'] * 300, -2),
                               "Expiration Date": Trade['EXPIR_DATE'].iloc[0], "Months Left": months_left,
                               'Potential Upside (%)': upside,
                               'Annual Potential Upside (%)': upside,
                               'Downside Before Protection (%)': '{:.2%}'.format(0),
                               'Protection (%)': '{:.2%}'.format(protection),
                               'Annual_Protection (%)': '{:.2%}'.format(annual_p),
                               'Spread (%)': 'NA', 'Collateral': 'BSCP'}, index=[0])  # creates line

    # begin work on plot data:
    chrt_rng = np.linspace(LCALL['STRIKE_PRC'] * 0.5, LCALL['STRIKE_PRC'] * 1.5, 50, dtype=float)  # set chart range
    chrt_rng = np.round(chrt_rng, 2)
    lcall_ev = [np.maximum(p - LCALL['STRIKE_PRC'], 0) * 4 for p in chrt_rng]  # calc lcall end val
    sput_ev = [np.maximum(SPUT['STRIKE_PRC'] - p, 0) * -3 for p in chrt_rng]  # calc sput end val
    perf_df = pd.DataFrame({"{} Price".format(sym): chrt_rng, "LCALL": lcall_ev, "SPUT": sput_ev,
                            "BOND": LCALL['STRIKE_PRC'] * 3})  # Create the df
    perf_df = perf_df.set_index("{} Price".format(sym))  # Set the mkt px as the index
    perf_df['Trade'] = perf_df.sum(axis=1)  # calculate total value
    cost = LCALL['MID'] * 4 - SPUT['MID'] * 3 + LCALL['STRIKE_PRC'] * 3 - bnd['TR'].iloc[0]  # trade cost
    perf_df['Trade Return'] = (perf_df['Trade'] / (cost)) - 1  # trade return
    perf_df = perf_df.sort_index(ascending=False)  # reorder in descending
    perf_df['{} Price Return'.format(sym)] = [(p / last) - 1 for p in perf_df.index]  # add underlying performance
    rets_tab = perf_df.filter(perf_df.columns[-2:]).reset_index()  # reset index
    rets_tab['Trade Return Net'] = [i - float(fee * time_to_bnd) for i in perf_df['Trade Return'].to_list()]
    fc_date = pd.to_datetime(Trade['EXPIR_DATE'].iloc[0]).strftime("%m-%d-%Y")

    if output == 'line':  # line output
        return trade_line
    elif output == 'trades':  # trades output
        print(trade_details)
        print(bnd.filter(['Instrument', "CF_LAST", "Shares", '1x Cost', ]))
    elif output == 'execute':
        op_exec = trade_details
        op_exec['units'] = [1.34, 1]  # numerical correction for rounding
        bnd_exec = bnd.filter(['Instrument', 'Shares', 'CF_LAST'])
        bnd_exec = bnd_exec.rename(columns={'Shares': 'units'})
        op_exec.to_csv('assets/temp_op.csv')
        bnd_exec.to_csv('assets/temp_bnd.csv')
    else:
        # Figure objects
        trade_title = "Uncapped 1.33x<br>Trade"
        trade_subtitle = f"Ending {fc_date}"
        table_description = f"The investor gets <b>UP TO 133%</b> of the price appreciation on {sym} beginning at a " \
                            f"price<br>of {int(LCALL['STRIKE_PRC'])}"
        chart_description = f"{sym} Price (Created on: {today} when {sym} was at {last})"
        trade_summary_data = [
            ('Structure', '1.33x levered upside with unlevered downside'),
            ('Term', f"{months_left} Months"),
            ('Underlying Asset', sym),
            ('Cap', 'Uncapped'),
            ('Downside Before Protection', 'No Gap'),
            ('Protection', 'No Protection'),
            ('Fee', f"{fee:.2%} Annually"),
        ]

        # Create the figure
        fig = plot_trade(chart_data=rets_tab, trade_title=trade_title, trade_subtitle=trade_subtitle,
                         table_description=table_description, chart_description=chart_description,
                         trade_summary_data=trade_summary_data, firm=firm, fee=fee)

        pio.write_html(fig, file='html/{}_{}_{}_unc133.html'.format(Date, firm, sym),
                       auto_open=False)
        pio.write_image(fig, file='Images/{}_{}_{}_uncapped133.png'.format(Date, firm, sym), format='png', scale=6)
        if firm == "L":
            prs = Presentation('rebrand_test.pptx')
        elif firm == "O":
            prs = Presentation('rebrand_test_oak.pptx')
        else:
            prs = Presentation('rebrand_test_cfm.pptx')
        prs.save('Images/{}_{}_{}_uncapped133.pptx'.format(Date, firm, sym))
        shapes = prs.slides[0].shapes
        pic = shapes.add_picture('Images/{}_{}_{}_uncapped133.png'.format(Date, firm, sym), Inches(0.001),
                                 Inches(0.001),
                                 height=Inches(10.99), width=Inches(8.5))
        prs.save('Images/{}_{}_{}_uncapped133.pptx'.format(Date, firm, sym))
        return


def accel_150(RIC, Date, get_data, firm, output, prot=0.15, fee=0.0125):
    """
    Accelerated upside (150%) trade with cap using corporate bond etf as collateral
    :param RIC: <str> the Reurters identification code for the security in question
    :param Date: <str> expiration date of the desired options chain.
    :param get_data: <str> y or n only - y will pull fresh data from tr, n will use stored data
    :param firm: <str> "L" for Lido, "O" for Oakhurst
    :param output: <str> "line" - will provide terms, "trades" will provide trades, "chart" will provide fact sheet
    :param fee: <int> annual managmeent fee set to 1.25%
    :return: based on output param
    """
    if Date == '2024-12-20':
        collat = 'BSCO'
    elif Date == '2025-12-19':
        collat = 'BSCP'
    else:
        print('ERROR - Please enter a valid expiration date for this trade.')

    if get_data == 'y':  # get data if necessary
        od.get_options([RIC])
        ytw = od.calc_credit(collat, get_data='y')

    else:  # otherwise pass
        ytw = od.calc_credit(collat, get_data='n')

    sym = pd.read_pickle('assets/{}_sym.pkl'.format(RIC)).iloc[0]['ticker']  # import the symbol
    und = pd.read_pickle('assets/{}_cached_und.pkl'.format(sym))  ##Gets Underlying Price###
    bnd_dat = pd.read_pickle('assets/bond_data.pkl')  # import bond data
    last = und.iloc[0]['CF_LAST']  # get most recent price

    # options chain set up
    chain = pd.read_pickle('assets/{}_cached_chain.pkl'.format(sym))  # read in the options chaing to pandas
    chain['MID'] = (chain['CF_BID'] + chain['CF_ASK']) / 2
    trade_chain = chain[chain['EXPIR_DATE'] == Date]  # isolate the chain we want
    sorted_trade_chain = trade_chain.sort_values(['PUTCALLIND', 'STRIKE_PRC'],
                                                 axis=0, )  # sort the chain by type and strike
    sorted_trade_chain_calls = sorted_trade_chain[sorted_trade_chain['PUTCALLIND'] == 'CALL']  # create a call df
    sorted_trade_chain_puts = sorted_trade_chain[sorted_trade_chain['PUTCALLIND'] == 'PUT ']  # create a put df
    LCALL = sorted_trade_chain_calls.loc[
        sorted_trade_chain_calls['STRIKE_PRC'].sub(last).abs().idxmin()]  # long ATM call
    SPUT = sorted_trade_chain_puts.loc[
        sorted_trade_chain_puts['STRIKE_PRC'].sub(last * (1 - prot)).abs().idxmin()]  # Short put at desired protection
    # Calculate bond total return
    bnd_tick = collat + ".O"
    date = datetime.datetime.strptime(Date, '%Y-%m-%d').date()  # convert date to datetime
    today = datetime.datetime.now().date()  # get today's date
    time_left = date - today  # days left
    adj_time_left = time_left / datetime.timedelta(days=1)  # convert to flt
    time_to_bnd = (adj_time_left / 365)  # convert to annual
    bnd = ek.get_data(bnd_tick, fields=['CF_LAST'])[0]  # retieves BSC collateral price

    notional = LCALL['STRIKE_PRC'] * 2  # trade's notional value

    bnd['1x Cost'] = (notional * 100) / (1 + (ytw * time_to_bnd))  # collateral value calculation
    bnd['Shares'] = (bnd['1x Cost'] / bnd['CF_LAST']).round(0)
    bond_val = bnd['1x Cost'].iloc[0] / 100
    bnd['TR'] = bond_val * ytw * time_to_bnd  # simple, non-compounded interest
    bond_tr = bnd['TR'].iloc[0]

    # Select SCALL strike
    combo = (LCALL['MID'] * 3) - (SPUT['MID'] * 2)  # (long ATM call premium x 3) - (Short Put premium x 2)
    scall_tgt = (combo - bond_tr) / 3  # target premium to be collected each Short Call is cost of combo minus bond TR
    SCALL = sorted_trade_chain_calls.loc[
        sorted_trade_chain_calls['MID'].sub(scall_tgt).abs().idxmin()]  # get short call strike

    SCALL['Trans'] = 'SCO'  # add trade type
    LCALL['Trans'] = 'BCO'  # add trade type
    SPUT['Trans'] = 'SPO'  # add trade type

    Trade = pd.concat([pd.DataFrame(SCALL).T,
                       pd.DataFrame(LCALL).T,
                       pd.DataFrame(SPUT).T]).reset_index(drop=True)
    exp_date = pd.to_datetime(Trade['EXPIR_DATE'].iloc[0]).strftime('%y%m%d')  # get options formated date
    option_type = [s[0] for s in Trade["PUTCALLIND"].to_list()]  # isolate first letter of option type
    strikes = Trade['STRIKE_PRC'].to_list()  # # isolate strikes in list
    option_sym = ['{}{}{}{}'.format(sym, exp_date, option_type[i], int(strikes[i])) for i in
                  range(0, len(Trade))]  # create symbols
    Trade['Symbol'] = option_sym  # add to df
    trade_details = Trade.filter(['Trans', 'Symbol', 'STRIKE_PRC', "MID"])  # create trade_det df
    trade_details['No. Contracts'] = [3, 3, 2]  # SCALL contracts, LCALL contracts, SPUT contracts
    protection = 1 - (SPUT['STRIKE_PRC'] / LCALL['STRIKE_PRC'])
    annual_p = protection / time_to_bnd
    upside = ((SCALL['STRIKE_PRC'] / LCALL['STRIKE_PRC']) - 1) * (3 / 2)
    annual_up = upside / time_to_bnd
    months_left = np.round((adj_time_left / 365) * 12, 1)

    trade_line = pd.DataFrame({"Underlying Asset": sym, "Asset Price": last,
                               'Minimum ($)': round(LCALL['STRIKE_PRC'] * 200, -2),
                               "Expiration Date": Trade['EXPIR_DATE'].iloc[0], "Months Left": months_left,
                               'Potential Upside (%)': '{:.2%}'.format(upside),
                               'Annual Potential Upside (%)': '{:.2%}'.format(annual_up),
                               'Downside Before Protection (%)': '{:.2%}'.format(0),
                               'Protection (%)': '{:.2%}'.format(protection),
                               'Annual_Protection (%)': '{:.2%}'.format(annual_p),
                               'Spread (%)': '{:.2%}'.format(upside + protection), 'Collateral': collat},
                              index=[0])  # creates line
    # begin work on plot data:
    chrt_rng = np.linspace(LCALL['STRIKE_PRC'] * 0.5, LCALL['STRIKE_PRC'] * 1.5, 50, dtype=float)  # set chart range
    chrt_rng = np.round(chrt_rng, 2)
    lcall_ev = [np.maximum(p - LCALL['STRIKE_PRC'], 0) * 3 for p in chrt_rng]  # calc atm_lcall end val
    sput_ev = [np.maximum(SPUT['STRIKE_PRC'] - p, 0) * -2 for p in chrt_rng]  # calc atm_sput end val
    scall_ev = [np.maximum(p - SCALL['STRIKE_PRC'], 0) * -3 for p in chrt_rng]  # calc scall end val
    perf_df = pd.DataFrame({"{} Price".format(sym): chrt_rng,
                            "SCALL": scall_ev,
                            "LCALL": lcall_ev,
                            "SPUT": sput_ev,
                            "BOND": LCALL['STRIKE_PRC'] * 2})  # Create the df
    perf_df = perf_df.set_index("{} Price".format(sym))  # Set the mkt px as the index
    perf_df['Trade'] = perf_df.sum(axis=1)  # calculate total value
    cost = (LCALL['MID'] * 3) - (SCALL['MID'] * 3) - (SPUT['MID'] * 2) + notional - bnd['TR'].iloc[0]  # trade cost
    perf_df['Trade Return'] = (perf_df['Trade'] / cost) - 1  # trade return
    perf_df = perf_df.sort_index(ascending=False)  # reorder in descending
    perf_df['{} Price Return'.format(sym)] = [(p / last) - 1 for p in perf_df.index]  # add underlying performance
    rets_tab = perf_df.filter(perf_df.columns[-2:]).reset_index()  # reset index
    rets_tab['Trade Return Net'] = [i - float(fee * time_to_bnd) for i in perf_df['Trade Return'].to_list()]
    fc_date = pd.to_datetime(Trade['EXPIR_DATE'].iloc[0]).strftime("%m-%d-%Y")
    upside = perf_df['Trade Return'].max()

    if output == 'line':  # line output
        return trade_line
    elif output == 'trades':  # trades output
        print(trade_details)
        print(bnd.filter(['Instrument', "CF_LAST", "Shares", '1x Cost', ]))
    elif output == 'execute':
        op_exec = trade_details
        op_exec['units'] = [1.5, 1.5, 1]  # numerical correction for rounding
        bnd_exec = bnd.filter(['Instrument', 'Shares', 'CF_LAST'])
        bnd_exec = bnd_exec.rename(columns={'Shares': 'units'})
        op_exec.to_csv('assets/temp_op.csv')
        bnd_exec.to_csv('assets/temp_bnd.csv')
    else:
        # Figure objects
        trade_title = f"{sym} Accelerator<br>1.5x Trade"
        trade_subtitle = f"Ending {fc_date}"
        table_description = f"The investor gets up to 150% of the price appreciation on {sym} beginning at a<br>" \
                            f"price of {int(LCALL['STRIKE_PRC'])}, capped at {int(SCALL['STRIKE_PRC'])}, and full protection between " \
                            f"{int(LCALL['STRIKE_PRC'])} and {int(SPUT['STRIKE_PRC'])} ({-protection:.2%})"
        chart_description = f"{sym} Price (Created on: {today} when {sym} was at {last})"
        trade_summary_data = [
            ('Structure', 'Accelerated Upside with Cap + Protection'),
            ('Term', f"{months_left} Months"),
            ('Underlying Asset', sym),
            ('Cap', f"{upside:.2%}"),
            ('Downside Before Protection', 'No Gap'),
            ('Protection', f"{protection:.2%}"),
            ('Fee', f"{fee:.2%} Annually"),
        ]

        # Create the figure
        fig = plot_trade(chart_data=rets_tab, trade_title=trade_title, trade_subtitle=trade_subtitle,
                         table_description=table_description, chart_description=chart_description,
                         trade_summary_data=trade_summary_data, firm=firm, fee=fee)

        pio.write_html(fig, file='html/{}_{}_{}_accel_150.html'.format(Date, firm, sym),
                       auto_open=False)
        pio.write_image(fig, file='Images/{}_{}_{}_accel_150.png'.format(Date, firm, sym), format='png', scale=6)
        if firm == "L":
            prs = Presentation('rebrand_test.pptx')
        elif firm == "O":
            prs = Presentation('rebrand_test_oak.pptx')
        else:
            prs = Presentation('rebrand_test_cfm.pptx')
        prs.save('Images/{}_{}_{}_accel_150.pptx'.format(Date, firm, sym))
        shapes = prs.slides[0].shapes
        pic = shapes.add_picture('Images/{}_{}_{}_accel_150.png'.format(Date, firm, sym), Inches(0.001), Inches(0.001),
                                 height=Inches(10.99), width=Inches(8.5))
        prs.save('Images/{}_{}_{}_accel_150.pptx'.format(Date, firm, sym))
        return


def kre_uncapped(get_data, firm, output, fee=0.0125):
    """
    KRE trade in December 2025 expiration using BSCP corporate bond ETF as collateral
    :param get_data: <str> y or n only - y will pull fresh data from tr, n will use stored data
    :param output: <str> "line" - will provide terms, "trades" will provide trades, "chart" will provide fact sheet
    :param firm: <str> "L" for Lido, "O" for Oakhurst
    :return: based on output param
    """
    if get_data == 'y':  # get data if necessary
        od.get_options(["KRE"])
        ytw = od.calc_credit('BSCP', get_data='y')

    else:  # otherwise pass
        ytw = od.calc_credit('BSCP', get_data='n')

    Date = '2025-12-19'
    sym = pd.read_pickle('assets/{}_sym.pkl'.format("KRE")).iloc[0]['ticker']  # import the symbol
    und = pd.read_pickle('assets/{}_cached_und.pkl'.format(sym))  ##Gets Underlying Price###
    bnd_dat = pd.read_pickle('assets/bond_data.pkl')  # import bond data
    last = und.iloc[0]['CF_LAST']  # get most recent price

    # Set up Bond Data

    # begin options work
    # options chain set up
    chain = pd.read_pickle('assets/{}_cached_chain.pkl'.format(sym))  # read in the options chaing to pandas
    chain['MID'] = (chain['CF_BID'] + chain['CF_ASK']) / 2
    trade_chain = chain[chain['EXPIR_DATE'] == Date]  # isolate the chain we want
    sorted_trade_chain = trade_chain.sort_values(['PUTCALLIND', 'STRIKE_PRC'],
                                                 axis=0, )  # sort the chain by type and strike
    sorted_trade_chain_calls = sorted_trade_chain[sorted_trade_chain['PUTCALLIND'] == 'CALL']  # create a call df
    sorted_trade_chain_puts = sorted_trade_chain[sorted_trade_chain['PUTCALLIND'] == 'PUT ']  # create a put df

    # ATM CALL
    LCALL = sorted_trade_chain_calls.loc[
        sorted_trade_chain_calls['STRIKE_PRC'].sub(last).abs().idxmin()]  # get atm call

    # Solve for the other two legs
    date = datetime.datetime.strptime(Date, '%Y-%m-%d').date()  # convert date to datetime
    today = datetime.datetime.now().date()  # get today's date
    time_left = date - today  # days left
    adj_time_left = time_left / datetime.timedelta(days=1)  # convert to flt
    time_to_bnd = (adj_time_left / 365)  # convert to annual
    bnd = ek.get_data("BSCP.O", fields=['CF_LAST'])[0]  ###Gets BSCN Price###
    rf = 1 + (ytw * time_to_bnd)
    bnd['1x Cost'] = (LCALL['STRIKE_PRC'] * 100) / rf
    bnd['Shares'] = (bnd['1x Cost'] / bnd['CF_LAST']).round(0)
    notional = LCALL['STRIKE_PRC']  # trade's notional is the strike of the option
    bnd['TR'] = (bnd['1x Cost'] / 100 * rf) - bnd['1x Cost'] / 100  # calculate total expected return from bond
    # begin work on put
    tgt = LCALL["MID"] - bnd['TR'].iloc[0]  # tgt put value
    SPUT = sorted_trade_chain_puts.loc[sorted_trade_chain_puts['MID'].sub(tgt).abs().idxmin()]  # get atm call
    LCALL['Trans'] = 'BCO'
    SPUT['Trans'] = 'SPO'  # add trade type
    Trade = pd.concat([LCALL.to_frame().T, SPUT.to_frame().T, ])  # create Trade df
    exp_date = pd.to_datetime(Trade['EXPIR_DATE'].iloc[0]).strftime('%y%m%d')  # get options formated date
    option_type = [s[0] for s in Trade["PUTCALLIND"].to_list()]  # isolate first letter of option type
    strikes = Trade['STRIKE_PRC'].to_list()  # # isolate strikes in list
    option_sym = ['{}{}{}{}'.format(sym, exp_date, option_type[i], int(strikes[i])) for i in
                  range(0, len(Trade))]  # create symbols
    Trade['Symbol'] = option_sym  # add to df
    trade_details = Trade.filter(['Trans', 'Symbol', 'STRIKE_PRC', 'CF_BID', 'CF_ASK'])  # create trade_det df
    upside = "Uncapped"  # upside calculation
    protection = np.abs((SPUT['STRIKE_PRC'] / LCALL['STRIKE_PRC']) - 1)  # protection calculation
    annual_p = protection * (365 / adj_time_left)
    months_left = np.round((adj_time_left / 365) * 12, 1)
    trade_line = pd.DataFrame({"Underlying Asset": sym, "Asset Price": last,
                               'Minimum ($)': round(LCALL['STRIKE_PRC'] * 100, -2),
                               "Expiration Date": Trade['EXPIR_DATE'].iloc[0], "Months Left": months_left,
                               'Potential Upside (%)': upside,
                               'Annual Potential Upside (%)': upside,
                               'Downside Before Protection (%)': '{:.2%}'.format(0),
                               'Protection (%)': '{:.2%}'.format(protection),
                               'Annual_Protection (%)': '{:.2%}'.format(annual_p),
                               'Spread (%)': 'NA', 'Collateral': 'BSCP'}, index=[0])  # creates line

    # begin work on plot data:
    chrt_rng = np.linspace(LCALL['STRIKE_PRC'] * 0.5, LCALL['STRIKE_PRC'] * 1.5, 50, dtype=float)  # set chart range
    chrt_rng = np.round(chrt_rng, 2)
    lcall_ev = [np.maximum(p - LCALL['STRIKE_PRC'], 0) for p in chrt_rng]  # calc lcall end val
    sput_ev = [np.maximum(SPUT['STRIKE_PRC'] - p, 0) * -1 for p in chrt_rng]  # calc sput end val
    perf_df = pd.DataFrame({"{} Price".format(sym): chrt_rng, "LCALL": lcall_ev, "SPUT": sput_ev,
                            "BOND": LCALL['STRIKE_PRC']})  # Create the df
    perf_df = perf_df.set_index("{} Price".format(sym))  # Set the mkt px as the index
    perf_df['Trade'] = perf_df.sum(axis=1)  # calculate total value
    cost = LCALL['MID'] - SPUT['MID'] + LCALL['STRIKE_PRC'] - bnd['TR'].iloc[0]  # trade cost
    perf_df['Trade Return'] = (perf_df['Trade'] / (cost)) - 1  # trade return
    perf_df = perf_df.sort_index(ascending=False)  # reorder in descending
    perf_df['{} Price Return'.format(sym)] = [(p / last) - 1 for p in perf_df.index]  # add underlying performance
    rets_tab = perf_df.filter(perf_df.columns[-2:]).reset_index()  # reset index
    rets_tab['Trade Return Net'] = [i - float(fee * time_to_bnd) for i in perf_df['Trade Return'].to_list()]
    fc_date = pd.to_datetime(Trade['EXPIR_DATE'].iloc[0]).strftime("%m-%d-%Y")

    if output == 'line':  # line output
        return trade_line
    elif output == 'trades':  # trades output
        print(trade_details)
        print(bnd.filter(['Instrument', "CF_LAST", "Shares", '1x Cost', ]))
    elif output == 'execute':
        op_exec = trade_details
        op_exec['MID'] = (op_exec['CF_BID'] + op_exec['CF_ASK']) / 2
        op_exec = op_exec.drop(['CF_BID', 'CF_ASK'], axis=1)
        op_exec['units'] = [1 for i in range(0, len(op_exec))]
        bnd_exec = bnd.filter(['Instrument', 'Shares', 'CF_LAST'])
        bnd_exec = bnd_exec.rename(columns={'Shares': 'units'})
        op_exec.to_csv('assets/temp_op.csv')
        bnd_exec.to_csv('assets/temp_bnd.csv')
    else:
        # Figure objects
        trade_title = "KRE Uncapped<br>Trade"
        trade_subtitle = f"Ending {fc_date}"
        table_description = f"The investor gets about 100% of the price appreciation on {sym} beginning at a<br>" \
                            f"price of {int(LCALL['STRIKE_PRC'])} and full protection between " \
                            f"{int(LCALL['STRIKE_PRC'])} and {int(SPUT['STRIKE_PRC'])} ({-protection:.2%})"
        chart_description = f"{sym} Price (Created on: {today} when {sym} was at {last})"
        trade_summary_data = [
            ('Structure', 'Unlimited Upside + Protection'),
            ('Term', f"{months_left} Months"),
            ('Underlying Asset', sym),
            ('Cap', 'Uncapped'),
            ('Downside Before Protection', 'No Gap'),
            ('Protection', f"{protection:.2%}"),
            ('Fee', f"{fee:.2%} Annually"),
        ]

        # Create the figure
        fig = plot_trade(chart_data=rets_tab, trade_title=trade_title, trade_subtitle=trade_subtitle,
                         table_description=table_description, chart_description=chart_description,
                         trade_summary_data=trade_summary_data, firm=firm, fee=fee)

        pio.write_html(fig, file='html/{}_{}_{}_kre_uncapped.html'.format(Date, firm, sym),
                       auto_open=False)
        pio.write_image(fig, file='Images/{}_{}_{}_kre_uncapped.png'.format(Date, firm, sym), format='png', scale=6)
        if firm == "L":
            prs = Presentation('rebrand_test.pptx')
        elif firm == "O":
            prs = Presentation('rebrand_test_oak.pptx')
        else:
            prs = Presentation('rebrand_test_cfm.pptx')
        prs.save('Images/{}_{}_{}_kre_uncapped.pptx'.format(Date, firm, sym))
        shapes = prs.slides[0].shapes
        pic = shapes.add_picture('Images/{}_{}_{}_kre_uncapped.png'.format(Date, firm, sym), Inches(0.001),
                                 Inches(0.001),
                                 height=Inches(10.99), width=Inches(8.5))
        prs.save('Images/{}_{}_{}_kre_uncapped.pptx'.format(Date, firm, sym))
        return


def xlf_uncapped(get_data, firm, output, fee=0.0125):
    """
    XLF trade in December 2025 expiration using BSCP corporate bond ETF as collateral
    :param get_data: <str> y or n only - y will pull fresh data from tr, n will use stored data
    :param output: <str> "line" - will provide terms, "trades" will provide trades, "chart" will provide fact sheet
    :param firm: <str> "L" for Lido, "O" for Oakhurst
    :return: based on output param
    """
    if get_data == 'y':  # get data if necessary
        od.get_options(["XLF"])
        ytw = od.calc_credit('BSCP', get_data='y')

    else:  # otherwise pass
        ytw = od.calc_credit('BSCP', get_data='n')

    Date = '2025-12-19'
    sym = pd.read_pickle('assets/{}_sym.pkl'.format("XLF")).iloc[0]['ticker']  # import the symbol
    und = pd.read_pickle('assets/{}_cached_und.pkl'.format(sym))  ##Gets Underlying Price###
    bnd_dat = pd.read_pickle('assets/bond_data.pkl')  # import bond data
    last = und.iloc[0]['CF_LAST']  # get most recent price

    # Set up Bond Data

    # begin options work
    # options chain set up
    chain = pd.read_pickle('assets/{}_cached_chain.pkl'.format(sym))  # read in the options chaing to pandas
    chain['MID'] = (chain['CF_BID'] + chain['CF_ASK']) / 2
    trade_chain = chain[chain['EXPIR_DATE'] == Date]  # isolate the chain we want
    sorted_trade_chain = trade_chain.sort_values(['PUTCALLIND', 'STRIKE_PRC'],
                                                 axis=0, )  # sort the chain by type and strike
    sorted_trade_chain_calls = sorted_trade_chain[sorted_trade_chain['PUTCALLIND'] == 'CALL']  # create a call df
    sorted_trade_chain_puts = sorted_trade_chain[sorted_trade_chain['PUTCALLIND'] == 'PUT ']  # create a put df

    # ATM CALL
    LCALL = sorted_trade_chain_calls.loc[
        sorted_trade_chain_calls['STRIKE_PRC'].sub(last).abs().idxmin()]  # get atm call

    # Solve for the other two legs
    date = datetime.datetime.strptime(Date, '%Y-%m-%d').date()  # convert date to datetime
    today = datetime.datetime.now().date()  # get today's date
    time_left = date - today  # days left
    adj_time_left = time_left / datetime.timedelta(days=1)  # convert to flt
    time_to_bnd = (adj_time_left / 365)  # convert to annual
    bnd = ek.get_data("BSCP.O", fields=['CF_LAST'])[0]  ###Gets BSCN Price###
    rf = 1 + (ytw * time_to_bnd)
    bnd['1x Cost'] = (LCALL['STRIKE_PRC'] * 100) / rf
    bnd['Shares'] = (bnd['1x Cost'] / bnd['CF_LAST']).round(0)
    notional = LCALL['STRIKE_PRC']  # trade's notional is the strike of the option
    bnd['TR'] = (bnd['1x Cost'] / 100 * rf) - bnd['1x Cost'] / 100  # calculate total expected return from bond
    # begin work on put
    tgt = LCALL["MID"] - bnd['TR'].iloc[0]  # tgt put value
    SPUT = sorted_trade_chain_puts.loc[sorted_trade_chain_puts['MID'].sub(tgt).abs().idxmin()]  # get atm call
    LCALL['Trans'] = 'BCO'
    SPUT['Trans'] = 'SPO'  # add trade type
    Trade = pd.concat([LCALL.to_frame().T, SPUT.to_frame().T, ])  # create Trade df
    exp_date = pd.to_datetime(Trade['EXPIR_DATE'].iloc[0]).strftime('%y%m%d')  # get options formated date
    option_type = [s[0] for s in Trade["PUTCALLIND"].to_list()]  # isolate first letter of option type
    strikes = Trade['STRIKE_PRC'].to_list()  # # isolate strikes in list
    option_sym = ['{}{}{}{}'.format(sym, exp_date, option_type[i], int(strikes[i])) for i in
                  range(0, len(Trade))]  # create symbols
    Trade['Symbol'] = option_sym  # add to df
    trade_details = Trade.filter(['Trans', 'Symbol', 'STRIKE_PRC', 'CF_BID', 'CF_ASK'])  # create trade_det df
    upside = "Uncapped"  # upside calculation
    protection = np.abs((SPUT['STRIKE_PRC'] / LCALL['STRIKE_PRC']) - 1)  # protection calculation
    annual_p = protection * (365 / adj_time_left)
    months_left = np.round((adj_time_left / 365) * 12, 1)
    trade_line = pd.DataFrame({"Underlying Asset": sym, "Asset Price": last,
                               'Minimum ($)': round(LCALL['STRIKE_PRC'] * 100, -2),
                               "Expiration Date": Trade['EXPIR_DATE'].iloc[0], "Months Left": months_left,
                               'Potential Upside (%)': upside,
                               'Annual Potential Upside (%)': upside,
                               'Downside Before Protection (%)': '{:.2%}'.format(0),
                               'Protection (%)': '{:.2%}'.format(protection),
                               'Annual_Protection (%)': '{:.2%}'.format(annual_p),
                               'Spread (%)': 'NA', 'Collateral': 'BSCP'}, index=[0])  # creates line

    # begin work on plot data:
    chrt_rng = np.linspace(LCALL['STRIKE_PRC'] * 0.5, LCALL['STRIKE_PRC'] * 1.5, 50, dtype=float)  # set chart range
    chrt_rng = np.round(chrt_rng, 2)
    lcall_ev = [np.maximum(p - LCALL['STRIKE_PRC'], 0) for p in chrt_rng]  # calc lcall end val
    sput_ev = [np.maximum(SPUT['STRIKE_PRC'] - p, 0) * -1 for p in chrt_rng]  # calc sput end val
    perf_df = pd.DataFrame({"{} Price".format(sym): chrt_rng, "LCALL": lcall_ev, "SPUT": sput_ev,
                            "BOND": LCALL['STRIKE_PRC']})  # Create the df
    perf_df = perf_df.set_index("{} Price".format(sym))  # Set the mkt px as the index
    perf_df['Trade'] = perf_df.sum(axis=1)  # calculate total value
    cost = LCALL['MID'] - SPUT['MID'] + LCALL['STRIKE_PRC'] - bnd['TR'].iloc[0]  # trade cost
    perf_df['Trade Return'] = (perf_df['Trade'] / (cost)) - 1  # trade return
    perf_df = perf_df.sort_index(ascending=False)  # reorder in descending
    perf_df['{} Price Return'.format(sym)] = [(p / last) - 1 for p in perf_df.index]  # add underlying performance
    rets_tab = perf_df.filter(perf_df.columns[-2:]).reset_index()  # reset index
    rets_tab['Trade Return Net'] = [i - float(fee * time_to_bnd) for i in perf_df['Trade Return'].to_list()]
    fc_date = pd.to_datetime(Trade['EXPIR_DATE'].iloc[0]).strftime("%m-%d-%Y")

    if output == 'line':  # line output
        return trade_line
    elif output == 'trades':  # trades output
        print(trade_details)
        print(bnd.filter(['Instrument', "CF_LAST", "Shares", '1x Cost', ]))
    elif output == 'execute':
        op_exec = trade_details
        op_exec['MID'] = (op_exec['CF_BID'] + op_exec['CF_ASK']) / 2
        op_exec = op_exec.drop(['CF_BID', 'CF_ASK'], axis=1)
        op_exec['units'] = [1 for i in range(0, len(op_exec))]
        bnd_exec = bnd.filter(['Instrument', 'Shares', 'CF_LAST'])
        bnd_exec = bnd_exec.rename(columns={'Shares': 'units'})
        op_exec.to_csv('assets/temp_op.csv')
        bnd_exec.to_csv('assets/temp_bnd.csv')
    else:
        # Figure objects
        trade_title = "XLF Uncapped<br>Trade"
        trade_subtitle = f"Ending {fc_date}"
        table_description = f"The investor gets about 100% of the price appreciation on {sym} beginning at a<br>" \
                            f"price of {int(LCALL['STRIKE_PRC'])} and full protection between " \
                            f"{int(LCALL['STRIKE_PRC'])} and {int(SPUT['STRIKE_PRC'])} ({-protection:.2%})"
        chart_description = f"{sym} Price (Created on: {today} when {sym} was at {last})"
        trade_summary_data = [
            ('Structure', 'Unlimited Upside + Protection'),
            ('Term', f"{months_left} Months"),
            ('Underlying Asset', sym),
            ('Cap', 'Uncapped'),
            ('Downside Before Protection', 'No Gap'),
            ('Protection', f"{protection:.2%}"),
            ('Fee', f"{fee:.2%} Annually"),
        ]

        # Create the figure
        fig = plot_trade(chart_data=rets_tab, trade_title=trade_title, trade_subtitle=trade_subtitle,
                         table_description=table_description, chart_description=chart_description,
                         trade_summary_data=trade_summary_data, firm=firm, fee=fee)

        pio.write_html(fig, file='html/{}_{}_{}_xlf_uncapped.html'.format(Date, firm, sym),
                       auto_open=False)
        pio.write_image(fig, file='Images/{}_{}_{}_xlf_uncapped.png'.format(Date, firm, sym), format='png', scale=6)
        if firm == "L":
            prs = Presentation('rebrand_test.pptx')
        elif firm == "O":
            prs = Presentation('rebrand_test_oak.pptx')
        else:
            prs = Presentation('rebrand_test_cfm.pptx')
        prs.save('Images/{}_{}_{}_xlf_uncapped.pptx'.format(Date, firm, sym))
        shapes = prs.slides[0].shapes
        pic = shapes.add_picture('Images/{}_{}_{}_xlf_uncapped.png'.format(Date, firm, sym), Inches(0.001),
                                 Inches(0.001),
                                 height=Inches(10.99), width=Inches(8.5))
        prs.save('Images/{}_{}_{}_xlf_uncapped.pptx'.format(Date, firm, sym))
        return
